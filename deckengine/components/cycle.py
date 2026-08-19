"""cycle — flywheel / reinforcing loop (canon: flywheel).

A PIL-drawn ring of clockwise arrow segments (core/rings.py, cached);
stage labels are measured pptx text boxes OUTSIDE the ring at segment
centroids (all text stays under the layout contract); optional hub label
centered in the hole. highlight_index accents one segment AND its label.

Geometry is a pure function of (width, n, label_h): ring diameter from
the fixed component height, label boxes on an ellipse just outside the
ring — the hub_spoke convention, parity by construction.
"""
from __future__ import annotations

import math
from dataclasses import dataclass

from ..core.bbox import BBox
from ..core.fit_text import Span, fit_text
from ..core.pptx_text import add_text_box, write_fit_result
from ..core.rings import get_cycle_ring
from ..core.units import inch, to_inch
from ..schema.components import CycleSpec
from ..schema.rich import parse_rich
from .base import Component, RenderContext, register

_RING_D = inch(2.7)        # ring image side
_LABEL_W_FRAC = 0.24
_LABEL_W_MAX = inch(2.0)
_MIN_LABEL = 7.0
_PROBE_H = 10_000_000


@dataclass(frozen=True)
class _Plan:
    label_w: int
    label_h: int
    total: int
    rx: int
    ry: int


@register("cycle")
class Cycle(Component):
    spec_model = CycleSpec

    def _plan(self, data: CycleSpec, width: int,
              ctx: RenderContext) -> _Plan:
        line_h = ctx.measurer.line_height_emu(
            [Span("Hg", bold=True)], ctx.font("body"), ctx.size("small"))
        label_h = 2 * line_h + ctx.theme.spacing(0.3)
        label_w = min(round(width * _LABEL_W_FRAC), _LABEL_W_MAX)
        ry = _RING_D // 2 + label_h // 2 + ctx.theme.spacing(0.5)
        rx = max(ry, min(width // 2 - label_w // 2,
                         _RING_D // 2 + label_w // 2
                         + ctx.theme.spacing(1.0)))
        return _Plan(label_w, label_h, 2 * ry + label_h, rx, ry)

    def measure(self, data: CycleSpec, width: int,
                ctx: RenderContext) -> int:
        return self._plan(data, width, ctx).total

    def render(self, slide, data: CycleSpec, bbox: BBox,
               ctx: RenderContext) -> int:
        theme = ctx.theme
        plan = self._plan(data, bbox.w, ctx)
        if plan.total > bbox.h:
            ctx.report.warn(
                f"cycle: needs {to_inch(plan.total):.2f}in but bbox is "
                f"{to_inch(bbox.h):.2f}in; rendering anyway (may clip)")
        n = len(data.stages)
        hi = data.highlight_index
        if hi is not None and not 0 <= hi < n:
            ctx.report.warn(f"cycle: highlight_index {hi} out of range for "
                            f"{n} stages; ignoring")
            hi = None
        family = ctx.font("body")
        cx = bbox.x + bbox.w // 2
        cy = bbox.y + plan.total // 2

        png = get_cycle_ring(n, theme.color("primary_dark"),
                             theme.color("accent"), hi)
        from pptx.util import Emu
        slide.shapes.add_picture(
            str(png), Emu(cx - _RING_D // 2), Emu(cy - _RING_D // 2),
            Emu(_RING_D), Emu(_RING_D))

        # hub label in the ring's hole
        if data.hub:
            hub_w = round(_RING_D * 0.46)
            fit = fit_text(
                [Span(sp.text, bold=True, italic=sp.italic,
                      color_role="ink")
                 for sp in parse_rich(data.hub, base_color_role="ink")],
                BBox(0, 0, hub_w, _PROBE_H), family,
                max_size=ctx.size("small"), min_size=_MIN_LABEL,
                max_lines=2, measurer=ctx.measurer)
            if fit.truncated:
                ctx.report.truncated(f"cycle hub: {data.hub[:40]!r}")
            box = add_text_box(
                slide, BBox(cx - hub_w // 2, cy - fit.height_emu // 2,
                            hub_w, fit.height_emu), align="center")
            write_fit_result(box.text_frame, fit, theme, family=family,
                             align="center", default_color_role="ink")

        # stage labels at segment centroids, pushed outside the ring
        seg = 360.0 / n
        for i, stage in enumerate(data.stages):
            mid = math.radians(-90 + (i + 0.5) * seg)
            px_ = cx + round(plan.rx * math.cos(mid))
            py_ = cy + round(plan.ry * math.sin(mid))
            ink = "accent" if i == hi else "ink"
            fit = fit_text(
                [Span(sp.text, bold=True, italic=sp.italic, color_role=ink)
                 for sp in parse_rich(stage, base_color_role=ink)],
                BBox(0, 0, plan.label_w, _PROBE_H), family,
                max_size=ctx.size("small"), min_size=_MIN_LABEL,
                max_lines=2, measurer=ctx.measurer)
            if fit.truncated:
                ctx.report.truncated(f"cycle stage: {stage[:40]!r}")
            # horizontal alignment follows which side of the ring we're on
            align = ("center" if abs(math.cos(mid)) < 0.35
                     else ("left" if math.cos(mid) > 0 else "right"))
            box = add_text_box(
                slide, BBox(px_ - plan.label_w // 2,
                            py_ - fit.height_emu // 2,
                            plan.label_w, fit.height_emu), align=align)
            write_fit_result(box.text_frame, fit, theme, family=family,
                             align=align, default_color_role=ink)
        return plan.total
