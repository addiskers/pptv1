"""data_table — flagship shape-grid table (the 'Priority Products' pattern).

Grouped rows with a vertically merged group-label column, badge pills in
cells, and heatmap readiness columns. Rendered as a SHAPE GRID — one rect per
cell with a pinned text frame — never a native OOXML table (native cells
cannot hold shapes like badge pills, and row heights are only minimums that
PowerPoint re-grows, which breaks the measure/render contract).

Geometry contract:
  header_h = tallest header fit (micro bold, <=2 lines) + spacing(0.4)
  row_h    = one micro line pitch + spacing(0.35)        (uniform, all rows)
  measure  = header_h + total_rows * row_h                (exact)
"""
from __future__ import annotations

from dataclasses import replace

# Emu is only used to re-apply the horizontal insets that make_text_frame
# zeroes; the inset width is always subtracted from the fit_text budget so
# wrapping can never disagree with what PowerPoint renders.
from pptx.util import Emu

from ..core.bbox import BBox
from ..core.fit_text import WRAP_SAFETY, FitResult, Span, fit_text
from ..core.pptx_shapes import add_shape
from ..core.pptx_text import make_text_frame, write_fit_result, write_spans_paragraph
from ..core.units import pt
from ..schema.components import DataColumnSpec, DataTableSpec
from ..schema.rich import parse_rich
from .base import Component, RenderContext, register

_CELL_MIN_PT = 6.0        # table cells shrink to 6pt before truncating
_HEATMAP_INK_FLIP = 0.55  # normalized value above which cell text flips to inverse ink
_GRID_W_PT = 0.5          # thin rule weight for all cell borders
_PROBE_H = 10_000_000     # generous probe height; real caps are row/pill boxes
_NA_TOKENS = frozenset({"N/A", "NA", "-", "–", "—"})


@register("data_table")
class DataTable(Component):
    spec_model = DataTableSpec

    # -- shared geometry (measure and render use exactly these) -----------

    def _col_widths(self, data: DataTableSpec, width: int,
                    ctx: RenderContext | None = None) -> list[int]:
        """Normalize column fracs into integer widths that sum to `width`.

        With ctx and auto_widen_label_col (default), col 0 grows toward the
        width its group labels need (2-line micro-bold estimate), capped at
        +25%; the other columns shrink proportionally. Deterministic pure
        math — measure(), row_offsets() and render() share it exactly.
        """
        total = sum(c.frac for c in data.columns)
        ws = [round(width * c.frac / total) for c in data.columns]
        ws[-1] = width - sum(ws[:-1])  # absorb rounding into the last column
        if (ctx is None or not data.auto_widen_label_col or len(ws) < 2
                or not data.groups):
            return ws
        pad = ctx.theme.spacing(0.25)
        family = ctx.font("body")
        size = ctx.size("micro")
        need = 0
        for g in data.groups:
            label = g.label.strip()
            if not label:
                continue
            lw = ctx.measurer.span_width_emu(Span(label, bold=True),
                                             family, size)
            word = max((ctx.measurer.span_width_emu(Span(w, bold=True),
                                                    family, size)
                        for w in label.split()), default=0)
            # a label gets two lines: need >= its widest word and roughly
            # half its full width (1.1 wrap-inefficiency factor)
            need = max(need, word, round(lw * 0.55))
        need = round(need / WRAP_SAFETY) + 2 * pad
        new0 = min(max(ws[0], need), round(ws[0] * 1.25))
        if new0 <= ws[0]:
            return ws
        delta = new0 - ws[0]
        rest = sum(ws[1:])
        out = [new0] + [w - round(delta * w / rest) for w in ws[1:]]
        out[-1] = width - sum(out[:-1])  # re-absorb rounding
        return out

    def _header_fits(self, data: DataTableSpec, col_ws: list[int],
                     ctx: RenderContext) -> list[FitResult]:
        pad = ctx.theme.spacing(0.2)
        fits: list[FitResult] = []
        for col, w in zip(data.columns, col_ws):
            spans = [replace(s, bold=True) for s in parse_rich(col.label)]
            probe = BBox(0, 0, max(1, w - 2 * pad), _PROBE_H)
            fits.append(fit_text(spans, probe, ctx.font("body"),
                                 max_size=ctx.size("micro"), min_size=_CELL_MIN_PT,
                                 max_lines=2, measurer=ctx.measurer))
        return fits

    def _header_h(self, fits: list[FitResult], ctx: RenderContext) -> int:
        return max(f.height_emu for f in fits) + ctx.theme.spacing(0.4)

    def _row_h(self, ctx: RenderContext) -> int:
        one_line = ctx.measurer.line_height_emu([], ctx.font("body"),
                                                ctx.size("micro"))
        return one_line + ctx.theme.spacing(0.35)

    # -- public contract ---------------------------------------------------

    def measure(self, data: DataTableSpec, width: int, ctx: RenderContext) -> int:
        col_ws = self._col_widths(data, width, ctx)
        header_h = self._header_h(self._header_fits(data, col_ws, ctx), ctx)
        total_rows = sum(len(g.rows) for g in data.groups)
        return header_h + total_rows * self._row_h(ctx)

    def row_offsets(self, data: DataTableSpec, width: int,
                    ctx: RenderContext) -> list[int]:
        """y offsets (EMU, from the component top) of each group boundary.

        offsets[i] is the top of group i (offsets[0] == header height);
        offsets[-1] is the table bottom, so len == len(groups) + 1. Future
        slide-splitting cuts the table at these offsets only.
        """
        col_ws = self._col_widths(data, width, ctx)
        header_h = self._header_h(self._header_fits(data, col_ws, ctx), ctx)
        row_h = self._row_h(ctx)
        offsets = [header_h]
        for group in data.groups:
            offsets.append(offsets[-1] + len(group.rows) * row_h)
        return offsets

    @staticmethod
    def _emphasis_targets(data: DataTableSpec) -> tuple[int, tuple[int, int]]:
        """(group_idx, (group_idx, row_idx)) for the emphasis matches,
        -1 / (-1,-1) when unset or unmatched. Case-insensitive on the
        group label / first data cell."""
        gi_hit, row_hit = -1, (-1, -1)
        want_g = (data.highlight_group or "").strip().casefold()
        want_r = (data.highlight_row or "").strip().casefold()
        for gi, group in enumerate(data.groups):
            if want_g and group.label.strip().casefold() == want_g:
                gi_hit = gi
            if want_r:
                for ri, row in enumerate(group.rows):
                    if row and str(row[0]).strip().casefold() == want_r:
                        row_hit = (gi, ri)
        return gi_hit, row_hit

    def render(self, slide, data: DataTableSpec, bbox: BBox,
               ctx: RenderContext) -> int:
        theme = ctx.theme
        col_ws = self._col_widths(data, bbox.w, ctx)
        header_fits = self._header_fits(data, col_ws, ctx)
        header_h = self._header_h(header_fits, ctx)
        hi_group, hi_row = self._emphasis_targets(data)
        if (data.highlight_group or "").strip() and hi_group < 0:
            ctx.report.warn(f"data_table: highlight_group "
                            f"{data.highlight_group!r} matches no group")
        if (data.highlight_row or "").strip() and hi_row == (-1, -1):
            ctx.report.warn(f"data_table: highlight_row "
                            f"{data.highlight_row!r} matches no row")
        row_h = self._row_h(ctx)
        total_rows = sum(len(g.rows) for g in data.groups)
        total_h = header_h + total_rows * row_h
        if total_h > bbox.h:
            ctx.report.warn(
                f"data_table: needs {total_h} EMU but bbox has {bbox.h}; "
                "split at row_offsets() upstream")
        elif ctx.fill_hint and bbox.h > total_h and total_rows:
            # flex: given extra height, relax row pitch (capped) to fill the zone
            row_h += min((bbox.h - total_h) // total_rows, round(row_h * 0.6))
            total_h = header_h + total_rows * row_h

        xs = [bbox.x]
        for w in col_ws:
            xs.append(xs[-1] + w)

        # header row: one rect per column, micro bold inverse ink, centered
        for i, fit in enumerate(header_fits):
            cell = BBox(xs[i], bbox.y, col_ws[i], header_h)
            s = add_shape(slide, cell, theme, fill_role=data.header_fill_role,
                          line_role="grid", line_w_pt=_GRID_W_PT)
            tf = make_text_frame(s, align="center", anchor="middle")
            write_fit_result(tf, fit, theme, family=ctx.font("body"),
                             align="center", default_color_role="inverse_ink")

        # body: per group, one merged label rect + one rect per data cell
        heat_bounds = self._heat_bounds(data)
        soft = theme.soft("accent")
        y = bbox.y + header_h
        table_w = sum(col_ws)
        for gi, group in enumerate(data.groups):
            group_h = len(group.rows) * row_h
            body_fill = "surface" if (data.zebra and gi % 2 == 1) else "bg"
            g_emph = gi == hi_group
            self._group_label_cell(slide, BBox(xs[0], y, col_ws[0], group_h),
                                   group.label, ctx,
                                   fill_hex=soft if g_emph else None,
                                   accent=g_emph)
            for ri, row in enumerate(group.rows):
                ry = y + ri * row_h
                r_emph = g_emph or (gi, ri) == hi_row
                for ci in range(1, len(data.columns)):
                    value = row[ci - 1] if ci - 1 < len(row) else ""
                    cell = BBox(xs[ci], ry, col_ws[ci], row_h)
                    self._render_cell(slide, cell, data.columns[ci], value,
                                      body_fill, data, heat_bounds, ctx,
                                      fill_hex=soft if r_emph else None,
                                      bold=r_emph)
            # emphasis ring: one no-fill accent outline over the group band
            if g_emph:
                add_shape(slide, BBox(bbox.x, y, table_w, group_h), theme,
                          shape="rect", line_role="accent", line_w_pt=1.5)
            y += group_h
        if hi_row != (-1, -1) and hi_row[0] != hi_group:
            gi, ri = hi_row
            ry = (bbox.y + header_h
                  + sum(len(g.rows) for g in data.groups[:gi]) * row_h
                  + ri * row_h)
            add_shape(slide, BBox(bbox.x, ry, table_w, row_h), theme,
                      shape="rect", line_role="accent", line_w_pt=1.5)
        return total_h

    def _heat_bounds(self, data: DataTableSpec) -> tuple[float, float]:
        """Auto-scale heatmap shading to the actual data range when the spec
        left the default 0-100 bounds — values clustered in a narrow band
        (e.g. 1-28% shares) would otherwise all render nearly white."""
        if (data.heatmap_lo, data.heatmap_hi) != (0.0, 100.0):
            return data.heatmap_lo, data.heatmap_hi
        heat_cols = [i for i, c in enumerate(data.columns)
                     if c.cell_kind == "heatmap"]
        vals = []
        for g in data.groups:
            for row in g.rows:
                for ci in heat_cols:
                    if ci - 1 < len(row):
                        v = self._parse_heat(row[ci - 1])
                        if v is not None:
                            vals.append(v)
        if len(set(vals)) < 2:
            return data.heatmap_lo, data.heatmap_hi
        return min(vals), max(vals)

    # -- cell renderers ------------------------------------------------------

    def _group_label_cell(self, slide, cell: BBox, label: str,
                          ctx: RenderContext, *,
                          fill_hex: str | None = None,
                          accent: bool = False) -> None:
        """Merged column-0 rect spanning the group's full pixel height."""
        theme = ctx.theme
        pad = theme.spacing(0.25)
        s = add_shape(slide, cell, theme,
                      fill_hex=fill_hex,
                      fill_role=None if fill_hex else "bg",
                      line_role="grid", line_w_pt=_GRID_W_PT)
        if not label.strip():
            return
        fit = fit_text([Span(label, bold=True)],
                       BBox(0, 0, max(1, cell.w - 2 * pad), cell.h),
                       ctx.font("body"), max_size=ctx.size("micro"),
                       min_size=_CELL_MIN_PT, max_lines=2, measurer=ctx.measurer)
        if fit.truncated:
            ctx.report.truncated(f"data_table group label: {label!r}")
        tf = make_text_frame(s, align="left", anchor="middle")
        tf.margin_left = Emu(pad)
        tf.margin_right = Emu(pad)
        write_fit_result(tf, fit, theme, family=ctx.font("body"), align="left",
                         default_color_role="accent" if accent else "primary")

    def _render_cell(self, slide, cell: BBox, col: DataColumnSpec, value: str,
                     body_fill: str, data: DataTableSpec,
                     heat_bounds: tuple[float, float],
                     ctx: RenderContext, *, fill_hex: str | None = None,
                     bold: bool = False) -> None:
        kind = col.cell_kind
        if kind == "badge":
            self._badge_cell(slide, cell, value, body_fill, ctx)
        elif kind == "heatmap":
            self._heatmap_cell(slide, cell, value, body_fill, heat_bounds, ctx)
        elif kind == "dot":
            self._dot_cell(slide, cell, value, body_fill, ctx)
        elif kind == "number":
            s = add_shape(slide, cell, ctx.theme, fill_hex=fill_hex,
                          fill_role=None if fill_hex else body_fill,
                          line_role="grid", line_w_pt=_GRID_W_PT)
            self._write_cell_text(s, cell, value, align="right",
                                  color_role="ink", ctx=ctx, bold=bold)
        else:  # "text"
            s = add_shape(slide, cell, ctx.theme, fill_hex=fill_hex,
                          fill_role=None if fill_hex else body_fill,
                          line_role="grid", line_w_pt=_GRID_W_PT)
            self._write_cell_text(s, cell, value, align=col.align,
                                  color_role="ink", ctx=ctx, bold=bold)

    def _write_cell_text(self, shape, cell: BBox, value: str, *, align: str,
                         color_role: str, ctx: RenderContext,
                         bold: bool = False) -> None:
        """Single-line micro text inside a cell rect, symmetric side insets."""
        if not value.strip():
            return
        theme = ctx.theme
        pad = theme.spacing(0.25)
        fit = fit_text([Span(value, bold=bold)],
                       BBox(0, 0, max(1, cell.w - 2 * pad), _PROBE_H),
                       ctx.font("body"), max_size=ctx.size("micro"),
                       min_size=_CELL_MIN_PT, max_lines=1, measurer=ctx.measurer)
        if fit.truncated:
            ctx.report.truncated(f"data_table cell: {value!r}")
        tf = make_text_frame(shape, align=align, anchor="middle")
        tf.margin_left = Emu(pad)
        tf.margin_right = Emu(pad)
        write_fit_result(tf, fit, theme, family=ctx.font("body"), align=align,
                         default_color_role=color_role)

    def _badge_cell(self, slide, cell: BBox, value: str, body_fill: str,
                    ctx: RenderContext) -> None:
        theme = ctx.theme
        add_shape(slide, cell, theme, fill_role=body_fill,
                  line_role="grid", line_w_pt=_GRID_W_PT)
        code = value.strip()
        if not code:
            return
        fill_hex = theme.badge_palette.get(code)
        if fill_hex is None:
            ctx.report.warn(
                f"data_table: badge code {code!r} not in theme badge_palette")
            fill_hex = theme.color("ink_muted")
        pad = theme.spacing(0.25)
        fit = fit_text([Span(code, bold=True)],
                       BBox(0, 0, max(1, cell.w - 2 * pad), _PROBE_H),
                       ctx.font("body"), max_size=ctx.size("micro"),
                       min_size=_CELL_MIN_PT, max_lines=1, measurer=ctx.measurer)
        text_w = max((ln.width_emu for ln in fit.lines), default=0)
        pill_h = max(pt(4), min(cell.h - theme.spacing(0.15),
                                fit.height_emu + theme.spacing(0.1)))
        pill_w = max(pill_h, min(cell.w - theme.spacing(0.2),
                                 text_w + theme.spacing(0.5)))
        pill = BBox(cell.x + (cell.w - pill_w) // 2,
                    cell.y + (cell.h - pill_h) // 2, pill_w, pill_h)
        p = add_shape(slide, pill, theme, shape="rounded", fill_hex=fill_hex,
                      corner_radius=0.5)
        tf = make_text_frame(p, align="center", anchor="middle")
        write_spans_paragraph(tf, [Span(code, bold=True)], fit.size_pt, theme,
                              family=ctx.font("body"), align="center",
                              default_color_role="inverse_ink")

    def _heatmap_cell(self, slide, cell: BBox, value: str, body_fill: str,
                      heat_bounds: tuple[float, float],
                      ctx: RenderContext) -> None:
        theme = ctx.theme
        v = self._parse_heat(value)
        if v is None:
            # 'N/A' / '' / unparseable: plain muted text cell on the body fill
            s = add_shape(slide, cell, theme, fill_role=body_fill,
                          line_role="grid", line_w_pt=_GRID_W_PT)
            self._write_cell_text(s, cell, value.strip(), align="center",
                                  color_role="ink_muted", ctx=ctx)
            return
        lo, hi = heat_bounds
        norm = 0.0 if hi <= lo else min(1.0, max(0.0, (v - lo) / (hi - lo)))
        s = add_shape(slide, cell, theme, fill_hex=theme.heatmap_color(v, lo, hi),
                      line_role="grid", line_w_pt=_GRID_W_PT)
        role = "inverse_ink" if norm > _HEATMAP_INK_FLIP else "ink"
        self._write_cell_text(s, cell, value.strip(), align="center",
                              color_role=role, ctx=ctx)

    def _dot_cell(self, slide, cell: BBox, value: str, body_fill: str,
                  ctx: RenderContext) -> None:
        theme = ctx.theme
        s = add_shape(slide, cell, theme, fill_role=body_fill,
                      line_role="grid", line_w_pt=_GRID_W_PT)
        if not value.strip():
            return
        tf = make_text_frame(s, align="center", anchor="middle")
        write_spans_paragraph(tf, [Span("•")], ctx.size("micro"), theme,
                              family=ctx.font("body"), align="center",
                              default_color_role="ink")

    @staticmethod
    def _parse_heat(value: str) -> float | None:
        txt = value.strip()
        if not txt or txt.upper() in _NA_TOKENS:
            return None
        try:
            return float(txt.replace("%", "").strip())
        except ValueError:
            return None
