"""xy_chart — scatter / bubble / quadrant-scatter (the relationship family).

Census: bubble 5+2, quadrant_scatter 2+2, dot_plot 3+1 — the "does X
relate to Y" forms, including the BCG growth-share engine. Points are
positioned data (own x/y, optional size -> bubble). quadrants=True draws
dashed midlines at the data midpoints and optional corner captions —
the strategy-scatter look.

The chart itself is a native pptx XY/bubble chart (fully editable);
per-point labels ride python-pptx data labels. highlight accents one
point and mutes the rest — emphasis from birth.

measure(): fixed natural height (like native_chart) + caption-free —
parity by arithmetic.
"""
from __future__ import annotations

from pptx.chart.data import BubbleChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.util import Emu, Pt

from ..core.bbox import BBox
from ..core.fit_text import Span
from ..core.pptx_shapes import add_hline, add_vline
from ..core.pptx_text import add_text_box, write_spans_paragraph
from ..core.units import inch, to_inch, to_pt
from ..schema.components import XYChartSpec
from .base import Component, RenderContext, register

_MUTE_ROLE = "ink_muted"
_QUAD_LABEL_INSET = 0.02   # corner caption inset as a plot fraction
# plot-area insets as frame fractions (benchmark-overlay convention;
# slightly larger left/bottom than native_chart's — axis TITLES here)
_PLOT_LEFT = 0.11
_PLOT_RIGHT = 0.03
_PLOT_TOP = 0.04
_PLOT_BOTTOM = 0.16


def _nice_bounds(vals: list[float]) -> tuple[float, float]:
    """Padded bounds snapped to a 'nice' step so axis ticks read clean
    (26.15 -> 30, 2.17 -> 2)."""
    import math
    lo, hi = min(vals), max(vals)
    span = (hi - lo) or abs(hi) or 1.0
    step = 10 ** math.floor(math.log10(span / 4))
    if span / step > 8:
        step *= 2
    lo = math.floor((lo - span * 0.12) / step) * step
    hi = math.ceil((hi + span * 0.12) / step) * step
    if min(vals) >= 0 and lo < 0:
        lo = 0.0
    return lo, hi


@register("xy_chart")
class XYChart(Component):
    spec_model = XYChartSpec

    def _label_h(self, data: XYChartSpec, ctx: RenderContext) -> int:
        if not data.quadrants or not data.quadrant_labels:
            return 0
        return 0  # captions live INSIDE the frame; no extra height

    def measure(self, data: XYChartSpec, width: int,
                ctx: RenderContext) -> int:
        natural = min(max(round(width * 0.52), inch(2.4)), inch(3.6))
        return natural

    def render(self, slide, data: XYChartSpec, bbox: BBox,
               ctx: RenderContext) -> int:
        theme = ctx.theme
        total = self.measure(data, bbox.w, ctx)
        if ctx.fill_hint and bbox.h > total:
            total = bbox.h
        if total > bbox.h:
            ctx.report.warn(
                f"xy_chart: needs {to_inch(total):.2f}in but bbox is "
                f"{to_inch(bbox.h):.2f}in; rendering anyway (may clip)")
        frame = BBox(bbox.x, bbox.y, bbox.w, total)

        bubble = any(p.size for p in data.points)
        hi = data.highlight
        if hi and not any(p.label == hi for p in data.points):
            ctx.report.warn(f"xy_chart: highlight {hi!r} matches no point; "
                            "no accent applied")
            hi = None

        # one series per point so each gets its own color + label —
        # the consulting scatter look (named entities, not a point cloud)
        if bubble:
            cd = BubbleChartData()
            for p in data.points:
                s = cd.add_series(p.label)
                s.add_data_point(p.x, p.y, p.size or 1.0)
            xl = XL_CHART_TYPE.BUBBLE
        else:
            cd = XyChartData()
            for p in data.points:
                s = cd.add_series(p.label)
                s.add_data_point(p.x, p.y)
            xl = XL_CHART_TYPE.XY_SCATTER

        gframe = slide.shapes.add_chart(
            xl, Emu(frame.x), Emu(frame.y), Emu(frame.w), Emu(frame.h), cd)
        chart = gframe.chart
        chart.has_title = False
        chart.has_legend = False
        size = max(6.0, theme.size_micro)
        chart.font.size = Pt(size)
        chart.font.name = ctx.font("body")
        chart.font.color.rgb = RGBColor.from_string(theme.color("ink"))

        for si, (series, p) in enumerate(zip(chart.series, data.points)):
            role = "accent" if (hi and p.label == hi) else \
                (_MUTE_ROLE if hi else "primary")
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor.from_string(
                theme.color(role))
            series.format.line.fill.background()
            try:
                series.data_labels.show_series_name = True
                series.data_labels.show_value = False
                series.data_labels.position = XL_LABEL_POSITION.ABOVE
                series.data_labels.font.size = Pt(size)
            except (AttributeError, ValueError):
                pass  # bubble label API varies; the chart stays valid

        # pin BOTH axis scales so the quadrant overlay's data->frame
        # mapping is exact (the benchmark-overlay convention)
        xlo, xhi = _nice_bounds([p.x for p in data.points])
        ylo, yhi = _nice_bounds([p.y for p in data.points])
        for axis, (lo, hi), label in (
                (chart.value_axis, (ylo, yhi), data.y_label),
                (getattr(chart, "category_axis", None), (xlo, xhi),
                 data.x_label)):
            if axis is None:
                continue
            try:
                axis.minimum_scale = lo
                axis.maximum_scale = hi
                axis.has_major_gridlines = True
                axis.major_gridlines.format.line.color.rgb = \
                    RGBColor.from_string(theme.color("grid"))
                axis.major_gridlines.format.line.width = Pt(0.5)
                axis.tick_labels.font.size = Pt(size)
                axis.has_title = True
                axis.axis_title.text_frame.text = label
                axis.axis_title.text_frame.paragraphs[0].runs[0] \
                    .font.size = Pt(size)
            except (AttributeError, ValueError):
                continue

        if data.quadrants:
            # dashed midlines at the DATA midpoints, mapped into the plot
            # rect via the pinned scales + plot-inset fractions
            plot = BBox(frame.x + round(frame.w * _PLOT_LEFT),
                        frame.y + round(frame.h * _PLOT_TOP),
                        round(frame.w * (1 - _PLOT_LEFT - _PLOT_RIGHT)),
                        round(frame.h * (1 - _PLOT_TOP - _PLOT_BOTTOM)))
            fx = ((xlo + xhi) / 2 - xlo) / (xhi - xlo)
            fy = ((ylo + yhi) / 2 - ylo) / (yhi - ylo)
            mx = plot.x + round(plot.w * fx)
            my = plot.y + round(plot.h * (1 - fy))
            add_vline(slide, mx, plot.y, plot.h, theme,
                      role="ink_muted", weight_pt=1.0, dash="dash")
            add_hline(slide, plot.x, my, plot.w, theme,
                      role="ink_muted", weight_pt=1.0, dash="dash")
            if data.quadrant_labels:
                inset = round(plot.w * _QUAD_LABEL_INSET)
                w = round(plot.w * 0.30)
                line_h = ctx.measurer.line_height_emu(
                    [Span("Ag", bold=True)], ctx.font("body"),
                    ctx.size("micro"))
                spots = (  # TL, TR, BL, BR
                    (plot.x + inset, plot.y + inset, "left"),
                    (plot.right - w - inset, plot.y + inset, "right"),
                    (plot.x + inset, plot.bottom - line_h - inset, "left"),
                    (plot.right - w - inset, plot.bottom - line_h - inset,
                     "right"))
                for (x, y, align), text in zip(spots,
                                               data.quadrant_labels):
                    box = add_text_box(slide, BBox(x, y, w, line_h),
                                       align=align)
                    write_spans_paragraph(
                        box.text_frame,
                        [Span(text, bold=True, color_role="ink_muted",
                              caps=True, spc_pts=0.6)],
                        ctx.size("micro"), theme, family=ctx.font("body"),
                        align=align, line_spacing_pt=to_pt(line_h),
                        default_color_role="ink_muted")
        return total
