"""image_block — full-width image frame with optional muted caption.

Reference look: the hero photograph band of a consulting deck — a fixed-
height frame spanning the whole content column, the image either aspect-fit
inside it ("contain", centered both axes) or center-cropped to fill it
edge-to-edge ("cover"), with a small centered ink_muted caption underneath.

measure(): frame height comes ONLY from the spec (inch(height_in)) plus the
caption fit — never from the image file, so parity and determinism hold even
when the asset is missing. render(): resolve_asset() miss draws a themed
surface_alt placeholder and warns instead of raising. Cover crops are baked
into cached PNGs under assets/cache_crop (no pptx srcRect surgery), keyed by
crop pixel size so the deck stays deterministic and edit-safe.
"""
from __future__ import annotations

from pathlib import Path

from pptx.util import Emu

from ..core.assets import ASSETS_DIR, image_size, resolve_asset
from ..core.bbox import BBox
from ..core.fit_text import FitResult, Span, fit_text
from ..core.pptx_shapes import add_shape
from ..core.pptx_text import add_text_box, make_text_frame, write_fit_result
from ..core.units import inch
from ..schema.components import ImageBlockSpec
from ..schema.rich import parse_rich
from .base import Component, RenderContext, register

_PROBE_H = 10_000_000  # caption fit decides its own height; never truncates
_MIN_CAPTION = 6.5
_MIN_PLACEHOLDER = 6.0
_CAPTION_ROLE = "ink_muted"
_CROP_CACHE = ASSETS_DIR / "cache_crop"


@register("image_block")
class ImageBlock(Component):
    spec_model = ImageBlockSpec

    # -- shared planning (side-effect free) --------------------------------

    def _plan(self, data: ImageBlockSpec, width: int, ctx: RenderContext
              ) -> tuple[int, FitResult | None, int, int]:
        frame_h = inch(data.height_in)
        gap = ctx.theme.spacing(0.4)
        caption_fit = None
        total = frame_h
        if data.caption:
            caption_fit = fit_text(
                parse_rich(data.caption, base_color_role=_CAPTION_ROLE),
                BBox(0, 0, width, _PROBE_H), ctx.font("body"),
                max_size=ctx.size("small"), min_size=_MIN_CAPTION,
                measurer=ctx.measurer)
            total += gap + caption_fit.height_emu
        return frame_h, caption_fit, gap, total

    # -- contract -----------------------------------------------------------

    def measure(self, data: ImageBlockSpec, width: int,
                ctx: RenderContext) -> int:
        _, _, _, total = self._plan(data, width, ctx)
        return total

    def render(self, slide, data: ImageBlockSpec, bbox: BBox,
               ctx: RenderContext) -> int:
        frame_h, caption_fit, gap, total = self._plan(data, bbox.w, ctx)
        if total > bbox.h:
            ctx.report.warn(
                f"image_block: height {total} exceeds bbox height {bbox.h}")
        if caption_fit is not None and caption_fit.truncated:
            ctx.report.truncated(f"image_block caption: {data.caption[:40]!r}")
        frame = bbox.with_height(frame_h)
        path = resolve_asset(data.src)
        if path is None:
            self._placeholder(slide, frame, data.src, ctx)
        elif data.fit == "cover":
            crop = _cover_crop(path, frame.w, frame.h)
            slide.shapes.add_picture(str(crop), Emu(frame.x), Emu(frame.y),
                                     Emu(frame.w), Emu(frame.h))
        else:
            iw, ih = image_size(path)
            scale = min(frame.w / iw, frame.h / ih)
            w, h = round(iw * scale), round(ih * scale)
            slide.shapes.add_picture(str(path),
                                     Emu(frame.x + (frame.w - w) // 2),
                                     Emu(frame.y + (frame.h - h) // 2),
                                     Emu(w), Emu(h))
        if caption_fit is not None:
            box = add_text_box(
                slide, BBox(bbox.x, frame.bottom + gap, bbox.w,
                            caption_fit.height_emu), align="center")
            write_fit_result(box.text_frame, caption_fit, ctx.theme,
                             family=ctx.font("body"), align="center",
                             default_color_role=_CAPTION_ROLE)
        return total

    # -- placeholder ---------------------------------------------------------

    def _placeholder(self, slide, frame: BBox, src: str,
                     ctx: RenderContext) -> None:
        shape = add_shape(slide, frame, ctx.theme, shape="rect",
                          fill_role="surface_alt")
        tf = make_text_frame(shape, align="center", anchor="middle")
        fit = fit_text([Span(src, color_role=_CAPTION_ROLE)],
                       frame.inset(ctx.theme.spacing(0.4)), ctx.font("body"),
                       max_size=ctx.size("micro"), min_size=_MIN_PLACEHOLDER,
                       max_lines=2, measurer=ctx.measurer)
        if fit.truncated:
            ctx.report.truncated(f"image_block placeholder: {src[:40]!r}")
        write_fit_result(tf, fit, ctx.theme, family=ctx.font("body"),
                         align="center", default_color_role=_CAPTION_ROLE)
        ctx.report.warn(
            f"image_block: asset {src!r} not found; placeholder rendered")


def _cover_crop(path: Path, frame_w: int, frame_h: int) -> Path:
    """Center-crop to the frame aspect, cached by crop pixel size."""
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
        cw, ch = iw, ih
        if iw * frame_h >= ih * frame_w:   # wider than frame: trim sides
            cw = max(1, round(ih * frame_w / frame_h))
        else:                              # taller than frame: trim top/bottom
            ch = max(1, round(iw * frame_h / frame_w))
        _CROP_CACHE.mkdir(parents=True, exist_ok=True)
        out = _CROP_CACHE / f"{path.stem}_{cw}x{ch}.png"
        if out.is_file() and out.stat().st_mtime >= path.stat().st_mtime:
            return out
        left = (iw - cw) // 2
        top = (ih - ch) // 2
        im.crop((left, top, left + cw, top + ch)).save(out)
    return out
