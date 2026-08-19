"""native_chart — real editable PowerPoint charts that ARGUE.

Bar / stacked bar / line / donut / waterfall via python-pptx's chart API
(editable in PowerPoint: right-click -> Edit Data). The spec forces
rhetorical decisions: sort order, one highlighted category, an annotation.

The `style` block (ChartStyleSpec) adds the top-firm rendering variants the
corpus mining found: value labels everywhere (91-100% of elite chart
slides), one highlighted series with the rest muted, endpoint labels + a
CAGR chip on growth lines, a dashed forecast tail, a benchmark reference
line, 100% stacking, horizontal bars, and a compact mode for multi-chart
slides. Overlays (benchmark line, CAGR chip) are drawn as measured shapes
ON TOP of the chart frame — the engine's home turf — because chart XML is
the flakiest OOXML there is; raw chart-XML surgery lives in
render/xml_utils.py.
"""
from __future__ import annotations

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Emu, Pt

from ..core.bbox import BBox
from ..core.fit_text import Span, fit_text
from ..core.semantic_roles import rag_map
from ..core.pptx_shapes import add_hline, add_shape
from ..core.pptx_text import (add_text_box, make_text_frame,
                              write_fit_result, write_spans_paragraph)
from ..core.units import inch
from ..render.xml_utils import (add_point_data_labels, set_series_line_dash,
                                split_combo)
from ..schema.components import NativeChartSpec
from ..schema.rich import parse_rich
from .base import Component, RenderContext, register

_SERIES_ROLES = ("primary", "accent", "positive", "primary_dark", "ink_muted",
                 "negative")
_ANNOT_H_MULT = 1.0  # annotation strip height in theme units (when present)

# plot-area insets as fractions of the chart FRAME — empirically tuned
# against PowerPoint's default column/line plot so the benchmark overlay
# lands on the right y-pixel.
_PLOT_LEFT = 0.10
_PLOT_RIGHT = 0.03
_PLOT_TOP = 0.05
_PLOT_BOTTOM = 0.14

_MUTE_ROLE = "grid"


@register("native_chart")
class NativeChart(Component):
    spec_model = NativeChartSpec

    # -- data shaping (pure) ------------------------------------------------

    def _sorted(self, data: NativeChartSpec):
        cats = list(data.categories)
        series = [(s.name, list(s.values)) for s in data.series]
        n = min([len(cats)] + [len(v) for _, v in series])
        cats, series = cats[:n], [(nm, v[:n]) for nm, v in series]
        if (data.sort in ("desc", "asc") and series
                and data.chart_type not in ("waterfall",)):
            order = sorted(range(n), key=lambda i: series[0][1][i],
                           reverse=data.sort == "desc")
            cats = [cats[i] for i in order]
            series = [(nm, [v[i] for i in order]) for nm, v in series]
        return cats, series

    def _annot_fit(self, data: NativeChartSpec, width: int,
                   ctx: RenderContext):
        """Fit the annotation to its real strip width (<=2 lines). The strip
        height DERIVES from this fit — the old fixed height was shorter than
        one line for compact strips, so text silently overflowed the box."""
        if not data.annotation:
            return None
        sp = ctx.theme.spacing
        size = ctx.size("micro") if data.style.compact else ctx.size("small")
        cell_w = max(1, width - sp(0.35) - sp(0.4))
        return fit_text([Span(data.annotation, bold=True)],
                        BBox(0, 0, cell_w, 10_000_000), ctx.font("body"),
                        max_size=size, min_size=6.5, max_lines=2,
                        measurer=ctx.measurer)

    def _annot_h(self, data: NativeChartSpec, width: int,
                 ctx: RenderContext) -> int:
        fit = self._annot_fit(data, width, ctx)
        if fit is None:
            return 0
        return fit.height_emu + ctx.theme.spacing(0.4) * 2

    def _source_fit(self, data: NativeChartSpec, width: int,
                    ctx: RenderContext):
        """Per-chart provenance micro line (annotation = rhetoric, source =
        provenance). Single line, ellipsized if needed."""
        if not data.source:
            return None
        return fit_text(parse_rich(data.source, base_color_role="ink_muted"),
                        BBox(0, 0, max(1, width), 10_000_000),
                        ctx.font("body"), max_size=ctx.size("micro"),
                        min_size=6.0, max_lines=1, measurer=ctx.measurer)

    def _source_h(self, data: NativeChartSpec, width: int,
                  ctx: RenderContext) -> int:
        fit = self._source_fit(data, width, ctx)
        if fit is None:
            return 0
        return fit.height_emu + ctx.theme.spacing(0.25)

    def measure(self, data: NativeChartSpec, width: int,
                ctx: RenderContext) -> int:
        natural = min(max(round(width * 0.48), inch(2.2)), inch(3.4))
        return (natural + self._annot_h(data, width, ctx)
                + self._source_h(data, width, ctx))

    # -- render --------------------------------------------------------------

    def render(self, slide, data: NativeChartSpec, bbox: BBox,
               ctx: RenderContext) -> int:
        theme = ctx.theme
        st = data.style
        annot_h = self._annot_h(data, bbox.w, ctx)
        source_h = self._source_h(data, bbox.w, ctx)
        total = self.measure(data, bbox.w, ctx)
        if ctx.fill_hint:
            if bbox.h > total:
                total = bbox.h  # charts fill their zone gracefully
            else:
                # ...and clamp DOWN too: charts scale, so a compressed
                # chart in its box strictly beats painting 1-2in over the
                # neighbour below (the canvas free-placement overlap class)
                total = max(bbox.h, inch(1.6) + annot_h + source_h)
        chart_h = max(inch(1.6), total - annot_h - source_h)
        frame = BBox(bbox.x, bbox.y, bbox.w, chart_h)

        cats, series = self._sorted(data)
        horizontal = st.direction == "horizontal" and data.chart_type in (
            "bar", "stacked_bar")
        if horizontal:
            # PowerPoint lists horizontal-bar categories bottom-up; reverse so
            # the reading order stays top-first.
            cats = list(reversed(cats))
            series = [(nm, list(reversed(v))) for nm, v in series]

        if data.chart_type == "waterfall":
            self._render_waterfall(slide, data, frame, cats, series, ctx)
        else:
            self._render_standard(slide, data, frame, cats, series,
                                   horizontal, ctx)

        # overlays (drawn above the chart frame)
        if st.benchmark is not None and data.chart_type in (
                "bar", "stacked_bar", "line"):
            self._draw_benchmark(slide, data, frame, cats, series, ctx)
        if st.cagr_chip:
            self._draw_cagr_chip(slide, data, frame, series, ctx)

        if data.annotation:
            # same fit that sized annot_h — measured wrap+shrink, never one
            # over-wide line for PowerPoint to wrap on its own terms
            fit = self._annot_fit(data, bbox.w, ctx)
            if fit.truncated:
                ctx.report.truncated(
                    f"chart annotation: {data.annotation[:40]!r}")
            strip = BBox(bbox.x, bbox.y + chart_h + theme.spacing(0.4),
                         bbox.w, max(1, annot_h - theme.spacing(0.4) * 2))
            bar, text = strip.take_left(theme.spacing(0.35))
            add_shape(slide, bar, theme, fill_role="accent")
            box = add_text_box(slide, text.inset(left=theme.spacing(0.4)),
                               anchor="middle")
            write_fit_result(box.text_frame, fit, theme,
                             family=ctx.font("body"),
                             default_color_role="ink")

        if data.source:
            sfit = self._source_fit(data, bbox.w, ctx)
            if sfit.truncated:
                ctx.report.truncated(f"chart source: {data.source[:40]!r}")
            sy = bbox.y + chart_h + annot_h + theme.spacing(0.25)
            sbox = add_text_box(slide, BBox(bbox.x, sy, bbox.w,
                                            max(1, source_h
                                                - theme.spacing(0.25))))
            write_fit_result(sbox.text_frame, sfit, theme,
                             family=ctx.font("body"),
                             default_color_role="ink_muted")
        return total

    # -- standard chart types -----------------------------------------------

    def _xl_type(self, data: NativeChartSpec, horizontal: bool):
        if data.chart_type == "line":
            return XL_CHART_TYPE.LINE
        if data.chart_type == "donut":
            return XL_CHART_TYPE.DOUGHNUT
        if data.chart_type == "stacked_bar":
            if data.style.percent_100:
                return (XL_CHART_TYPE.BAR_STACKED_100 if horizontal
                        else XL_CHART_TYPE.COLUMN_STACKED_100)
            return (XL_CHART_TYPE.BAR_STACKED if horizontal
                    else XL_CHART_TYPE.COLUMN_STACKED)
        if data.chart_type == "combo":
            # built as clustered columns, then split_combo moves the line
            # series into its own plot with a secondary axis
            return XL_CHART_TYPE.COLUMN_CLUSTERED
        return (XL_CHART_TYPE.BAR_CLUSTERED if horizontal
                else XL_CHART_TYPE.COLUMN_CLUSTERED)

    def _render_standard(self, slide, data, frame, cats, series, horizontal,
                         ctx):
        st = data.style
        # forecast split: single-series line -> solid + dashed halves
        forecast = (st.forecast_from and data.chart_type == "line"
                    and len(series) == 1 and st.forecast_from in cats)
        if st.forecast_from and not forecast:
            ctx.report.warn(
                f"native_chart: forecast_from {st.forecast_from!r} ignored "
                "(needs a single-series line whose category exists)")

        cd = CategoryChartData()
        cd.categories = cats
        if forecast:
            name, vals = series[0]
            idx = cats.index(st.forecast_from)
            solid = [v if i <= idx else None for i, v in enumerate(vals)]
            dashed = [v if i >= idx else None for i, v in enumerate(vals)]
            cd.add_series(name, solid)
            cd.add_series(name + " (forecast)", dashed)
        else:
            for name, values in series:
                cd.add_series(name, values)

        chart = slide.shapes.add_chart(
            self._xl_type(data, horizontal), Emu(frame.x), Emu(frame.y),
            Emu(frame.w), Emu(frame.h), cd).chart
        combo_lines: list[str] = []
        if data.chart_type == "combo" and len(series) >= 2:
            want = st.combo_line_series or series[-1][0]
            if not any(nm == want for nm, _ in series):
                ctx.report.warn(
                    f"native_chart: combo_line_series {want!r} not found; "
                    "using the last series")
                want = series[-1][0]
            combo_lines = split_combo(chart, [want])
        elif data.chart_type == "combo":
            ctx.report.warn("native_chart: combo needs >=2 series; "
                            "rendering as columns")
        self._style(chart, data, cats, series, forecast, horizontal, ctx,
                    combo_lines=combo_lines)

    def _style(self, chart, data, cats, series, forecast, horizontal, ctx,
               combo_lines: list[str] = ()):
        theme = ctx.theme
        st = data.style
        chart.has_title = False
        size = theme.size_micro - 1.5 if st.compact else theme.size_micro
        size = max(6.0, size)
        chart.font.size = Pt(size)
        chart.font.name = ctx.font("body")
        chart.font.color.rgb = RGBColor.from_string(theme.color("ink"))

        multi = len(series) > 1
        # legend: off for compact, forecast split (2 fake series) or single
        show_legend = multi and not st.compact and not forecast
        chart.has_legend = show_legend
        if show_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(size)

        highlight_ser = (st.highlight_series
                         if st.highlight_series
                         and any(nm == st.highlight_series for nm, _ in series)
                         else None)
        if st.highlight_series and highlight_ser is None:
            ctx.report.warn(
                f"native_chart: highlight_series {st.highlight_series!r} "
                "not found; no muting applied")

        plot_names = ([series[0][0], series[0][0] + " (forecast)"]
                      if forecast else [nm for nm, _ in series])
        # semantic colour by series NAME (red=stress/green=healthy) — only
        # when >=2 names match >=2 distinct buckets; explicit highlight wins
        series_rag = (rag_map([nm for nm, _ in series])
                      if len(series) > 1 and highlight_ser is None else None)
        for si, plot_series in enumerate(chart.series):
            if combo_lines:
                # after split_combo the plot order shifted: read each
                # series' real name from its XML instead of by position
                tx = plot_series._element.find(qn("c:tx"))
                nm = "".join(v.text or "" for v in tx.iter(qn("c:v"))) \
                    if tx is not None else ""
            else:
                nm = plot_names[si] if si < len(plot_names) else ""
            base_nm = nm.replace(" (forecast)", "")
            if nm in combo_lines:
                # the combo's line series: accent line on the secondary axis
                plot_series.format.line.color.rgb = RGBColor.from_string(
                    theme.color("accent"))
                plot_series.format.line.width = Pt(2.25)
                continue
            if highlight_ser is not None:
                role = "accent" if base_nm == highlight_ser else _MUTE_ROLE
            elif series_rag and base_nm in series_rag:
                role = series_rag[base_nm]
            elif combo_lines:
                role = "primary"   # combo columns stay calm under the line
            else:
                role = _SERIES_ROLES[si % len(_SERIES_ROLES)]
            hexc = RGBColor.from_string(theme.color(role))
            if data.chart_type == "line":
                plot_series.format.line.color.rgb = hexc
                plot_series.format.line.width = Pt(2.25)
                plot_series.smooth = False
                if forecast and si == 1:
                    set_series_line_dash(plot_series, "dash")
            else:
                plot_series.format.fill.solid()
                plot_series.format.fill.fore_color.rgb = hexc

        self._point_colors(chart, data, cats, series, ctx)
        self._value_labels(chart, data, series, forecast, ctx)
        self._axes(chart, data, ctx)

    def _point_colors(self, chart, data, cats, series, ctx):
        """Per-point emphasis: highlight one category, diverging by sign,
        donut per-slice — single-series only."""
        theme = ctx.theme
        if len(series) != 1:
            return
        pts = chart.series[0].points
        hi_idx = cats.index(data.highlight) if data.highlight in cats else -1
        if data.chart_type in ("bar",) and any(
                v < 0 for v in series[0][1]):
            # diverging: color by sign (highlight still wins its point)
            for i, v in enumerate(series[0][1]):
                if i == hi_idx:
                    role = "accent"
                else:
                    role = "positive" if v >= 0 else "negative"
                pts[i].format.fill.solid()
                pts[i].format.fill.fore_color.rgb = RGBColor.from_string(
                    theme.color(role))
            return
        if data.chart_type == "bar":
            # category RAG: labels like 'Overdrawn'/'Moderate'/'Healthy'
            # colour by MEANING (>=2 labels across >=2 buckets; the
            # explicit highlight keeps its accent point)
            cat_rag = rag_map(list(cats))
            if cat_rag:
                for i, cat in enumerate(cats):
                    role = "accent" if i == hi_idx else cat_rag.get(cat)
                    if role is None:
                        continue
                    pts[i].format.fill.solid()
                    pts[i].format.fill.fore_color.rgb = RGBColor.from_string(
                        theme.color(role))
                return
        if data.chart_type in ("bar", "donut") and hi_idx >= 0:
            pts[hi_idx].format.fill.solid()
            pts[hi_idx].format.fill.fore_color.rgb = RGBColor.from_string(
                theme.color("accent"))
        if data.chart_type == "donut":
            for pi, point in enumerate(pts):
                if pi == hi_idx:
                    continue
                role = _SERIES_ROLES[pi % len(_SERIES_ROLES)]
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = RGBColor.from_string(
                    theme.color(role))

    def _num_format(self, suffix: str) -> str:
        return f'#,##0.#"{suffix}"' if suffix else '#,##0.#'

    def _value_labels(self, chart, data, series, forecast, ctx):
        theme = ctx.theme
        st = data.style
        multi = len(series) > 1
        # resolve auto: single-series bar keeps the historic default ON
        want = st.value_labels
        if want is None:
            want = data.chart_type == "bar" and not multi
        if st.endpoint_labels and data.chart_type == "line":
            # per-point first+last on each plotted series
            fmt = self._num_format(data.value_suffix)
            csize = max(6.0, theme.size_micro - (1.5 if st.compact else 0))
            for ser in chart.series:
                n = len(list(ser.values))
                idxs = sorted({0, n - 1}) if n else []
                if idxs:
                    add_point_data_labels(
                        ser, idxs, fmt, csize, theme.color("ink"))
            return
        if not want:
            return
        fmt = f'0"{data.value_suffix}"' if (
            data.chart_type == "bar" and st.value_labels is None
        ) else self._num_format(data.value_suffix)
        for plot in chart.plots:
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.font.size = Pt(max(6.0, theme.size_micro
                                  - (1.5 if st.compact else 0)))
            dl.font.bold = True
            dl.font.color.rgb = RGBColor.from_string(theme.color("ink"))
            dl.number_format = fmt
            dl.number_format_is_linked = False

    def _axes(self, chart, data, ctx):
        theme = ctx.theme
        if data.chart_type not in ("bar", "stacked_bar", "line"):
            return
        va = chart.value_axis
        va.has_major_gridlines = True
        va.major_gridlines.format.line.color.rgb = RGBColor.from_string(
            theme.color("grid"))
        va.major_gridlines.format.line.width = Pt(0.5)
        va.format.line.fill.background()
        va.tick_labels.font.size = Pt(max(6.0, theme.size_micro))
        # pin the scale when a benchmark overlay needs an exact mapping
        if data.style.benchmark is not None:
            lo, hi = self._axis_bounds(data)
            va.minimum_scale = lo
            va.maximum_scale = hi
        elif data.chart_type in ("bar", "stacked_bar") and all(
                v >= 0 for s in data.series for v in s.values):
            # truncated bar axes lie about magnitude: pin zero explicitly
            # so no viewer 'optimises' the scale
            va.minimum_scale = 0.0
        ca = chart.category_axis
        ca.format.line.color.rgb = RGBColor.from_string(theme.color("grid"))
        ca.tick_labels.font.size = Pt(max(6.0, theme.size_micro))

    # -- axis bounds + overlays ---------------------------------------------

    def _axis_bounds(self, data) -> tuple[float, float]:
        vals = [v for _, s in [(s.name, s.values) for s in data.series]
                for v in s]
        dmax = max(vals) if vals else 1.0
        dmin = min(vals) if vals else 0.0
        if data.style.benchmark is not None:
            dmax = max(dmax, data.style.benchmark.value)
            dmin = min(dmin, data.style.benchmark.value)
        lo = min(0.0, dmin * 1.1)
        hi = self._nice_ceiling(max(dmax, 0.0) * 1.05)
        if hi <= lo:
            hi = lo + 1.0
        return lo, hi

    def _nice_ceiling(self, v: float) -> float:
        import math
        if v <= 0:
            return 1.0
        exp = math.floor(math.log10(v))
        base = 10 ** exp
        for m in (1, 2, 2.5, 5, 10):
            if m * base >= v:
                return m * base
        return 10 * base

    def _draw_benchmark(self, slide, data, frame, cats, series, ctx):
        theme = ctx.theme
        bm = data.style.benchmark
        lo, hi = self._axis_bounds(data)
        plot_x = frame.x + round(frame.w * _PLOT_LEFT)
        plot_w = round(frame.w * (1 - _PLOT_LEFT - _PLOT_RIGHT))
        plot_y = frame.y + round(frame.h * _PLOT_TOP)
        plot_h = round(frame.h * (1 - _PLOT_TOP - _PLOT_BOTTOM))
        frac = (bm.value - lo) / (hi - lo) if hi > lo else 0.0
        frac = min(1.0, max(0.0, frac))
        y = plot_y + round(plot_h * (1 - frac))
        add_hline(slide, plot_x, y, plot_w, theme, role="ink_muted",
                  weight_pt=1.25, dash="dash")
        lbl = f"{bm.label} {bm.value:g}{data.value_suffix}"
        box = add_text_box(slide, BBox(plot_x, y - theme.spacing(1.8),
                                       plot_w, theme.spacing(1.6)),
                           align="right", anchor="bottom")
        write_spans_paragraph(box.text_frame, [Span(lbl, bold=True)],
                              ctx.size("micro"), theme, family=ctx.font("body"),
                              align="right", default_color_role="ink_muted")

    def _draw_cagr_chip(self, slide, data, frame, series, ctx):
        theme = ctx.theme
        if not series:
            return
        vals = series[0][1]
        if len(vals) < 2 or vals[0] <= 0:
            ctx.report.warn("native_chart: cagr_chip needs a positive first "
                            "value and >=2 periods; skipped")
            return
        periods = len(vals) - 1
        cagr = ((vals[-1] / vals[0]) ** (1 / periods) - 1) * 100
        text = f"CAGR {cagr:+.1f}%"
        w = theme.spacing(9)
        h = theme.spacing(3.2)
        chip = BBox(frame.right - w - theme.spacing(1.5),
                    frame.y + theme.spacing(1.2), w, h)
        s = add_shape(slide, chip, theme, shape="rounded",
                      fill_role="surface_alt", corner_radius=0.3)
        tf = make_text_frame(s, align="center", anchor="middle")
        write_spans_paragraph(tf, [Span(text, bold=True)], ctx.size("small"),
                              theme, family=ctx.font("body"), align="center",
                              default_color_role="ink")

    # -- waterfall -----------------------------------------------------------

    def _render_waterfall(self, slide, data, frame, cats, series, ctx):
        theme = ctx.theme
        if data.sort != "none":
            ctx.report.warn("native_chart: waterfall ignores sort (order is "
                            "the bridge sequence)")
        vals = series[0][1] if series else []
        n = len(vals)
        # compute base (invisible) and visible height + a color role per bar
        bases: list[float] = []
        visibles: list[float] = []
        roles: list[str] = []
        cum = 0.0
        for i, v in enumerate(vals):
            first, last = i == 0, i == n - 1
            if first:
                bases.append(0.0)
                visibles.append(v)
                roles.append("primary_dark")
                cum = v
            elif last:
                bases.append(0.0)
                visibles.append(v)
                roles.append("primary_dark")
            else:
                after = cum + v
                bases.append(min(cum, after))
                visibles.append(abs(v))
                roles.append("positive" if v >= 0 else "negative")
                cum = after
        # arithmetic check: last should equal the running total of deltas
        if n >= 2:
            reconstructed = sum(vals[:-1])
            if abs(vals[-1] - reconstructed) > 0.02 * max(1.0, abs(vals[-1])):
                ctx.report.warn(
                    f"native_chart: waterfall does not balance "
                    f"(sum of steps {reconstructed:g} vs closing {vals[-1]:g})")

        cd = CategoryChartData()
        cd.categories = cats
        cd.add_series("base", bases)
        cd.add_series("value", visibles)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_STACKED, Emu(frame.x), Emu(frame.y),
            Emu(frame.w), Emu(frame.h), cd).chart
        chart.has_title = False
        chart.has_legend = False
        size = max(6.0, theme.size_micro - (1.5 if data.style.compact else 0))
        chart.font.size = Pt(size)
        chart.font.name = ctx.font("body")
        chart.font.color.rgb = RGBColor.from_string(theme.color("ink"))
        base_ser, val_ser = chart.series[0], chart.series[1]
        base_ser.format.fill.background()      # invisible pedestal
        base_ser.format.line.fill.background()
        for i, point in enumerate(val_ser.points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = RGBColor.from_string(
                theme.color(roles[i]))
        # signed value labels on the visible series (original deltas)
        fmt = self._num_format(data.value_suffix)
        add_point_data_labels(val_ser, list(range(n)), fmt, size,
                              theme.color("ink"))
        self._axes(chart, data, ctx)
