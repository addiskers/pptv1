"""donut_stat — donut ring showing value_pct filled, a stat in the hole, label.

The ring is THREE stacked shapes: a full 'surface_alt' oval (track), a
'primary' PIE wedge spanning value_pct (starting at 12 o'clock, sweeping
clockwise), and a 'bg' oval at 55% diameter that punches the hole.
center_text sits middle-anchored over the hole; the label fits in the width
right of the ring, or below it when the bbox is narrow.

PIE adjustments (verified empirically against the written XML): python-pptx
normalizes the raw OOXML adjustment value by /100000, while preset-geometry
angles are stored in 60000ths of a degree — so ONE DEGREE == 0.6 adjustment
units. _add_pie converts; naive `adjustments[0] = -90` would land at -150°.
value_pct >= 99.5 renders as a full 'primary' oval (a ~360° pie degenerates).
"""
from __future__ import annotations

from dataclasses import dataclass

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu

from ..core.bbox import BBox
from ..core.fit_text import FitResult, Span, fit_text
from ..core.pptx_shapes import add_shape
from ..core.pptx_text import add_text_box, write_fit_result
from ..core.theme import Theme
from ..core.units import inch
from ..schema.components import DonutStatSpec
from ..schema.rich import parse_rich
from .base import Component, RenderContext, register

_RING_W_FRAC = 0.4       # ring outer diameter as a share of bbox width…
_RING_MAX_D = inch(1.1)  # …capped at 1.1in
_HOLE_FRAC = 0.55        # hole diameter as a share of the ring diameter
_SIDE_MIN_W = inch(0.7)  # min width beside the ring for a side label
_GAP_MULT = 0.4          # spacing multiple between ring and side label
_PAD_MULT = 0.2          # trailing spacing multiple (see height contract)
_START_DEG = -90.0       # 12 o'clock; the wedge sweeps clockwise
_DEG_PER_PCT = 3.6
_ADJ_PER_DEG = 0.6       # python-pptx adjustment units per degree (see above)
_FULL_PCT = 99.5         # >= this: full 'primary' oval instead of a pie
_MIN_CENTER_PT = 7.5
_MIN_LABEL_PT = 7.5
_PROBE_H = inch(11)      # generous probe: fit decides its own natural height


def _add_pie(slide, bbox: BBox, theme: Theme, *, fill_role: str,
             start_deg: float, end_deg: float):
    """Themed pie wedge. Fill/line resolve through theme roles only."""
    # raw shape: PIE not in core.pptx_shapes (its SHAPES dict is frozen)
    s = slide.shapes.add_shape(MSO_SHAPE.PIE, Emu(bbox.x), Emu(bbox.y),
                               Emu(bbox.w), Emu(bbox.h))
    s.adjustments[0] = start_deg * _ADJ_PER_DEG
    s.adjustments[1] = end_deg * _ADJ_PER_DEG
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor.from_string(theme.color(fill_role))
    s.line.fill.background()
    s.shadow.inherit = False
    return s


@dataclass(frozen=True)
class _Layout:
    d: int                # ring outer diameter
    side: bool            # True: label beside the ring; False: below
    label_w: int
    label_fit: FitResult
    center_fit: FitResult
    height: int           # total consumed height


@register("donut_stat")
class DonutStat(Component):
    spec_model = DonutStatSpec

    # -- pure layout math (no side effects) --------------------------------

    def _layout(self, data: DonutStatSpec, width: int,
                ctx: RenderContext) -> _Layout:
        ink = "inverse_ink" if data.inverse else "ink"
        d = min(round(width * _RING_W_FRAC), _RING_MAX_D)
        gap = ctx.theme.spacing(_GAP_MULT)
        pad = ctx.theme.spacing(_PAD_MULT)
        side_w = width - d - gap
        side = side_w >= _SIDE_MIN_W
        label_w = side_w if side else width
        label_fit = fit_text(parse_rich(data.label, base_color_role=ink),
                             BBox(0, 0, max(1, label_w), _PROBE_H),
                             ctx.font("body"), max_size=ctx.size("small"),
                             min_size=_MIN_LABEL_PT, measurer=ctx.measurer)
        hole_w = round(d * _HOLE_FRAC)
        center_fit = fit_text([Span(data.center_text, bold=True,
                                    color_role=ink)],
                              BBox(0, 0, max(1, hole_w), d), ctx.font("body"),
                              max_size=ctx.size("stat"),
                              min_size=_MIN_CENTER_PT, max_lines=1,
                              measurer=ctx.measurer)
        if side:
            height = max(d, label_fit.height_emu) + pad
        else:
            height = d + pad + label_fit.height_emu
        return _Layout(d, side, label_w, label_fit, center_fit, height)

    # -- contract -----------------------------------------------------------

    def measure(self, data: DonutStatSpec, width: int,
                ctx: RenderContext) -> int:
        return self._layout(data, width, ctx).height

    def render(self, slide, data: DonutStatSpec, bbox: BBox,
               ctx: RenderContext) -> int:
        theme = ctx.theme
        lay = self._layout(data, bbox.w, ctx)
        gap = theme.spacing(_GAP_MULT)
        pad = theme.spacing(_PAD_MULT)

        if lay.side:
            content_h = max(lay.d, lay.label_fit.height_emu)
            ring = BBox(bbox.x, bbox.y + (content_h - lay.d) // 2,
                        lay.d, lay.d)
            label_box = BBox(bbox.x + lay.d + gap, bbox.y,
                             lay.label_w, content_h)
            label_align, label_anchor = "left", "middle"
        else:
            ring = BBox(bbox.x + (bbox.w - lay.d) // 2, bbox.y, lay.d, lay.d)
            label_box = BBox(bbox.x, bbox.y + lay.d + pad,
                             lay.label_w, lay.label_fit.height_emu)
            label_align, label_anchor = "center", "top"

        # (1) track, (2) primary wedge, (3) hole — stacked in z-order
        add_shape(slide, ring, theme, shape="oval", fill_role="surface_alt")
        if data.value_pct >= _FULL_PCT:
            add_shape(slide, ring, theme, shape="oval", fill_role="primary")
        elif data.value_pct > 0:
            _add_pie(slide, ring, theme, fill_role="primary",
                     start_deg=_START_DEG,
                     end_deg=_START_DEG + data.value_pct * _DEG_PER_PCT)
        hole_d = round(lay.d * _HOLE_FRAC)
        hole = BBox(ring.x + (lay.d - hole_d) // 2,
                    ring.y + (lay.d - hole_d) // 2, hole_d, hole_d)
        add_shape(slide, hole, theme, shape="oval",
                  fill_role=data.hole_fill_role)

        # center stat over the hole
        ink = "inverse_ink" if data.inverse else "ink"
        if lay.center_fit.truncated:
            ctx.report.truncated(f"donut_stat center: {data.center_text!r}")
        cbox = add_text_box(slide, ring, align="center", anchor="middle")
        write_fit_result(cbox.text_frame, lay.center_fit, theme,
                         family=ctx.font("body"), align="center",
                         default_color_role=ink)

        # label
        if lay.label_fit.truncated:
            ctx.report.truncated(f"donut_stat label: {data.label[:40]!r}")
        lbox = add_text_box(slide, label_box, align=label_align,
                            anchor=label_anchor)
        write_fit_result(lbox.text_frame, lay.label_fit, theme,
                         family=ctx.font("body"), align=label_align,
                         default_color_role=ink)

        if lay.height > bbox.h:
            ctx.report.warn("donut_stat: content taller than provided bbox")
        return lay.height
