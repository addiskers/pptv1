"""venn — 2-3 overlapping translucent circles.

Census: the overlap/synergy/sweet-spot form. Circles fill primary /
accent / positive at ~45% opacity so intersections visibly blend
(python-pptx has no alpha API — we write <a:alpha> into the solid fill).
Circle labels sit in each circle's OUTER third; the optional
intersection label takes the shared center with a bold ink chip so it
stays legible over the blend.

measure(): pure function of width (circle diameter + optional label
headroom) shared with render() via _plan.
"""
from __future__ import annotations

from dataclasses import dataclass

from pptx.oxml.ns import qn

from ..core.bbox import BBox
from ..core.fit_text import Span, fit_text
from ..core.pptx_shapes import add_shape
from ..core.pptx_text import add_text_box, write_fit_result
from ..core.units import inch, to_inch
from ..schema.components import VennSpec
from ..schema.rich import parse_rich
from .base import Component, RenderContext, register

_D_FRAC = 0.40            # circle diameter as a share of component width
_D_MAX = inch(2.9)
_OVERLAP = 0.34           # center offset = D * (1 - overlap) / ... see _plan
_ALPHA_PCT = 45           # fill opacity per circle
_MIN_LABEL = 7.0
_PROBE_H = 10_000_000
_FILLS = ("primary", "accent", "positive")


def _set_fill_alpha(shape, pct: int) -> None:
    """Append <a:alpha> to the shape's solid fill (pct = opacity %)."""
    spPr = shape.fill._xPr
    solid = spPr.find(qn("a:solidFill"))
    if solid is None:
        return
    clr = solid.find(qn("a:srgbClr"))
    if clr is None:
        return
    alpha = clr.makeelement(qn("a:alpha"), {"val": str(pct * 1000)})
    clr.append(alpha)


@dataclass(frozen=True)
class _Plan:
    d: int
    total: int
    centers: tuple[tuple[int, int], ...]   # offsets from (cx, top)


@register("venn")
class Venn(Component):
    spec_model = VennSpec

    def _plan(self, data: VennSpec, width: int) -> _Plan:
        d = min(round(width * _D_FRAC), _D_MAX)
        off = round(d * (1 - _OVERLAP) / 2 + d * _OVERLAP / 2 * 0.6)
        r = d // 2
        if len(data.circles) == 2:
            total = d
            centers = ((-off, r), (off, r))
        else:
            drop = round(d * 0.52)
            total = d + drop
            centers = ((-off, r), (off, r), (0, r + drop))
        return _Plan(d, total, centers)

    def measure(self, data: VennSpec, width: int,
                ctx: RenderContext) -> int:
        return self._plan(data, width).total

    def render(self, slide, data: VennSpec, bbox: BBox,
               ctx: RenderContext) -> int:
        theme = ctx.theme
        plan = self._plan(data, bbox.w)
        if plan.total > bbox.h:
            ctx.report.warn(
                f"venn: needs {to_inch(plan.total):.2f}in but bbox is "
                f"{to_inch(bbox.h):.2f}in; rendering anyway (may clip)")
        family = ctx.font("body")
        n = len(data.circles)
        cx = bbox.x + bbox.w // 2
        r = plan.d // 2

        for i, (dx, dy) in enumerate(plan.centers):
            circle = BBox(cx + dx - r, bbox.y + dy - r, plan.d, plan.d)
            s = add_shape(slide, circle, theme, shape="oval",
                          fill_role=_FILLS[i])
            _set_fill_alpha(s, _ALPHA_PCT)

        # labels in each circle's outer third (away from the overlap)
        label_w = round(plan.d * 0.52)
        for i, ((dx, dy), label) in enumerate(zip(plan.centers,
                                                  data.circles)):
            fit = fit_text(
                [Span(sp.text, bold=True, italic=sp.italic,
                      color_role="ink")
                 for sp in parse_rich(label, base_color_role="ink")],
                BBox(0, 0, label_w, _PROBE_H), family,
                max_size=ctx.size("small"), min_size=_MIN_LABEL,
                max_lines=2, measurer=ctx.measurer)
            if fit.truncated:
                ctx.report.truncated(f"venn: {label[:40]!r}")
            ccx, ccy = cx + dx, bbox.y + dy
            if n == 2 or i < 2:
                # left/right circles: label pushed outward horizontally
                lx = ccx - r + round(plan.d * 0.06) if dx < 0 else \
                    ccx + r - label_w - round(plan.d * 0.06)
                ly = ccy - fit.height_emu // 2
            else:
                # bottom circle: label in its lower half
                lx = ccx - label_w // 2
                ly = ccy + round(r * 0.35)
            box = add_text_box(slide, BBox(lx, ly, label_w,
                                           fit.height_emu),
                               align="center")
            write_fit_result(box.text_frame, fit, theme, family=family,
                             align="center", default_color_role="ink")

        if data.intersection:
            # shared center: mean of the circle centers
            mx = sum(dx for dx, _ in plan.centers) // n
            my = sum(dy for _, dy in plan.centers) // n
            ifit = fit_text(
                [Span(sp.text, bold=True, color_role="ink")
                 for sp in parse_rich(data.intersection,
                                      base_color_role="ink")],
                BBox(0, 0, round(plan.d * 0.5), _PROBE_H), family,
                max_size=ctx.size("small"), min_size=_MIN_LABEL,
                max_lines=2, measurer=ctx.measurer)
            if ifit.truncated:
                ctx.report.truncated(
                    f"venn intersection: {data.intersection[:40]!r}")
            pad = theme.spacing(0.25)
            chip = BBox(cx + mx - round(plan.d * 0.25) - pad,
                        bbox.y + my - ifit.height_emu // 2 - pad,
                        round(plan.d * 0.5) + 2 * pad,
                        ifit.height_emu + 2 * pad)
            cs = add_shape(slide, chip, theme, shape="rounded",
                           fill_role="bg", corner_radius=0.35,
                           line_role="grid", line_w_pt=0.75)
            from ..core.pptx_text import make_text_frame
            tf = make_text_frame(cs, align="center", anchor="middle")
            write_fit_result(tf, ifit, theme, family=family, align="center",
                             default_color_role="ink")
        return plan.total
