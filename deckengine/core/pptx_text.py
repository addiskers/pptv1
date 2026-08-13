"""The ONLY sanctioned path for putting text on a slide.

python-pptx's add_textbox defaults are hostile to measured layout:
wrap="none" + spAutoFit (shape re-grows on first click) + hidden 0.1in/0.05in
insets. Every text frame in the engine goes through make_text_frame, which pins
wrap ON, autofit OFF, margins to zero, and writes the exact line spacing that
fit_text measured as <a:spcPts> — so the value we measured is the value
PowerPoint must render.
"""
from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from .bbox import BBox
from .fit_text import FitResult, Line, Span
from .theme import Theme

ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
         "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
          "bottom": MSO_ANCHOR.BOTTOM}


def make_text_frame(shape, *, align: str = "left", anchor: str = "top"):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = ANCHOR[anchor]
    tf.paragraphs[0].alignment = ALIGN[align]
    return tf


def add_text_box(slide, bbox: BBox, *, align: str = "left", anchor: str = "top"):
    box = slide.shapes.add_textbox(Emu(bbox.x), Emu(bbox.y), Emu(bbox.w), Emu(bbox.h))
    make_text_frame(box, align=align, anchor=anchor)
    return box


def set_notes(slide, text: str) -> None:
    """Speaker notes. Notes pages are laid out by PowerPoint and exempt from
    the measured-layout contract; never called with empty text (python-pptx
    would still create the notes part)."""
    slide.notes_slide.notes_text_frame.text = text


def write_fit_result(tf, fit: FitResult, theme: Theme, *, family: str,
                     align: str = "left", default_color_role: str = "ink") -> None:
    """Write measured lines as one paragraph per line with exact point spacing.

    fit_text already wrapped the text; writing each measured line as its own
    paragraph means PowerPoint can never re-wrap differently than we measured.
    """
    first = True
    for line in fit.lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = ALIGN[align]
        p.line_spacing = Pt(fit.line_spacing_pt)
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        _write_spans(p, line.spans, fit.size_pt, theme, family, default_color_role)


def write_spans_paragraph(tf, spans: list[Span], size_pt: float, theme: Theme, *,
                          family: str, align: str = "left",
                          line_spacing_pt: float | None = None,
                          default_color_role: str = "ink",
                          paragraph=None) -> None:
    """Single unwrapped paragraph (labels, headers — content known to fit)."""
    p = paragraph if paragraph is not None else tf.paragraphs[0]
    p.alignment = ALIGN[align]
    if line_spacing_pt is not None:
        p.line_spacing = Pt(line_spacing_pt)
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    _write_spans(p, spans, size_pt, theme, family, default_color_role)


# CT_TextCharacterProperties child order: ln, (fill group), effectLst/Dag,
# highlight, uLn*, uFill*, latin, ea, cs, ... — highlight must precede latin.
_AFTER_HIGHLIGHT = ("a:uLnTx", "a:uLn", "a:uFillTx", "a:uFill", "a:latin",
                    "a:ea", "a:cs", "a:sym", "a:hlinkClick", "a:hlinkMouseOver",
                    "a:rtl", "a:extLst")


def _set_highlight(rpr, hex_color: str) -> None:
    for old in rpr.findall(qn("a:highlight")):
        rpr.remove(old)
    hl = rpr.makeelement(qn("a:highlight"), {})
    clr = hl.makeelement(qn("a:srgbClr"), {"val": hex_color})
    hl.append(clr)
    anchor = None
    for tag in _AFTER_HIGHLIGHT:
        found = rpr.find(qn(tag))
        if found is not None:
            anchor = found
            break
    if anchor is not None:
        anchor.addprevious(hl)
    else:
        rpr.append(hl)


def _write_spans(p, spans: list[Span], base_size: float, theme: Theme,
                 family: str, default_color_role: str) -> None:
    for span in spans:
        r = p.add_run()
        r.text = span.text
        f = r.font
        f.size = Pt(span.sized(base_size))
        f.bold = span.bold
        f.italic = span.italic
        f.name = span.font or family
        f.color.rgb = RGBColor.from_string(
            theme.color(span.color_role or default_color_role))
        if span.superscript or span.highlight_role:
            rpr = r._r.get_or_add_rPr()
            if span.superscript:
                # baseline in thousandths of a percent (30% raise)
                rpr.set("baseline", "30000")
            if span.highlight_role:
                _set_highlight(rpr, theme.color(span.highlight_role))
