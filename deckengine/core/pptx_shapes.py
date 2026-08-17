"""Themed shape helpers. All fills/lines resolve through Theme color roles."""
from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from .bbox import BBox
from .theme import Theme

SHAPES = {
    "rect": MSO_SHAPE.RECTANGLE,
    "rounded": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "pentagon": MSO_SHAPE.PENTAGON,
    "chevron": MSO_SHAPE.CHEVRON,
    "right_arrow": MSO_SHAPE.RIGHT_ARROW,
    "down_arrow": MSO_SHAPE.DOWN_ARROW,
    "right_brace": MSO_SHAPE.RIGHT_BRACE,
    "left_brace": MSO_SHAPE.LEFT_BRACE,
    "notched_arrow": MSO_SHAPE.NOTCHED_RIGHT_ARROW,
    "trapezoid": MSO_SHAPE.TRAPEZOID,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "diamond": MSO_SHAPE.DIAMOND,
}


def add_shape(slide, bbox: BBox, theme: Theme, *, shape: str = "rect",
              fill_role: str | None = None, fill_hex: str | None = None,
              line_role: str | None = None, line_w_pt: float = 0.75,
              line_dash: str | None = None, corner_radius: float | None = None,
              shadow: bool = False):
    s = slide.shapes.add_shape(SHAPES[shape], Emu(bbox.x), Emu(bbox.y),
                               Emu(bbox.w), Emu(bbox.h))
    if fill_hex is not None:
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    elif fill_role is not None:
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor.from_string(theme.color(fill_role))
    else:
        s.fill.background()
    if line_role is not None:
        s.line.color.rgb = RGBColor.from_string(theme.color(line_role))
        s.line.width = Pt(line_w_pt)
        if line_dash:
            _set_dash(s.line, line_dash)
    else:
        s.line.fill.background()
    if corner_radius is not None and shape == "rounded":
        try:
            s.adjustments[0] = corner_radius
        except (IndexError, ValueError):
            pass
    s.shadow.inherit = False
    if shadow:
        from ..render.xml_utils import add_soft_shadow
        add_soft_shadow(s)
    return s


def add_hline(slide, x: int, y: int, w: int, theme: Theme, *,
              role: str = "grid", weight_pt: float = 0.75, dash: str | None = None):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(x), Emu(y),
                                    Emu(x + w), Emu(y))
    ln.line.color.rgb = RGBColor.from_string(theme.color(role))
    ln.line.width = Pt(weight_pt)
    if dash:
        _set_dash(ln.line, dash)
    return ln


def add_vline(slide, x: int, y: int, h: int, theme: Theme, *,
              role: str = "grid", weight_pt: float = 0.75, dash: str | None = None):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(x), Emu(y),
                                    Emu(x), Emu(y + h))
    ln.line.color.rgb = RGBColor.from_string(theme.color(role))
    ln.line.width = Pt(weight_pt)
    if dash:
        _set_dash(ln.line, dash)
    return ln


def add_icon(slide, bbox: BBox, name: str, theme: Theme, *,
             color_role: str = "inverse_ink") -> None:
    """Place a themed monochrome icon centered in bbox (square, fit to min side)."""
    from .icons import get_icon
    side = min(bbox.w, bbox.h)
    path = get_icon(name, theme.color(color_role), px=128)
    slide.shapes.add_picture(str(path), Emu(bbox.x + (bbox.w - side) // 2),
                             Emu(bbox.y + (bbox.h - side) // 2),
                             Emu(side), Emu(side))


def _set_dash(line_format, dash: str) -> None:
    # python-pptx has no dash API on LineFormat in all versions; write prstDash.
    ln = line_format._get_or_add_ln()
    for existing in ln.findall(qn("a:prstDash")):
        ln.remove(existing)
    el = ln.makeelement(qn("a:prstDash"), {"val": dash})
    # prstDash must come after fill props but before cap/join children if present.
    ln.append(el)
