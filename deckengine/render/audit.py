"""Rendered-output audit — the alignment engine's measuring stick.

Spec-level rails catch declared geometry; this audits what actually got
PAINTED. Two checks per slide, both feeding the candidate defect score
(a clipped or ragged design must LOSE to its competitor) and the eval
harness:

- text-on-text overlap: two non-empty text frames intersecting by more
  than a sliver — catches every clip/spill case regardless of cause
  (component overflow, z collisions, freehand geometry).
- left-edge raggedness: distinct left edges (quantized) among content
  shapes. Hand-built slides share 4-8 edges; freehand LLM output showed
  18-30. High counts are the "visually not aligned" tell.
"""
from __future__ import annotations

_EDGE_STEP = 45720          # 0.05in in EMU — edges within this are "same"
_OVERLAP_MIN_FRAC = 0.15    # intersection / smaller area above this = clip
_FOOTER_BAND = 457200       # bottom 0.5in excluded (footer strip)
_FULL_WIDTH_FRAC = 0.95     # full-bleed panels don't count as edges


_EXEMPT_TAGS = ("de:tab", "de:wordmark", "de:overlay")


def _shape_info(shape):
    try:
        box = (int(shape.left), int(shape.top),
               int(shape.left) + int(shape.width),
               int(shape.top) + int(shape.height))
    except Exception:  # noqa: BLE001
        return None
    text, align = "", None
    try:
        if shape.has_text_frame:
            text = (shape.text_frame.text or "").strip()
            if text:
                align = shape.text_frame.paragraphs[0].alignment
    except Exception:  # noqa: BLE001 — exotic shapes never crash audits
        pass
    exempt = any((shape.name or "").startswith(t) for t in _EXEMPT_TAGS)
    return box, text, align, exempt


def audit_slide(slide, slide_w: int, slide_h: int) -> dict:
    """{'overlaps': [(text_a, text_b, frac)...], 'left_edges': int}

    Edges count only LEFT-aligned text shapes — centered/right content
    (badge pills, stat values, chevron labels) legitimately varies its
    left and is not a raggedness signal."""
    infos = [i for i in (_shape_info(s) for s in slide.shapes) if i]
    texted = [(b, t, al) for b, t, al, ex in infos if t and not ex]

    overlaps = []
    for i in range(len(texted)):
        for j in range(i + 1, len(texted)):
            (ax0, ay0, ax1, ay1), ta, aa = texted[i]
            (bx0, by0, bx1, by1), tb, ab = texted[j]
            iw = min(ax1, bx1) - max(ax0, bx0)
            ih = min(ay1, by1) - max(ay0, by0)
            if iw <= 0 or ih <= 0:
                continue
            a_area = max(1, (ax1 - ax0) * (ay1 - ay0))
            b_area = max(1, (bx1 - bx0) * (by1 - by0))
            frac = (iw * ih) / min(a_area, b_area)
            if frac <= _OVERLAP_MIN_FRAC:
                continue
            if frac >= 0.85 and aa != ab:
                # co-located label/value with opposing alignments (a value
                # riding a track, a right-aligned target on a band) — the
                # deliberate consulting gesture, not a clip
                continue
            overlaps.append((ta[:36], tb[:36], round(frac, 2)))

    from pptx.enum.text import PP_ALIGN
    edges = set()
    for (x0, y0, x1, y1), t, al, _ in infos:
        if not t or al in (PP_ALIGN.CENTER, PP_ALIGN.RIGHT):
            continue
        if y0 > slide_h - _FOOTER_BAND:      # footer strip
            continue
        if (x1 - x0) > slide_w * _FULL_WIDTH_FRAC:  # full-bleed panel/rule
            continue
        edges.add(x0 // _EDGE_STEP)

    return {"overlaps": overlaps, "left_edges": len(edges)}


def audit_presentation(prs) -> list[dict]:
    w, h = int(prs.slide_width), int(prs.slide_height)
    return [audit_slide(s, w, h) for s in prs.slides]
