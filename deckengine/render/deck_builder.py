"""DeckSpec -> .pptx. The single entry point for rendering."""
from __future__ import annotations

import logging
import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

from ..components.base import BuildReport, RenderContext
from ..core.bbox import BBox
from ..core.fit_text import Span, TextMeasurer
from ..core.fonts import default_registry
from ..core.pptx_shapes import add_hline
from ..core.pptx_text import add_text_box, write_spans_paragraph
from ..core.theme import load_theme
from ..core.units import SLIDE_H_16_9, SLIDE_W_16_9, inch
from ..schema.slide_types import DeckSpec
# importing these modules registers all components / assemblers
from .. import components as _components  # noqa: F401
from ..slides import base as slide_base
from .. import slides as _slides  # noqa: F401

log = logging.getLogger("deckengine")


def build_deck(spec: DeckSpec, out_path: str | Path) -> BuildReport:
    theme = load_theme(spec.theme)
    default_registry().validate_theme(theme)  # fail loudly before drawing anything
    report = BuildReport()
    ctx = RenderContext(theme=theme, measurer=TextMeasurer(), report=report)

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_16_9)
    prs.slide_height = Emu(SLIDE_H_16_9)
    blank = prs.slide_layouts[6]

    # pre-pass: split any deep-dive slide whose table exceeds one page
    from ..layout.pagination import paginate_deep_dive
    from ..llm.format_rules import enrich_slide_charts
    expanded = []
    for s in spec.slides:
        # deterministic chart enrichment: rich, content-driven chart styling
        # (endpoint labels, CAGR chip, highlight, horizontal, 100% stack) —
        # applied to every deck, respecting any explicit author style choices
        enrich_slide_charts(s)
        if s.slide_type == "data_deep_dive":
            expanded.extend(paginate_deep_dive(s, ctx))
        else:
            expanded.append(s)

    from pptx.dml.color import RGBColor
    for i, slide_spec in enumerate(expanded, start=1):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(theme.color_bg)
        if getattr(slide_spec, "bg_image", None):
            _paint_bg_image(slide, slide_spec.bg_image,
                            slide_spec.bg_image_opacity, report)
        try:
            slide_base.get_assembler(slide_spec.slide_type).assemble(
                slide, slide_spec, ctx)
        except Exception:
            log.exception("slide %d (%s) failed", i, slide_spec.slide_type)
            raise
        if slide_spec.slide_type not in ("title", "section_divider"):
            _footer(slide, spec, i, ctx)
        _logo(slide, spec.meta.logo, report)

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    if os.environ.get("DECKENGINE_EMBED_FONTS") == "1":
        # a viewer without our fonts must never substitute-and-rewrap;
        # embed whatever the registry resolved (real faces on Windows,
        # metric clones on Linux)
        from .embed_fonts import embed_fonts
        try:
            n = embed_fonts(out_path, warn=report.warn)
            log.info("embedded %d font file(s) into %s", n, out_path.name)
        except Exception as exc:  # embedding is best-effort, never fatal
            report.warn(f"font embedding failed: {exc}")
    log.info("wrote %s (%d slides, %d warnings)", out_path, len(expanded),
             len(report.warnings))
    return report


def _paint_bg_image(slide, src: str, opacity: float,
                    report: BuildReport) -> None:
    """Faint watermark, aspect-contained and centered in the body band,
    painted BEFORE content so everything renders above it."""
    from ..core.assets import faded_copy, image_size, resolve_asset
    path = resolve_asset(src)
    if path is None:
        report.warn(f"bg_image {src!r} not found under assets/; skipped")
        return
    faded = faded_copy(path, opacity)
    band = BBox(inch(0.45), inch(1.4), SLIDE_W_16_9 - inch(0.9),
                SLIDE_H_16_9 - inch(2.0))
    iw, ih = image_size(faded)
    scale = min(band.w / iw, band.h / ih)
    w, h = round(iw * scale), round(ih * scale)
    slide.shapes.add_picture(str(faded),
                             Emu(band.x + (band.w - w) // 2),
                             Emu(band.y + (band.h - h) // 2),
                             Emu(w), Emu(h))


def _logo(slide, src: str | None, report: BuildReport) -> None:
    """Brand logo, aspect-contained in a small box pinned to the top-right
    corner. Runs on EVERY slide (the cover's date moved left to clear it).
    A missing asset warns + skips — a logo never blocks a render."""
    if not src:
        return
    from ..core.assets import image_size, resolve_asset
    path = resolve_asset(src)
    if path is None:
        report.warn(f"logo {src!r} not found under assets/; skipped")
        return
    box_w, box_h = inch(1.4), inch(0.42)
    iw, ih = image_size(path)
    scale = min(box_w / iw, box_h / ih)
    w, h = round(iw * scale), round(ih * scale)
    x = SLIDE_W_16_9 - inch(0.4) - w  # right edge, matches slide margins
    slide.shapes.add_picture(str(path), Emu(x), Emu(inch(0.22)), Emu(w), Emu(h))


def _footer(slide, spec: DeckSpec, page: int, ctx: RenderContext) -> None:
    y = SLIDE_H_16_9 - inch(0.42)
    x0, x1 = inch(0.45), SLIDE_W_16_9 - inch(0.45)
    add_hline(slide, x0, y, x1 - x0, ctx.theme, role="grid", weight_pt=1.0)
    size = ctx.theme.size_micro
    ty = y + inch(0.07)
    if spec.meta.date:
        b = add_text_box(slide, BBox(x0, ty, inch(2.0), inch(0.25)))
        write_spans_paragraph(b.text_frame, [Span(spec.meta.date)], size,
                              ctx.theme, family=ctx.font("body"),
                              default_color_role="ink_muted")
    if spec.meta.confidentiality:
        b = add_text_box(slide, BBox(SLIDE_W_16_9 // 2 - inch(2), ty, inch(4),
                                     inch(0.25)), align="center")
        write_spans_paragraph(b.text_frame, [Span(spec.meta.confidentiality)],
                              size, ctx.theme, family=ctx.font("body"),
                              align="center", default_color_role="ink_muted")
    right_bits = " ".join(b for b in (spec.meta.footer_org, str(page)) if b)
    if right_bits:
        b = add_text_box(slide, BBox(x1 - inch(3), ty, inch(3), inch(0.25)),
                         align="right")
        write_spans_paragraph(b.text_frame, [Span(right_bits)], size, ctx.theme,
                              family=ctx.font("body"), align="right",
                              default_color_role="ink_muted")
