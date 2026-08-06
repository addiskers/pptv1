"""Corpus inventory: walk consulting decks, normalize, dedupe, index.

CLI:  python corpus/ingest.py [--src DIR] [--limit N] [--decks GLOB]
                              [--work DIR]

Walks a source directory (env DECKENGINE_CORPUS or --src) for
*.pptx|*.ppt|*.pdf, normalizes every slide/page to a 1280px PNG plus
metadata, dedupes via a 256-bit dHash (exact + hamming<=4 near-dupe),
and appends SlideRecord lines to work/index.jsonl — the shared inventory
consumed by the style-mining and template-induction pipelines.

Design: pure helpers (dhash, hamming, chart_types_from_blob,
shape_stats_for_slide, iter_source_files, manifest IO) are separated
from COM/PDF IO so tests can run without Office; COM calls go through
corpus.com_convert (monkeypatchable module attributes). Resume skips
decks whose (deck_id, mtime, size) already succeeded per
work/manifest.jsonl; any per-deck failure is appended to
work/ingest_errors.jsonl and the walk continues — one bad deck never
crashes the run.
"""
from __future__ import annotations

import argparse
import fnmatch
import hashlib
import json
import os
import sys
from datetime import datetime, timezone
from pathlib import Path

from lxml import etree
from PIL import Image

if __package__ in (None, ""):  # allow `python corpus/ingest.py`
    sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from corpus import com_convert
from corpus.models import DeckManifestEntry, ShapeStats, SlideRecord

SOURCE_SUFFIXES = {".pptx", ".ppt", ".pdf"}
TARGET_WIDTH = 1280
TARGET_HEIGHT = 720
NEAR_DUP_HAMMING = 4
TEXT_CAP = 2000
# Edge snapping: 0.05in grid = 45720 EMU = 3.6pt per grid step.
EDGE_GRID_EMU = 45720
EDGE_GRID_PT = 3.6
EDGE_CAP = 60


# --- pure helpers -----------------------------------------------------------

def dhash(img: Image.Image) -> str:
    """16x16 difference hash (256 bits) over the grayscale image, as hex."""
    g = img.convert("L").resize((17, 16), Image.Resampling.LANCZOS)
    px = g.tobytes()  # mode L: one byte per pixel, row-major
    bits = 0
    for row in range(16):
        for col in range(16):
            left = px[row * 17 + col]
            right = px[row * 17 + col + 1]
            bits = (bits << 1) | (1 if left < right else 0)
    return f"{bits:064x}"


def hamming(a: str, b: str) -> int:
    return bin(int(a, 16) ^ int(b, 16)).count("1")


def deck_id_for(src_root: Path, path: Path) -> str:
    rel = path.relative_to(src_root).as_posix()
    return hashlib.sha1(rel.encode("utf-8")).hexdigest()[:12]


def iter_source_files(src: Path, decks_glob: str | None = None,
                      exclude_dirs: set[str] | None = None) -> list[Path]:
    """Deck files under src, sorted; skips Office temp ~$ files and any
    file with an excluded directory name anywhere in its relative path
    (the corpus root often also holds the engine repo / preview dirs)."""
    out: list[Path] = []
    for p in sorted(src.rglob("*")):
        if not p.is_file() or p.suffix.lower() not in SOURCE_SUFFIXES:
            continue
        if p.name.startswith("~$"):
            continue
        rel_parts = p.relative_to(src).parts[:-1]
        if exclude_dirs and any(part in exclude_dirs for part in rel_parts):
            continue
        if decks_glob:
            rel = p.relative_to(src).as_posix()
            if not (fnmatch.fnmatch(rel, decks_glob)
                    or fnmatch.fnmatch(p.name, decks_glob)):
                continue
        out.append(p)
    return out


def chart_types_from_blob(blob: bytes) -> list[str]:
    """Distinct plot-type local names (barChart, lineChart, ...) in a
    chartSpace part XML."""
    root = etree.fromstring(blob)
    out: list[str] = []
    for el in root.iter():
        if not isinstance(el.tag, str):
            continue
        name = etree.QName(el).localname
        if name.endswith("Chart") and name not in out:
            out.append(name)
    return out


def _walk_shapes(shapes):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for sh in shapes:
        yield sh
        try:
            is_group = sh.shape_type == MSO_SHAPE_TYPE.GROUP
        except Exception:
            is_group = False
        if is_group:
            yield from _walk_shapes(sh.shapes)


def chart_types_from_pptx(prs) -> list[str]:
    """Distinct chart plot types across an opened python-pptx Presentation."""
    out: list[str] = []
    for slide in prs.slides:
        for sh in _walk_shapes(slide.shapes):
            if getattr(sh, "has_chart", False):
                for t in chart_types_from_blob(sh.chart.part.blob):
                    if t not in out:
                        out.append(t)
    return out


def _edge_pts(grid_indices: set[int]) -> list[int]:
    """Grid indices -> sorted int point values, capped at EDGE_CAP."""
    return [int(round(g * EDGE_GRID_PT))
            for g in sorted(grid_indices)][:EDGE_CAP]


def shape_stats_for_slide(slide, slide_width: int | None = None,
                          slide_height: int | None = None) -> ShapeStats:
    from pptx.enum.dml import MSO_FILL
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    n_shapes = n_textboxes = n_pictures = n_tables = 0
    n_charts = n_groups = 0
    chart_types: list[str] = []
    fonts: list[str] = []
    texts: list[str] = []
    has_smartart = False
    x_grid: set[int] = set()
    y_grid: set[int] = set()
    fills: set[str] = set()
    area_emu = 0
    for sh in _walk_shapes(slide.shapes):
        n_shapes += 1
        try:
            left, top = sh.left, sh.top
            width, height = sh.width, sh.height
        except Exception:
            left = top = width = height = None
        if None not in (left, top, width, height):
            x_grid.add(round(left / EDGE_GRID_EMU))
            x_grid.add(round((left + width) / EDGE_GRID_EMU))
            y_grid.add(round(top / EDGE_GRID_EMU))
            y_grid.add(round((top + height) / EDGE_GRID_EMU))
            # Naive: overlaps and group children double-count; the
            # consumer clamps and treats this as a coarse signal only.
            area_emu += width * height
        try:
            fill = sh.fill
            if fill.type == MSO_FILL.SOLID:
                fills.add(str(fill.fore_color.rgb))
        except Exception:
            # Inherited/theme fills raise on .rgb; skip, keep count of
            # explicit solid fills only.
            pass
        try:
            stype = sh.shape_type
        except Exception:
            stype = None
        if stype == MSO_SHAPE_TYPE.GROUP:
            n_groups += 1
        if stype == MSO_SHAPE_TYPE.PICTURE:
            n_pictures += 1
        if sh.has_table:
            n_tables += 1
        if sh.has_chart:
            n_charts += 1
            try:
                for t in chart_types_from_blob(sh.chart.part.blob):
                    if t not in chart_types:
                        chart_types.append(t)
            except Exception:
                pass
        elif (sh.element.tag.endswith("}graphicFrame")
              and not sh.has_table):
            has_smartart = True
        if sh.has_text_frame:
            n_textboxes += 1
            tf = sh.text_frame
            if tf.text:
                texts.append(tf.text)
            for para in tf.paragraphs:
                for run in para.runs:
                    name = run.font.name
                    if name and name not in fonts:
                        fonts.append(name)
    title_text = ""
    try:
        title = slide.shapes.title
        if title is not None:
            title_text = title.text or ""
    except Exception:
        pass
    covered = None
    if slide_width and slide_height:
        covered = min(1.0, max(0.0,
                               area_emu / (slide_width * slide_height)))
    return ShapeStats(n_shapes=n_shapes, n_textboxes=n_textboxes,
                      n_pictures=n_pictures, n_tables=n_tables,
                      n_charts=n_charts, n_groups=n_groups,
                      chart_xml_types=chart_types,
                      has_smartart=has_smartart, title_text=title_text,
                      fonts=fonts, all_text="\n".join(texts),
                      x_edges=_edge_pts(x_grid),
                      y_edges=_edge_pts(y_grid),
                      n_distinct_fills=min(len(fills), EDGE_CAP),
                      covered_frac=covered)


# --- jsonl / manifest IO ----------------------------------------------------

def _read_jsonl(path: Path) -> list[dict]:
    if not path.exists():
        return []
    rows = []
    for line in path.read_text(encoding="utf-8").splitlines():
        line = line.strip()
        if not line:
            continue
        try:
            rows.append(json.loads(line))
        except json.JSONDecodeError:
            # Torn line from a run that crashed mid-append. Skipping is
            # safe in both files: a lost manifest line just re-ingests
            # that deck; a lost index line just weakens dedupe by one
            # slide. Raising here would brick every future run instead.
            continue
    return rows


def _append_jsonl(path: Path, rows: list[dict]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("a", encoding="utf-8") as f:
        for row in rows:
            f.write(json.dumps(row, ensure_ascii=False) + "\n")


def load_manifest(work: Path) -> dict[str, DeckManifestEntry]:
    """Last manifest line per deck_id wins (manifest is append-only)."""
    entries: dict[str, DeckManifestEntry] = {}
    for row in _read_jsonl(work / "manifest.jsonl"):
        entry = DeckManifestEntry.model_validate(row)
        entries[entry.deck_id] = entry
    return entries


def append_manifest(work: Path, entry: DeckManifestEntry) -> None:
    _append_jsonl(work / "manifest.jsonl", [entry.model_dump()])


def load_dedupe_reference(work: Path,
                          exclude_decks: set[str]) -> list[tuple[int, str]]:
    """(phash-int, slide_id) of previously indexed canonical slides,
    excluding decks about to be re-ingested in this run."""
    ref: list[tuple[int, str]] = []
    for row in _read_jsonl(work / "index.jsonl"):
        if row.get("deck_id") in exclude_decks or row.get("dup_of"):
            continue
        ref.append((int(row["phash"], 16), row["slide_id"]))
    return ref


def _match_ref(ref: list[tuple[int, str]], phash_hex: str) -> str | None:
    h = int(phash_hex, 16)
    for known, slide_id in ref:
        if bin(h ^ known).count("1") <= NEAR_DUP_HAMMING:
            return slide_id
    return None


# --- per-deck ingest --------------------------------------------------------

def _now_iso() -> str:
    return datetime.now(timezone.utc).isoformat(timespec="seconds")


def _record(deck_id: str, idx: int, png: Path, source_type: str,
            source_path: Path, text: str, stats: ShapeStats | None,
            ref: list[tuple[int, str]]) -> SlideRecord:
    with Image.open(png) as img:
        phash = dhash(img)
        if stats is not None and img.height:
            stats.aspect = round(img.width / img.height, 4)
    dup_of = _match_ref(ref, phash)
    slide_id = f"{deck_id}:{idx:03d}"
    if dup_of is None:
        ref.append((int(phash, 16), slide_id))
    return SlideRecord(deck_id=deck_id, slide_id=slide_id, slide_index=idx,
                       png_path=Path(png).resolve().as_posix(),
                       source_type=source_type,
                       source_path=Path(source_path).resolve().as_posix(),
                       extracted_text=text[:TEXT_CAP], shape_stats=stats,
                       phash=phash, dup_of=dup_of, ingested_at=_now_iso())


def _ingest_pptx(deck_id: str, pptx: Path, source_path: Path,
                 source_type: str, work: Path,
                 ref: list[tuple[int, str]]) -> list[SlideRecord]:
    from pptx import Presentation
    png_dir = work / "png" / deck_id
    if png_dir.exists():
        # Stale PNGs from a prior ingest of this deck (deck shrank, or a
        # timed-out export left partial output) must not be indexed as
        # current slides; exporters glob/overwrite this directory.
        for old in png_dir.glob("slide*.png"):
            old.unlink()
    pngs = sorted(com_convert.export_slide_pngs(
        pptx, png_dir, TARGET_WIDTH, TARGET_HEIGHT))
    prs = Presentation(str(pptx))
    stats = [shape_stats_for_slide(s, prs.slide_width, prs.slide_height)
             for s in prs.slides]
    records = []
    for idx, png in enumerate(pngs):
        ss = stats[idx] if idx < len(stats) else None
        text = ss.all_text if ss is not None else ""
        records.append(_record(deck_id, idx, png, source_type,
                               source_path, text, ss, ref))
    return records


def _ingest_pdf(deck_id: str, path: Path, work: Path,
                ref: list[tuple[int, str]]) -> list[SlideRecord]:
    import pypdfium2 as pdfium
    png_dir = work / "png" / deck_id
    png_dir.mkdir(parents=True, exist_ok=True)
    records = []
    pdf = pdfium.PdfDocument(str(path))
    try:
        for idx in range(len(pdf)):
            page = pdf[idx]
            width_pt, _ = page.get_size()
            scale = TARGET_WIDTH / width_pt if width_pt else 1.0
            pil = page.render(scale=scale).to_pil()
            png = png_dir / f"slide{idx + 1:03d}.png"
            pil.save(png)
            try:
                text = page.get_textpage().get_text_range()
            except Exception:
                text = ""
            records.append(_record(deck_id, idx, png, "pdf", path,
                                   text, None, ref))
    finally:
        pdf.close()
    return records


def ingest_deck(deck_id: str, path: Path, work: Path,
                ref: list[tuple[int, str]]) -> list[SlideRecord]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _ingest_pdf(deck_id, path, work, ref)
    if suffix == ".ppt":
        converted = com_convert.ppt_to_pptx(
            path, work / "converted" / deck_id)
        return _ingest_pptx(deck_id, converted, path, "ppt", work, ref)
    return _ingest_pptx(deck_id, path, path, "pptx", work, ref)


# --- run --------------------------------------------------------------------

def run_ingest(src: str | Path, work: str | Path,
               limit: int | None = None,
               decks: str | None = None,
               exclude_dirs: set[str] | None = None) -> dict:
    src = Path(src).resolve()
    work = Path(work).resolve()
    work.mkdir(parents=True, exist_ok=True)
    manifest = load_manifest(work)

    todo: list[tuple[str, Path, os.stat_result]] = []
    skipped = 0
    for path in iter_source_files(src, decks, exclude_dirs):
        if work == path or work in path.parents:
            # If work lives under src, derived files (work/converted/
            # <deck_id>/*.pptx) would be walked as brand-new source
            # decks on the next run, compounding every run.
            continue
        deck_id = deck_id_for(src, path)
        st = path.stat()
        prev = manifest.get(deck_id)
        if (prev is not None and prev.ok
                and abs(prev.mtime - st.st_mtime) < 1e-6
                and prev.size == st.st_size):
            skipped += 1
            continue
        todo.append((deck_id, path, st))
    if limit is not None:
        todo = todo[:limit]

    ref = load_dedupe_reference(work, {d for d, _, _ in todo})
    processed = errors = slides = 0
    for deck_id, path, st in todo:
        try:
            records = ingest_deck(deck_id, path, work, ref)
        except Exception as exc:  # never crash the walk on one bad deck
            errors += 1
            _append_jsonl(work / "ingest_errors.jsonl", [{
                "deck_id": deck_id,
                "source_path": path.resolve().as_posix(),
                "error": f"{type(exc).__name__}: {exc}",
            }])
            append_manifest(work, DeckManifestEntry(
                deck_id=deck_id, source_path=path.resolve().as_posix(),
                mtime=st.st_mtime, size=st.st_size, n_slides=0, ok=False))
            continue
        _append_jsonl(work / "index.jsonl",
                      [r.model_dump() for r in records])
        append_manifest(work, DeckManifestEntry(
            deck_id=deck_id, source_path=path.resolve().as_posix(),
            mtime=st.st_mtime, size=st.st_size, n_slides=len(records),
            ok=True))
        processed += 1
        slides += len(records)
    return {"processed": processed, "skipped": skipped,
            "errors": errors, "slides": slides}


def main(argv: list[str] | None = None) -> None:
    ap = argparse.ArgumentParser(
        description="Ingest a directory of decks into the corpus index.")
    ap.add_argument("--src", default=os.environ.get("DECKENGINE_CORPUS"),
                    help="source dir (default: env DECKENGINE_CORPUS)")
    ap.add_argument("--limit", type=int, default=None,
                    help="max decks to (re)process this run")
    ap.add_argument("--decks", default=None,
                    help="glob filter on deck filename/relpath")
    ap.add_argument("--work",
                    default=str(Path(__file__).resolve().parent / "work"),
                    help="derived-data dir (default: corpus/work)")
    ap.add_argument("--exclude", default="deckengine",
                    help="comma-separated directory names to skip anywhere "
                         "under src (default: deckengine)")
    args = ap.parse_args(argv)
    if not args.src:
        ap.error("--src required (or set DECKENGINE_CORPUS)")
    exclude = {d.strip() for d in args.exclude.split(",") if d.strip()}
    summary = run_ingest(args.src, args.work, args.limit, args.decks,
                         exclude_dirs=exclude or None)
    print(json.dumps(summary))


if __name__ == "__main__":
    main()
