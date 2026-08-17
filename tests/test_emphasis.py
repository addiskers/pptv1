"""The emphasis contract: renderer treatments (parity) + the title-entity
cross-check + the field-exists guard (provenance-lesson regression net)."""
from __future__ import annotations

from pptx import Presentation

from deckengine.components.base import BuildReport, RenderContext, \
    get_component
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch
from deckengine.llm.emphasis import _CAPABLE, check_slide_emphasis
from deckengine.schema.components import (ComparisonColumnsSpec,
                                          FunnelSpec, GanttRowSpec,
                                          KpiCardStripSpec, MiniTableSpec,
                                          PyramidSpec)
from deckengine.schema.slide_types import ChartSlideSpec


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


# -- Theme.soft --------------------------------------------------------------

def test_theme_soft_blends_toward_bg():
    t = load_theme("consulting_navy")
    tint = t.soft("accent")
    assert len(tint) == 6 and tint != t.color("accent") != t.color("bg")


# -- renderer parity: emphasis never changes heights -------------------------

def _parity(kind, plain_spec, emphasized_spec):
    ctx = make_ctx()
    comp = get_component(kind)
    w = inch(8)
    assert comp.measure(plain_spec, w, ctx) == \
        comp.measure(emphasized_spec, w, ctx)
    _, s1 = blank_slide()
    _, s2 = blank_slide()
    h1 = comp.render(s1, plain_spec, BBox(0, 0, w, inch(5)), ctx)
    h2 = comp.render(s2, emphasized_spec, BBox(0, 0, w, inch(5)), ctx)
    assert h1 == h2


def test_mini_table_emphasis_parity():
    base = dict(headers=["State", "Yield"],
                rows=[["Gujarat", "3.2"], ["Punjab", "4.1"]])
    _parity("mini_table", MiniTableSpec(**base),
            MiniTableSpec(**base, highlight_row="Punjab"))


def test_comparison_columns_emphasis_parity():
    base = dict(columns=[
        {"header": "Indonesia", "cells": [{"kind": "text_block",
                                           "text": "Biggest TAM."}]},
        {"header": "Vietnam", "cells": [{"kind": "text_block",
                                         "text": "Crowded."}]}])
    _parity("comparison_columns", ComparisonColumnsSpec(**base),
            ComparisonColumnsSpec(**base, highlight_column="Indonesia"))


def test_kpi_and_funnel_emphasis_parity():
    kpi = dict(cards=[{"title": "**42**"}, {"title": "**7**"}])
    _parity("kpi_card_strip", KpiCardStripSpec(**kpi),
            KpiCardStripSpec(**kpi, highlight_index=1))
    fun = dict(stages=[{"label": "Leads", "value": "900"},
                       {"label": "Quals", "value": "300"},
                       {"label": "Wins", "value": "60"}])
    _parity("funnel", FunnelSpec(**fun),
            FunnelSpec(**fun, highlight_index=0))


# -- new primitives render + parity ------------------------------------------

def test_pyramid_parity_and_render():
    ctx = make_ctx()
    comp = get_component("pyramid")
    spec = PyramidSpec(tiers=[{"label": "Apex", "value": "12"},
                              {"label": "Middle", "value": "240"},
                              {"label": "Base", "value": "9,000"}],
                       highlight_index=1)
    w = inch(7)
    measured = comp.measure(spec, w, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, spec, BBox(0, 0, w, inch(4)), ctx)
    assert abs(measured - consumed) <= 12700  # 1pt
    texts = " ".join("".join(r.text for p in s.text_frame.paragraphs
                             for r in p.runs)
                     for s in slide.shapes if s.has_text_frame)
    assert "Middle" in texts and "9,000" in texts


def test_gantt_parity_render_and_clamp():
    ctx = make_ctx()
    comp = get_component("gantt_row")
    spec = GanttRowSpec(
        periods=["Q1", "Q2", "Q3", "Q4"],
        items=[{"label": "Pilot", "start": 0, "end": 2,
                "milestone": "gate"},
               {"label": "Scale", "start": 1, "end": 9}],  # clamps to 4
        today_index=1, highlight_index=0)
    w = inch(8)
    measured = comp.measure(spec, w, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, spec, BBox(0, 0, w, inch(3)), ctx)
    assert abs(measured - consumed) <= 12700
    assert any("clamped" in x for x in ctx.report.warnings)


# -- the cross-check ---------------------------------------------------------

def _chart_slide(title, highlight=None):
    return ChartSlideSpec(title=title, chart={
        "kind": "native_chart", "chart_type": "bar",
        "categories": ["Indonesia", "Vietnam", "Thailand"],
        "series": [{"name": "TAM", "values": [9.0, 6.0, 4.0]}],
        "sort": "desc", "highlight": highlight, "annotation": None})


def test_single_entity_title_without_emphasis_fires():
    problems = check_slide_emphasis(
        _chart_slide("**Indonesia** is the best market to enter"))
    assert any("Indonesia" in p and "highlight" in p for p in problems)


def test_emphasized_slide_passes_and_ambiguous_skips():
    assert check_slide_emphasis(
        _chart_slide("**Indonesia** is the best market", "Indonesia")) == []
    # two entities named -> ambiguity never enters the repair loop
    assert check_slide_emphasis(
        _chart_slide("Indonesia outpaces Vietnam on growth")) == []


def test_bad_highlight_gets_inverse_check():
    problems = check_slide_emphasis(_chart_slide("Growth varies by market",
                                                 "Indonesa"))
    assert any("matches none" in p for p in problems)


def test_index_field_repair_names_the_index():
    s = FunnelSpec(stages=[{"label": "Awareness", "value": "9,000"},
                           {"label": "Consideration", "value": "1,200"},
                           {"label": "Purchase", "value": "310"}])

    class Holder:  # minimal slide-shaped object
        title = "Consideration is where the funnel leaks"

        @staticmethod
        def model_dump():
            return {"kind_holder": s.model_dump()}

    problems = check_slide_emphasis(Holder)
    assert any("highlight_index=1" in p for p in problems)


def test_guard_every_capable_field_exists():
    """The provenance lesson: a repair sentence must never name a field
    the schema lacks."""
    from deckengine.schema import components as C
    models = {"native_chart": C.NativeChartSpec,
              "mini_table": C.MiniTableSpec,
              "data_table": C.DataTableSpec,
              "comparison_columns": C.ComparisonColumnsSpec,
              "kpi_card_strip": C.KpiCardStripSpec,
              "funnel": C.FunnelSpec,
              "chevron_pathway": C.ChevronPathwaySpec,
              "matrix_2x2": C.Matrix2x2Spec}
    for kind, (field, _) in _CAPABLE.items():
        assert field in models[kind].model_fields, (kind, field)
