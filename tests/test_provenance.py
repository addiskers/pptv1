"""Marker system: token parsing, plain() semantics, enforcement, and the
structural guard that markers are never demanded in PlainStr fields."""
from __future__ import annotations

import typing

import pytest
from pydantic import BaseModel

from deckengine.llm.provenance import (MARKER_KEYS, check_slide_markers,
                                       collect_marked_figures,
                                       marker_coverage)
from deckengine.schema import components as comp_mod
from deckengine.schema import slide_types as slides_mod
from deckengine.schema.components import ChartSeriesSpec, NativeChartSpec
from deckengine.schema.rich import (MARKER_FONT, parse_rich, plain,
                                    strip_markers)
from deckengine.schema.slide_types import BulletContentSpec, ChartSlideSpec


# -- token parsing ---------------------------------------------------------

def test_marker_tokens_parse_to_glyph_spans():
    spans = parse_rich("Income rose 34% [[src:recon]] since 2019.")
    glyphs = [s for s in spans if s.marker]
    assert len(glyphs) == 1
    g = glyphs[0]
    assert g.text == "◐" and g.superscript and g.font == MARKER_FONT


@pytest.mark.parametrize("tier,glyph", [("official", "●"), ("recon", "◐"),
                                        ("est", "○")])
def test_all_three_tiers(tier, glyph):
    spans = parse_rich(f"$4.2B [[src:{tier}]]")
    assert [s.text for s in spans if s.marker] == [glyph]


def test_unknown_tier_stays_literal():
    spans = parse_rich("42 [[src:vibes]]")
    assert not any(s.marker for s in spans)
    assert "[[src:vibes]]" in "".join(s.text for s in spans)


def test_plain_drops_markers():
    assert plain("34% [[src:est]] of farms") == "34%  of farms"


def test_marker_inherits_base_color_role():
    spans = parse_rich("90% [[src:official]]", base_color_role="inverse_ink")
    assert [s.color_role for s in spans if s.marker] == ["inverse_ink"]


def test_strip_markers_for_display_surfaces():
    t = "Market reached **0.95M** [[src:recon]] registrations in 2024"
    assert strip_markers(t) == \
        "Market reached **0.95M** registrations in 2024"
    # other markup untouched; unknown tiers stay literal
    assert strip_markers("42 [[src:vibes]] **x**") == "42 [[src:vibes]] **x**"


def test_title_block_renders_without_marker_glyphs(tmp_path):
    from pptx import Presentation
    from deckengine.render.deck_builder import build_deck
    from deckengine.schema.slide_types import DeckSpec
    deck = DeckSpec.model_validate({
        "schema_version": 1, "theme": "consulting_navy",
        "meta": {"title": "t", "date": "19 Aug 2026", "footer_org": "DE"},
        "slides": [{"slide_type": "bullet_content",
                    "title": "Sales hit **948k** [[src:recon]] this year",
                    "subtitle": "Up 10% [[src:est]] vs plan",
                    "bullets": [{"text": "Body keeps its marker: 5% "
                                         "[[src:est]] uplift"}]}]})
    out = tmp_path / "t.pptx"
    build_deck(deck, out, embed_fonts=False)
    slide = Presentation(str(out)).slides[0]
    texts = [r.text for s in slide.shapes if s.has_text_frame
             for p in s.text_frame.paragraphs for r in p.runs]
    joined = "".join(texts)
    assert "948k" in joined
    # exactly one glyph survives — the BODY marker; none from title/subtitle
    assert joined.count("◐") == 0 and joined.count("○") == 1


# -- enforcement -----------------------------------------------------------

def _bullets(*texts, footnote=None, title="Groundwater stress is the binding constraint"):
    return BulletContentSpec(
        title=title, bullets=[{"text": t} for t in texts], footnote=footnote)


def test_unmarked_prose_figure_flagged():
    s = _bullets("Tubewell depth grew 40% in a decade.")
    problems = check_slide_markers(s)
    assert any("'40'" in p and "provenance marker" in p for p in problems)


def test_marked_prose_figure_passes():
    s = _bullets("Tubewell depth grew 40% [[src:recon]] in a decade.")
    assert check_slide_markers(s) == []


def test_illustrative_banned():
    s = _bullets("Numbers are illustrative only.")
    assert any("banned" in p for p in check_slide_markers(s))


def test_benign_numbers_need_no_marker():
    s = _bullets("Three phases across 2024-2026, five districts each.")
    assert check_slide_markers(s) == []


def test_chart_values_demand_marked_footnote():
    s = ChartSlideSpec(
        title="Output doubled in five seasons",
        chart=NativeChartSpec(
            chart_type="line", categories=["2020", "2024"],
            series=[ChartSeriesSpec(name="Output", values=[210.0, 428.0])],
            sort="none", highlight=None, annotation=None),
        footnote="Source: state agriculture dept.")
    problems = check_slide_markers(s)
    assert any("marked source" in p for p in problems)
    s2 = s.model_copy(update={
        "footnote": "Source: Gujarat agriculture dept 2025 [[src:official]]"})
    assert check_slide_markers(s2) == []


def test_marker_coverage_counts():
    s = _bullets("40% [[src:recon]] marked, another 62% unmarked.")
    assert marker_coverage(s) == (1, 2)


def test_collect_marked_figures_skips_official():
    s1 = _bullets("Income rose 34% [[src:recon]] on drip plots.")
    s2 = _bullets("Subsidy outlay hit $310M [[src:official]] in FY25.")
    rows = collect_marked_figures([s1, s2])
    assert [r["tier"] for r in rows] == ["recon"]
    assert rows[0]["slide"] == 1 and rows[0]["figure"] == "34"


# -- end-to-end render -------------------------------------------------------

def test_marker_renders_in_symbol_font(tmp_path):
    from pptx import Presentation

    from deckengine.render.deck_builder import build_deck
    from deckengine.schema.slide_types import DeckSpec
    spec = DeckSpec.model_validate({
        "schema_version": 1, "theme": "consulting_navy",
        "meta": {"title": "m", "date": "13 Aug 2026", "footer_org": "DE"},
        "slides": [{
            "slide_type": "bullet_content",
            "title": "Drip adoption tripled [[src:recon]] since 2019",
            "bullets": [{"text": "Coverage reached 34% [[src:est]] of "
                                 "irrigated area."}]}]})
    out = tmp_path / "markers.pptx"
    report = build_deck(spec, out)
    assert not report.warnings
    runs = [r for s in Presentation(str(out)).slides[0].shapes
            if s.has_text_frame for p in s.text_frame.paragraphs
            for r in p.runs]
    glyphs = [r for r in runs if r.text in "●◐○"]
    # title glyph is stripped for display (stray-dot fix); the BODY glyph
    # renders in the symbol face
    assert len(glyphs) == 1
    assert all(r.font.name == "Segoe UI Symbol" for r in glyphs)
    assert not any("[[src:" in r.text for r in runs)


# -- the structural guard --------------------------------------------------

def _plain_field_names(mod) -> set[tuple[str, str]]:
    """(model, field) pairs whose annotation is PlainStr (strips markup)."""
    out = set()
    for name in dir(mod):
        cls = getattr(mod, name)
        if not (isinstance(cls, type) and issubclass(cls, BaseModel)):
            continue
        hints = typing.get_type_hints(cls, include_extras=True)
        for fname, hint in hints.items():
            for part in typing.get_args(hint) or (hint,):
                meta = getattr(part, "__metadata__", ())
                if any("BeforeValidator" in type(m).__name__ for m in meta):
                    out.add((name, fname))
    return out


def test_marker_keys_are_rich_everywhere():
    """A MARKER_KEY that is PlainStr on ANY model would create an unfixable
    repair loop (the model adds the marker, validation strips it, the check
    fires again). Structurally impossible while this passes."""
    plain_fields = (_plain_field_names(comp_mod)
                    | _plain_field_names(slides_mod))
    offenders = [(m, f) for m, f in plain_fields if f in MARKER_KEYS]
    assert not offenders, (
        f"MARKER_KEYS demanded in PlainStr fields: {offenders}")
