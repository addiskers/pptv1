"""Blank-payload guards: an all-blank band measures 0, renders nothing, and
the stacker drops it with its gap — no more stray grey tracks / empty panels
(the evaluated deck's slide-23 bug class)."""
from pptx import Presentation

from deckengine.components.base import BuildReport, RenderContext, \
    get_component
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch
from deckengine.layout import stacker
from deckengine.schema.components import (CalloutBandSpec, KpiCardSpec,
                                          KpiCardStripSpec, ProgressPillSpec,
                                          TwoToneHeaderSpec)


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _assert_skipped(kind: str, data, needle: str) -> None:
    ctx = make_ctx()
    comp = get_component(kind)
    assert comp.measure(data, inch(9), ctx) == 0
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, inch(9), inch(3)), ctx)
    assert consumed == 0
    assert len(slide.shapes) == 0          # renders NOTHING
    assert any(needle in w for w in ctx.report.warnings)


def test_progress_pill_all_blank_skipped():
    _assert_skipped("progress_pill",
                    ProgressPillSpec(label=" ", value_pct=0, display=""),
                    "progress_pill: all-blank")


def test_callout_band_all_blank_skipped():
    _assert_skipped("callout_band", CalloutBandSpec(segments=["", "  "]),
                    "callout_band: all-blank")


def test_two_tone_header_all_blank_skipped():
    _assert_skipped("two_tone_header", TwoToneHeaderSpec(left="", right=" "),
                    "two_tone_header: all-blank")


def test_kpi_card_strip_all_blank_skipped():
    _assert_skipped("kpi_card_strip",
                    KpiCardStripSpec(cards=[KpiCardSpec(title=""),
                                            KpiCardSpec(title=" ")]),
                    "kpi_card_strip: all-blank")


def test_progress_pill_with_content_still_renders():
    ctx = make_ctx()
    comp = get_component("progress_pill")
    data = ProgressPillSpec(label="Coverage", value_pct=0, display="0%")
    assert comp.measure(data, inch(9), ctx) > 0
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, inch(9), inch(2)), ctx)
    assert consumed > 0 and len(slide.shapes) > 0


def test_stacker_drops_zero_height_item_and_gap():
    ctx = make_ctx()
    items = [
        stacker.item("progress_pill",
                     ProgressPillSpec(label="Adoption", value_pct=62,
                                      display="62%")),
        stacker.item("progress_pill",                 # blank -> dropped
                     ProgressPillSpec(label="", value_pct=0, display=""),
                     gap_before=1.0),
        stacker.item("progress_pill",
                     ProgressPillSpec(label="Retention", value_pct=81,
                                      display="81%"), gap_before=1.0),
    ]
    p = stacker.plan(items, BBox(0, 0, inch(9), inch(5)), ctx)
    assert len(p.items) == 2                          # blank item gone
    assert all(h > 0 for h in p.heights)
    assert any("dropped zero-height" in w for w in ctx.report.warnings)
    _, slide = blank_slide()
    stacker.render(slide, p, BBox(0, 0, inch(9), inch(5)), ctx)
    texts = [r.text for s in slide.shapes if s.has_text_frame
             for para in s.text_frame.paragraphs for r in para.runs]
    assert any("Adoption" in t for t in texts)
    assert any("Retention" in t for t in texts)
