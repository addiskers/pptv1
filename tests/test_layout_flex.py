"""Flex/page-fill cascade tests."""
from pptx import Presentation

from deckengine.components.base import BuildReport, RenderContext
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch
from deckengine.layout import stacker
from deckengine.schema.components import BulletItemSpec, BulletListSpec, TextBlockSpec
import deckengine.components  # noqa: F401  (registration)


def make_ctx():
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def test_underfull_stack_expands_and_centers():
    ctx = make_ctx()
    items = [stacker.item("text_block", TextBlockSpec(text="short line")),
             stacker.item("text_block", TextBlockSpec(text="another line"))]
    zone = BBox(0, 0, inch(8), inch(5))
    plan = stacker.plan(items, zone, ctx)
    assert not plan.overflows
    # gaps expanded beyond base level and/or centered
    assert plan.center_offset > 0
    assert plan.center_offset <= stacker.CENTER_CAP


def test_flex_item_absorbs_surplus():
    ctx = make_ctx()
    items = [stacker.item("bullet_list", BulletListSpec(
        items=[BulletItemSpec(text=f"point {i}") for i in range(3)]),
        flex=1.0)]
    zone = BBox(0, 0, inch(8), inch(4))
    plan = stacker.plan(items, zone, ctx)
    assert plan.assigned[0] > plan.heights[0]
    # growth is capped
    assert plan.assigned[0] <= plan.heights[0] * (1 + stacker.FLEX_GROWTH_CAP) + 1


def test_overflow_still_compresses():
    ctx = make_ctx()
    items = [stacker.item("text_block",
                          TextBlockSpec(text="word " * 200), gap_before=1.0)
             for _ in range(3)]
    zone = BBox(0, 0, inch(3), inch(2))
    plan = stacker.plan(items, zone, ctx)
    assert plan.overflows
    assert plan.gap_level == stacker.GAP_LEVELS[-1]


def test_fill_ratio_recorded():
    ctx = make_ctx()
    slide = blank_slide()
    stacker.stack_into(slide, BBox(0, 0, inch(8), inch(5)), ctx,
                       [stacker.item("text_block", TextBlockSpec(text="x"))])
    assert ctx.report.fills and 0 < ctx.report.fills[0] <= 1.0
