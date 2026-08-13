"""Round-trip slack audit — proof the wrap contract survives into the file.

Builds the two pinned demo decks, then re-measures EVERY written run at its
written (family, size, bold, italic) and asserts each paragraph's text fits
inside its frame's writable width. This is the regression net for the
evaluated deck's failure mode: text measured one way, rendered another.
"""
from __future__ import annotations

import json
from pathlib import Path

import pytest
from pptx import Presentation

from deckengine.core.fit_text import TextMeasurer
from deckengine.core.units import EMU_PER_PT
from deckengine.render.deck_builder import build_deck
from deckengine.schema.slide_types import DeckSpec

REPO = Path(__file__).resolve().parents[1]
SPECS = sorted((REPO / "examples" / "specs").glob("*.json"))
TOL = 2 * EMU_PER_PT  # 2pt rounding tolerance across EMU/pt conversions


def _decorative(text: str) -> bool:
    """Icon/emoji/marker glyphs measure differently per platform; the audit
    covers prose, numbers and labels — the layout-contract surface."""
    return all(not ch.isalnum() for ch in text) or any(
        ord(ch) > 0x2500 for ch in text)


@pytest.mark.parametrize("spec_path", SPECS, ids=lambda p: p.stem)
def test_every_paragraph_fits_its_frame(spec_path, tmp_path):
    spec = DeckSpec.model_validate(
        json.loads(spec_path.read_text(encoding="utf-8")))
    out = tmp_path / f"{spec_path.stem}.pptx"
    build_deck(spec, out)

    m = TextMeasurer()
    prs = Presentation(str(out))
    checked = 0
    failures: list[str] = []
    for si, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            writable = int(shape.width) - int(tf.margin_left) \
                - int(tf.margin_right)
            for para in tf.paragraphs:
                width = 0.0
                skip = False
                for run in para.runs:
                    if not run.text:
                        continue
                    if _decorative(run.text) or run.font.size is None:
                        skip = True
                        break
                    width += m.width_pt(
                        run.text, run.font.name or "Segoe UI",
                        run.font.size.pt, bold=bool(run.font.bold),
                        italic=bool(run.font.italic))
                if skip or width == 0.0:
                    continue
                checked += 1
                if width * EMU_PER_PT > writable + TOL:
                    text = "".join(r.text for r in para.runs)
                    failures.append(
                        f"slide {si}: {text[:60]!r} measures "
                        f"{width:.1f}pt but frame is "
                        f"{writable / EMU_PER_PT:.1f}pt wide")
    assert not failures, "\n".join(failures[:12])
    assert checked > 50  # the audit actually saw the deck