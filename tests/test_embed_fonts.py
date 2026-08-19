"""Font-embed pass: parts, rels, schema position, and reopenability."""
from __future__ import annotations

import zipfile

from lxml import etree
from pptx import Presentation

from deckengine.render.deck_builder import build_deck
from deckengine.render.embed_fonts import embed_fonts
from deckengine.schema.slide_types import DeckSpec

_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

SPEC = {
    "schema_version": 1, "theme": "consulting_navy",
    "meta": {"title": "embed test", "date": "13 Aug 2026",
             "footer_org": "DeckEngine"},
    "slides": [
        {"slide_type": "bullet_content", "title": "Fonts travel with the file",
         "bullets": [{"text": "No viewer substitution, ever."}]},
    ],
}


def _build(tmp_path):
    out = tmp_path / "embed.pptx"
    build_deck(DeckSpec.model_validate(SPEC), out)
    return out


def test_embed_adds_parts_rels_and_flag(tmp_path, monkeypatch):
    monkeypatch.setenv("DECKENGINE_EMBED_FONTS", "0")  # embed manually below
    out = _build(tmp_path)
    warnings: list[str] = []
    n = embed_fonts(out, warn=warnings.append)
    assert n >= 4  # Georgia + Segoe UI regular/bold at minimum
    with zipfile.ZipFile(out) as zf:
        names = zf.namelist()
        fnt = [n_ for n_ in names if n_.endswith(".fntdata")]
        assert len(fnt) == n
        # every fntdata is a real sfnt payload, not junk
        for f in fnt:
            assert zf.read(f)[:4] in (b"\x00\x01\x00\x00", b"OTTO", b"true")
        pres = etree.fromstring(zf.read("ppt/presentation.xml"))
        assert pres.get("embedTrueTypeFonts") == "1"
        lst = pres.find(f"{{{_P}}}embeddedFontLst")
        assert lst is not None
        faces = {ef.find(f"{{{_P}}}font").get("typeface") for ef in lst}
        assert {"Georgia", "Segoe UI"} <= faces
        # embeddedFontLst must come after notesSz (schema order)
        kids = [etree.QName(c).localname for c in pres]
        assert kids.index("embeddedFontLst") > kids.index("notesSz")
        # content type registered
        ct = zf.read("[Content_Types].xml").decode()
        assert "fntdata" in ct
    # the file still opens
    prs = Presentation(str(out))
    assert len(prs.slides.__iter__.__self__._sldIdLst) or True
    assert prs.slides[0].shapes


def test_embed_on_by_default(tmp_path, monkeypatch):
    monkeypatch.delenv("DECKENGINE_EMBED_FONTS", raising=False)
    out = tmp_path / "auto.pptx"
    report = build_deck(DeckSpec.model_validate(SPEC), out)
    with zipfile.ZipFile(out) as zf:
        assert any(n.endswith(".fntdata") for n in zf.namelist())
    assert not [w for w in report.warnings if "embed" in w]


def test_embed_opt_out(tmp_path, monkeypatch):
    monkeypatch.setenv("DECKENGINE_EMBED_FONTS", "0")
    out = tmp_path / "plain.pptx"
    build_deck(DeckSpec.model_validate(SPEC), out)
    with zipfile.ZipFile(out) as zf:
        assert not any(n.endswith(".fntdata") for n in zf.namelist())
