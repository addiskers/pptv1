"""S5 craft parity: kicker ladder, semantic colour, chart density, sparse
lint, zero-based axes."""
from __future__ import annotations

from pptx import Presentation
from pptx.oxml.ns import qn

from deckengine.components.base import BuildReport, RenderContext, \
    get_component
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import Span, TextMeasurer
from deckengine.core.semantic_roles import bucket, rag_map
from deckengine.core.theme import load_theme
from deckengine.core.units import inch
from deckengine.llm.format_rules import check_outline_chart_density
from deckengine.llm.story import Outline, OutlineSlide
from deckengine.llm.writing import check_slide_writing
from deckengine.render.deck_builder import build_deck
from deckengine.schema.components import (ChartSeriesSpec, NativeChartSpec,
                                          SectionHeaderSpec)
from deckengine.schema.slide_types import BulletContentSpec, DeckSpec


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


# -- kicker ladder -----------------------------------------------------------

def test_kicker_measures_taller_and_renders_caps_spc():
    ctx = make_ctx()
    comp = get_component("section_header")
    plain_h = comp.measure(SectionHeaderSpec(title="T"), inch(9), ctx)
    kicked = SectionHeaderSpec(title="T", kicker="The opportunity")
    assert comp.measure(kicked, inch(9), ctx) > plain_h
    _, slide = blank_slide()
    comp.render(slide, kicked, BBox(0, 0, inch(9), inch(1.2)), ctx)
    rprs = [r._r.get_or_add_rPr()
            for s in slide.shapes if s.has_text_frame
            for p in s.text_frame.paragraphs for r in p.runs
            if r.text == "The opportunity"]
    assert rprs, "kicker run not found"
    assert rprs[0].get("cap") == "all"
    assert int(rprs[0].get("spc")) == 90  # 0.9pt tracking in 1/100 pt


def test_caps_span_measured_uppercase_with_tracking():
    m = TextMeasurer()
    base = m.span_width_emu(Span("wide open"), "Segoe UI", 10)
    caps = m.span_width_emu(Span("wide open", caps=True, spc_pts=0.9),
                            "Segoe UI", 10)
    assert caps > base  # uppercase glyphs + tracking are wider


def test_kicker_flows_from_spec_and_bumps_zone(tmp_path):
    spec = DeckSpec.model_validate({
        "schema_version": 1, "theme": "consulting_navy",
        "meta": {"title": "k", "date": "13 Aug 2026", "footer_org": "DE"},
        "slides": [{"slide_type": "bullet_content",
                    "title": "Momentum is compounding",
                    "kicker": "the opportunity",
                    "bullets": [{"text": "evidence line"}]}]})
    out = tmp_path / "kicker.pptx"
    build_deck(spec, out)
    runs = [r for s in Presentation(str(out)).slides[0].shapes
            if s.has_text_frame for p in s.text_frame.paragraphs
            for r in p.runs]
    kick = [r for r in runs if r.text == "the opportunity"]
    assert kick and kick[0]._r.get_or_add_rPr().get("cap") == "all"


# -- semantic colour ---------------------------------------------------------

def test_bucket_and_ambiguity():
    assert bucket("Overdrawn basins") == "negative"
    assert bucket("Moderate stress zones") is None  # two buckets -> refuse
    assert bucket("Healthy recharge") == "positive"
    assert bucket("North region") is None


def test_rag_map_needs_two_buckets():
    assert rag_map(["Overdrawn", "Critical", "Depleted"]) is None  # 1 bucket
    m = rag_map(["Overdrawn", "Moderate", "Healthy"])
    assert m == {"Overdrawn": "negative", "Moderate": "warning",
                 "Healthy": "positive"}


def test_bar_categories_colored_by_meaning():
    ctx = make_ctx()
    comp = get_component("native_chart")
    data = NativeChartSpec(
        chart_type="bar", categories=["Overdrawn", "Moderate", "Healthy"],
        series=[ChartSeriesSpec(name="Blocks", values=[48.0, 30.0, 22.0])],
        sort="none", highlight=None, annotation=None)
    _, slide = blank_slide()
    comp.render(slide, data, BBox(0, 0, inch(8), inch(4)), ctx)
    chart = next(s.chart for s in slide.shapes if s.has_chart)
    pts = chart.series[0].points
    theme = ctx.theme
    assert str(pts[0].format.fill.fore_color.rgb) == theme.color("negative")
    assert str(pts[1].format.fill.fore_color.rgb) == theme.color("warning")
    assert str(pts[2].format.fill.fore_color.rgb) == theme.color("positive")


def test_zero_based_axis_on_positive_bars():
    ctx = make_ctx()
    comp = get_component("native_chart")
    data = NativeChartSpec(
        chart_type="bar", categories=["A", "B"],
        series=[ChartSeriesSpec(name="s", values=[80.0, 95.0])],
        sort="none", highlight=None, annotation=None)
    _, slide = blank_slide()
    comp.render(slide, data, BBox(0, 0, inch(8), inch(4)), ctx)
    chart = next(s.chart for s in slide.shapes if s.has_chart)
    assert chart.value_axis.minimum_scale == 0.0


# -- chart density + sparse lint ----------------------------------------------

def _outline(types):
    return Outline(governing_thought="Act now.",
                   slides=[OutlineSlide(slide_type=t, claim=f"Claim {i} holds")
                           for i, t in enumerate(types)])


def test_chart_density_flags_prose_heavy_outline():
    o = _outline(["title", "bullet_content", "bullet_content",
                  "framework_slide", "timeline_slide", "exec_summary"])
    problems = check_outline_chart_density(o)
    assert problems and "60%" in problems[0]


def test_chart_density_passes_visual_outline():
    o = _outline(["title", "chart_slide", "custom_layout", "chart_slide",
                  "bullet_content", "exec_summary"])
    assert check_outline_chart_density(o) == []


def test_sparse_dense_body_slide_flagged():
    s = BulletContentSpec(title="Five moves win the market",
                          bullets=[{"text": "Secure supply first."}])
    assert any("sparse" in p or "words" in p for p in check_slide_writing(s))
