"""Corpus ingest tests — no COM/Office: com_convert is monkeypatched."""
import json
import os
import random
import time
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from corpus import com_convert, ingest
from corpus.models import SlideRecord


# --- image fixtures ----------------------------------------------------------

def gradient_img(w=320, h=180, noise_seed=None):
    img = Image.new("L", (w, h))
    px = img.load()
    for x in range(w):
        val = int(x * 255 / (w - 1))
        for y in range(h):
            px[x, y] = val
    if noise_seed is not None:
        rng = random.Random(noise_seed)
        for _ in range(w * h // 20):
            x, y = rng.randrange(w), rng.randrange(h)
            px[x, y] = max(0, min(255, px[x, y] + rng.choice((-4, 4))))
    return img.convert("RGB")


def block_img(seed, w=320, h=180):
    rng = random.Random(seed)
    cells = Image.new("L", (17, 16))
    cells.putdata([rng.choice((0, 255)) for _ in range(17 * 16)])
    return cells.resize((w, h), Image.Resampling.NEAREST).convert("RGB")


# --- dhash -------------------------------------------------------------------

def test_dhash_identical_images_identical_hash():
    a, b = block_img(1), block_img(1)
    assert ingest.dhash(a) == ingest.dhash(b)
    assert len(ingest.dhash(a)) == 64  # 256-bit hex


def test_dhash_noised_copy_is_near_dupe():
    a = gradient_img()
    b = gradient_img(noise_seed=7)
    assert ingest.hamming(ingest.dhash(a), ingest.dhash(b)) <= 4


def test_dhash_different_images_far_apart():
    d = ingest.hamming(ingest.dhash(block_img(1)),
                       ingest.dhash(block_img(2)))
    assert d > 4


# --- corpus fixtures ---------------------------------------------------------

def build_pptx(path: Path) -> None:
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    s1.shapes.title.text = "Corpus Deck Title"
    tb = s1.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(1))
    tb.text_frame.text = "Alpha body text"
    for _ in range(2):  # slides 2 and 3 text-identical
        s = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        tb = s.shapes.add_textbox(Inches(1), Inches(1),
                                  Inches(5), Inches(1))
        tb.text_frame.text = "Repeated body"
    prs.save(str(path))


def build_pdf(path: Path) -> None:
    p1, p2 = block_img(300), block_img(400)
    p1.save(str(path), "PDF", save_all=True, append_images=[p2])


def build_text_pdf(path: Path, text: str = "CorpusPdfMarker") -> None:
    """Minimal one-page PDF with a real Helvetica text layer (PIL-saved
    PDFs are image-only, so they can never verify text extraction)."""
    objs = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>",
    ]
    stream = ("BT /F1 24 Tf 72 700 Td (%s) Tj ET" % text).encode("ascii")
    objs.append(b"<< /Length %d >>\nstream\n" % len(stream) + stream
                + b"\nendstream")
    objs.append(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")
    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for i, body in enumerate(objs, start=1):
        offsets.append(len(out))
        out += b"%d 0 obj\n" % i + body + b"\nendobj\n"
    xref_at = len(out)
    out += b"xref\n0 %d\n" % (len(objs) + 1)
    out += b"0000000000 65535 f \n"
    for off in offsets:
        out += b"%010d 00000 n \n" % off
    out += (b"trailer\n<< /Size %d /Root 1 0 R >>\nstartxref\n%d\n%%%%EOF\n"
            % (len(objs) + 1, xref_at))
    path.write_bytes(bytes(out))


def fake_export(pptx, out_dir, width=1280, height=720, timeout=300):
    """Slide 1 -> unique image; slides 2+ -> one identical image."""
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(pptx))  # raises on corrupt files, like COM
    paths = []
    for i in range(len(prs.slides)):
        img = block_img(100) if i == 0 else block_img(200)
        png = out_dir / f"slide{i + 1:03d}.png"
        img.save(png)
        paths.append(png)
    return paths


def fake_ppt_to_pptx(src, out_dir, timeout=180):
    """Stand-in for COM SaveAs: 'converts' a legacy ppt by building a
    real pptx of the same stem in out_dir."""
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    dst = out_dir / (Path(src).stem + ".pptx")
    build_pptx(dst)
    return dst


def setup_corpus(tmp_path, monkeypatch, with_corrupt=False):
    src = tmp_path / "src"
    src.mkdir()
    build_pptx(src / "deck.pptx")
    build_pdf(src / "doc.pdf")
    (src / "~$deck.pptx").write_bytes(b"office temp junk")  # must skip
    if with_corrupt:
        (src / "bad.pptx").write_bytes(b"this is not a zip archive")
    monkeypatch.setattr(com_convert, "export_slide_pngs", fake_export)
    return src, tmp_path / "work"


def read_index(work: Path) -> list[SlideRecord]:
    lines = (work / "index.jsonl").read_text(encoding="utf-8").splitlines()
    return [SlideRecord.model_validate_json(li) for li in lines
            if li.strip()]


# --- ingest ------------------------------------------------------------------

def test_ingest_pptx_and_pdf_with_dedupe(tmp_path, monkeypatch):
    src, work = setup_corpus(tmp_path, monkeypatch)
    summary = ingest.run_ingest(src, work)
    assert summary["processed"] == 2
    assert summary["errors"] == 0
    assert not (work / "ingest_errors.jsonl").exists()  # ~$ was skipped

    recs = read_index(work)
    assert len(recs) == 5  # 3 pptx slides + 2 pdf pages
    pptx_recs = [r for r in recs if r.source_type == "pptx"]
    pdf_recs = [r for r in recs if r.source_type == "pdf"]
    assert len(pptx_recs) == 3
    assert len(pdf_recs) == 2

    # slide 3 renders identically to slide 2 -> near-dupe chain to first
    assert pptx_recs[0].dup_of is None
    assert pptx_recs[1].dup_of is None
    assert pptx_recs[2].dup_of == pptx_recs[1].slide_id

    for r in pdf_recs:
        assert isinstance(r.extracted_text, str)
        assert r.shape_stats is None
        png = Image.open(r.png_path)
        assert 1275 <= png.width <= 1285
        png.close()
    for r in recs:
        assert r.slide_id == f"{r.deck_id}:{r.slide_index:03d}"
        assert len(r.phash) == 64


def test_structural_stats_captured(tmp_path, monkeypatch):
    src, work = setup_corpus(tmp_path, monkeypatch)
    ingest.run_ingest(src, work)
    first = [r for r in read_index(work) if r.source_type == "pptx"][0]
    stats = first.shape_stats
    assert stats is not None
    assert stats.n_textboxes > 0
    assert stats.title_text == "Corpus Deck Title"
    assert "Alpha body text" in stats.all_text
    assert stats.n_charts == 0
    assert not stats.has_smartart


def test_resume_skips_then_retouch_reingests_one_deck(tmp_path,
                                                      monkeypatch):
    src, work = setup_corpus(tmp_path, monkeypatch)
    ingest.run_ingest(src, work)
    n1 = len(read_index(work))

    second = ingest.run_ingest(src, work)
    assert second["processed"] == 0
    assert second["skipped"] == 2
    assert len(read_index(work)) == n1  # zero new index lines

    future = time.time() + 30
    os.utime(src / "deck.pptx", (future, future))
    third = ingest.run_ingest(src, work)
    assert third["processed"] == 1  # the touched pptx only
    recs = read_index(work)
    assert len(recs) == n1 + 3
    new = recs[n1:]
    assert all(r.source_type == "pptx" for r in new)
    # The deck's own stale index lines must be excluded from the dedupe
    # reference: otherwise every re-ingested slide would dup_of its own
    # previous incarnation (same slide_id!) instead of being canonical.
    assert new[0].dup_of is None
    assert new[1].dup_of is None
    assert new[2].dup_of == new[1].slide_id


def test_corrupt_deck_logged_and_run_continues(tmp_path, monkeypatch):
    src, work = setup_corpus(tmp_path, monkeypatch, with_corrupt=True)
    summary = ingest.run_ingest(src, work)
    assert summary["errors"] == 1
    errs = [json.loads(li) for li in
            (work / "ingest_errors.jsonl")
            .read_text(encoding="utf-8").splitlines() if li.strip()]
    assert len(errs) == 1
    assert errs[0]["source_path"].endswith("bad.pptx")
    assert errs[0]["error"]
    assert len(read_index(work)) == 5  # good decks still indexed


def test_com_unavailable_logged_and_run_continues(tmp_path, monkeypatch):
    src, work = setup_corpus(tmp_path, monkeypatch)

    def no_com(*args, **kwargs):
        raise com_convert.ComUnavailable("no Office on this box")

    monkeypatch.setattr(com_convert, "export_slide_pngs", no_com)
    summary = ingest.run_ingest(src, work)
    assert summary["errors"] == 1  # the pptx deck
    assert summary["processed"] == 1  # the pdf still ingested
    errs = (work / "ingest_errors.jsonl").read_text(encoding="utf-8")
    assert "ComUnavailable" in errs


def test_pdf_text_layer_actually_extracted(tmp_path):
    """Guards the pypdfium2 get_text_range call: an image-only fixture
    would pass even if extraction always raised into the except-fallback."""
    src = tmp_path / "src"
    src.mkdir()
    build_text_pdf(src / "report.pdf")
    work = tmp_path / "work"
    summary = ingest.run_ingest(src, work)  # no COM needed for pdf
    assert summary["processed"] == 1
    recs = read_index(work)
    assert len(recs) == 1
    assert "CorpusPdfMarker" in recs[0].extracted_text


def test_reingest_cleans_stale_pngs(tmp_path, monkeypatch):
    """A deck that shrinks between runs must not keep stale PNGs in its
    png dir where glob-based exporters would index them as slides."""
    src, work = setup_corpus(tmp_path, monkeypatch)
    ingest.run_ingest(src, work)
    deck_id = ingest.deck_id_for(src, src / "deck.pptx")
    stale = work / "png" / deck_id / "slide999.png"
    block_img(999).save(stale)
    future = time.time() + 60
    os.utime(src / "deck.pptx", (future, future))
    ingest.run_ingest(src, work)
    assert not stale.exists()


def test_torn_trailing_jsonl_does_not_brick_resume(tmp_path, monkeypatch):
    """A crash mid-append leaves a partial last line; every later run
    must still load the manifest/index instead of raising forever."""
    src, work = setup_corpus(tmp_path, monkeypatch)
    ingest.run_ingest(src, work)
    with (work / "index.jsonl").open("a", encoding="utf-8") as f:
        f.write('{"deck_id": "abc", "slide_id')  # torn, no newline
    with (work / "manifest.jsonl").open("a", encoding="utf-8") as f:
        f.write('{"deck_id": "de')
    summary = ingest.run_ingest(src, work)
    assert summary["errors"] == 0
    assert summary["skipped"] == 2  # intact earlier lines still resume


def test_legacy_ppt_converted_and_indexed(tmp_path, monkeypatch):
    src = tmp_path / "src"
    src.mkdir()
    (src / "old.ppt").write_bytes(b"legacy binary blob")
    monkeypatch.setattr(com_convert, "ppt_to_pptx", fake_ppt_to_pptx)
    monkeypatch.setattr(com_convert, "export_slide_pngs", fake_export)
    work = tmp_path / "work"
    summary = ingest.run_ingest(src, work)
    assert summary["processed"] == 1
    recs = read_index(work)
    assert len(recs) == 3
    assert all(r.source_type == "ppt" for r in recs)
    assert recs[0].source_path.endswith("old.ppt")  # original, not converted
    assert recs[0].shape_stats is not None  # stats from the converted pptx


def test_work_under_src_derived_files_not_rewalked(tmp_path, monkeypatch):
    """work/converted/<id>/*.pptx must never be picked up as new source
    decks when the work dir happens to live under the source dir."""
    src = tmp_path / "src"
    src.mkdir()
    (src / "old.ppt").write_bytes(b"legacy binary blob")
    monkeypatch.setattr(com_convert, "ppt_to_pptx", fake_ppt_to_pptx)
    monkeypatch.setattr(com_convert, "export_slide_pngs", fake_export)
    work = src / "work"
    first = ingest.run_ingest(src, work)
    assert first["processed"] == 1
    assert (work / "converted").exists()  # derived pptx now sits under src
    second = ingest.run_ingest(src, work)
    assert second["processed"] == 0
    assert second["skipped"] == 1
    assert len(read_index(work)) == 3  # no duplicate records appended
