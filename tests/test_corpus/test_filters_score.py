"""Filter + score tests — records built directly, no COM/Office."""
import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

from corpus import com_convert, filters, ingest, score
from corpus.models import ShapeStats, SlideRecord

WORDS10 = " ".join(f"w{i}" for i in range(10))
WORDS40 = " ".join(f"w{i}" for i in range(40))
WORDS80 = " ".join(f"w{i}" for i in range(80))
WORDS300 = " ".join(f"w{i}" for i in range(300))


# --- record fixtures ---------------------------------------------------------

def make_rec(slide_id="deck1:000", *, text="", pdf=False, dup_of=None,
             title="", n_shapes=5, n_pictures=0, n_tables=0, n_charts=0,
             fonts=(), x_edges=(), covered=None, aspect=None, n_fills=0):
    stats = None
    if not pdf:
        stats = ShapeStats(
            n_shapes=n_shapes, n_textboxes=2, n_pictures=n_pictures,
            n_tables=n_tables, n_charts=n_charts, n_groups=0,
            chart_xml_types=[], has_smartart=False, title_text=title,
            fonts=list(fonts), all_text=text, x_edges=list(x_edges),
            y_edges=[], n_distinct_fills=n_fills, covered_frac=covered,
            aspect=aspect)
    return SlideRecord(
        deck_id="deck1", slide_id=slide_id, slide_index=0,
        png_path="slide.png", source_type="pdf" if pdf else "pptx",
        source_path="deck.pptx", extracted_text=text, shape_stats=stats,
        phash="0" * 64, dup_of=dup_of,
        ingested_at="2026-01-01T00:00:00+00:00")


def survivor_rec(slide_id="deck1:000"):
    return make_rec(slide_id, text=WORDS40, n_charts=1,
                    fonts=("Arial", "Calibri"), covered=0.6,
                    x_edges=(72, 144, 216, 288, 360, 432),
                    aspect=round(16 / 9, 4),
                    title="Revenue bridge FY26")


# --- filter reasons ----------------------------------------------------------

def test_dup_reason():
    rec = make_rec(dup_of="deck1:001", text=WORDS40, n_charts=1)
    v = filters.filter_slide(rec, {})
    assert not v.keep
    assert "dup" in v.reasons


def test_placeholder_reason():
    rec = make_rec(text=WORDS40 + " CLICK TO ADD title", n_charts=1)
    v = filters.filter_slide(rec, {})
    assert not v.keep
    assert "placeholder" in v.reasons


def test_too_sparse_reason():
    rec = make_rec(text=WORDS10, n_shapes=6,
                   title="Deep dive metrics detail")
    v = filters.filter_slide(rec, {})
    assert not v.keep
    assert "too_sparse" in v.reasons
    assert "boilerplate" not in v.reasons


def test_wall_of_text_reason():
    rec = make_rec(text=WORDS300, n_pictures=1)  # picture is no excuse
    v = filters.filter_slide(rec, {})
    assert not v.keep
    assert "wall_of_text" in v.reasons


def test_era_by_aspect_and_by_font():
    by_aspect = make_rec(text=WORDS40, n_charts=1, aspect=1.334)
    v = filters.filter_slide(by_aspect, {})
    assert not v.keep
    assert "era" in v.reasons

    by_font = make_rec(text=WORDS40, n_charts=1,
                       aspect=round(16 / 9, 4),
                       fonts=("Comic Sans MS",))
    v = filters.filter_slide(by_font, {})
    assert not v.keep
    assert "era" in v.reasons


def test_survivor_passes_all_filters():
    v = filters.filter_slide(survivor_rec(), {})
    assert v.keep
    assert v.reasons == []


def test_boilerplate_keeps_first_two_agenda_kills_third():
    seen = {}
    verdicts = [
        filters.filter_slide(
            make_rec(f"deck1:{i:03d}", title="Agenda", n_shapes=2,
                     text="Agenda\nIntro\nMarket"),
            seen)
        for i in range(3)
    ]
    # Exemplars survive even though a 3-word agenda is sparse.
    assert verdicts[0].keep and verdicts[0].reasons == []
    assert verdicts[1].keep and verdicts[1].reasons == []
    assert not verdicts[2].keep
    assert verdicts[2].reasons == ["boilerplate"]


def test_table_excuses_wall_and_picture_excuses_sparse():
    # A dense but GOOD table slide must not die to wall_of_text.
    dense_table = make_rec(text=WORDS300, n_tables=1)
    v = filters.filter_slide(dense_table, {})
    assert v.keep and v.reasons == []
    # A sparse slide carried by a picture must not die too_sparse.
    pic_slide = make_rec(text=WORDS10, n_pictures=1,
                         title="Deep dive metrics detail")
    v = filters.filter_slide(pic_slide, {})
    assert v.keep and v.reasons == []


def test_blank_slide_dies_sparse_not_divider_exemplar():
    seen = {}
    v = filters.filter_slide(make_rec(n_shapes=0, text=""), seen)
    assert not v.keep
    assert "too_sparse" in v.reasons
    assert "boilerplate" not in v.reasons
    assert seen == {}  # blank slides never consume divider quota


def test_pdf_record_never_boilerplate_killed():
    sparse_pdf = make_rec(pdf=True, text="Thank you")
    v = filters.filter_slide(sparse_pdf, {})
    assert not v.keep
    assert "too_sparse" in v.reasons
    assert "boilerplate" not in v.reasons

    wordy_pdf = make_rec(pdf=True, text="Agenda\n" + WORDS40)
    v = filters.filter_slide(wordy_pdf, {})
    assert v.keep
    assert v.reasons == []


# --- deterministic score -----------------------------------------------------

def test_score_crafted_good_record_high():
    rec = make_rec(text=WORDS80, n_charts=1, fonts=("Arial", "Calibri"),
                   covered=0.6, x_edges=(72, 144, 216, 288, 360, 432),
                   n_fills=3)
    ds = score.score_slide(rec)
    assert ds.score >= 80
    assert ds.parts["fill"] == 25.0
    assert ds.parts["alignment"] == 25.0
    assert ds.parts["font_discipline"] == 20.0
    assert ds.parts["balance"] == 15.0


def test_score_careless_record_low():
    rec = make_rec(text=WORDS300, fonts=tuple(f"F{i}" for i in range(5)),
                   covered=0.15, x_edges=tuple(range(0, 300, 10)),
                   n_fills=10)
    assert len(rec.shape_stats.x_edges) == 30
    ds = score.score_slide(rec)
    assert ds.score <= 40
    assert ds.parts["alignment"] == 0.0
    assert ds.parts["balance"] == 3.0


def test_score_pdf_gets_neutral_geometry_parts():
    ds = score.score_slide(make_rec(pdf=True, text=WORDS40))
    assert ds.parts["fill"] == 12.0
    assert ds.parts["alignment"] == 12.0
    # No stats: 0 fonts / 0 fills are free, 40 words with no visual
    # is word-range-only balance. 12 + 12 + 20 + 15 + 8 = 67.
    assert ds.parts["font_discipline"] == 20.0
    assert ds.parts["fill_discipline"] == 15.0
    assert ds.parts["balance"] == 8.0
    assert ds.score == 67


# --- CLI end-to-end ----------------------------------------------------------

def test_filters_and_score_cli_end_to_end(tmp_path):
    work = tmp_path / "work"
    work.mkdir()
    recs = [
        survivor_rec("deck1:000"),
        make_rec("deck1:001", dup_of="deck1:000", text=WORDS40),
        make_rec("deck1:002", text=WORDS10, n_shapes=6,
                 title="Deep dive metrics detail"),
        make_rec("deck1:003", text=WORDS300),
        make_rec("deck1:004", text="Click to add text\n" + WORDS40),
    ]
    with (work / "index.jsonl").open("w", encoding="utf-8") as f:
        for r in recs:
            f.write(json.dumps(r.model_dump()) + "\n")

    filters.main(["--work", str(work)])
    lines = (work / "filtered.jsonl").read_text(
        encoding="utf-8").splitlines()
    verdicts = [filters.FilterVerdict.model_validate_json(li)
                for li in lines if li.strip()]
    assert [v.slide_id for v in verdicts] == [r.slide_id for r in recs]
    assert [v.keep for v in verdicts] == [True, False, False,
                                          False, False]

    score.main(["--work", str(work)])
    lines = (work / "scores.jsonl").read_text(
        encoding="utf-8").splitlines()
    scores = [score.DetScore.model_validate_json(li)
              for li in lines if li.strip()]
    assert len(scores) == 1  # keep=true slides only
    assert scores[0].slide_id == "deck1:000"
    assert scores[0].score >= 80
    assert set(scores[0].parts) == {"fill", "alignment",
                                    "font_discipline",
                                    "fill_discipline", "balance"}


# --- ingest populates the new ShapeStats fields ------------------------------

def build_geo_pptx(path: Path) -> None:
    prs = Presentation()  # default 10 x 7.5in canvas
    s = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    s.shapes.title.text = "Revenue Bridge FY26"
    tb = s.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(1))
    tb.text_frame.text = "Alpha body text"
    tb.fill.solid()
    tb.fill.fore_color.rgb = RGBColor(0x11, 0x22, 0x33)
    prs.save(str(path))


def fake_export_720p(pptx, out_dir, width=1280, height=720, timeout=300):
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(pptx))
    paths = []
    for i in range(len(prs.slides)):
        png = out_dir / f"slide{i + 1:03d}.png"
        Image.new("RGB", (1280, 720), (30, 60, 90)).save(png)
        paths.append(png)
    return paths


def test_ingest_populates_new_shape_stats(tmp_path, monkeypatch):
    src = tmp_path / "src"
    src.mkdir()
    build_geo_pptx(src / "deck.pptx")
    monkeypatch.setattr(com_convert, "export_slide_pngs",
                        fake_export_720p)
    work = tmp_path / "work"
    ingest.run_ingest(src, work)
    lines = (work / "index.jsonl").read_text(
        encoding="utf-8").splitlines()
    recs = [SlideRecord.model_validate_json(li) for li in lines
            if li.strip()]
    assert len(recs) == 1
    st = recs[0].shape_stats
    assert st is not None
    # Textbox at 1in..6in horizontally -> 72pt and 432pt edges.
    assert 72 in st.x_edges
    assert 432 in st.x_edges
    assert st.x_edges == sorted(st.x_edges)
    # Textbox at 2in..3in vertically -> 144pt and 216pt edges.
    assert 144 in st.y_edges
    assert 216 in st.y_edges
    assert st.n_distinct_fills == 1  # only the explicit solid fill
    assert st.covered_frac is not None
    assert 0.0 < st.covered_frac < 1.0
    # Aspect comes from the monkeypatched 1280x720 PNG export.
    assert abs(st.aspect - 1280 / 720) < 1e-3
