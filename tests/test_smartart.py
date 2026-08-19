"""True-SmartArt clone: parts, data model, recolor, fallback, toggle."""
from __future__ import annotations

import zipfile

from lxml import etree
from pptx import Presentation

from deckengine.components.base import BuildReport, RenderContext, \
    get_component
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch
from deckengine.schema.components import SmartDiagramSpec, SmartNodeSpec

_DGM = "http://schemas.openxmlformats.org/drawingml/2006/diagram"


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def _spec(layout="cycle", **kw):
    kw.setdefault("nodes", [SmartNodeSpec(label=f"N{i}") for i in range(4)])
    return SmartDiagramSpec(layout=layout, **kw)


def _render_deck(tmp_path, specs):
    ctx = make_ctx()
    prs = Presentation()
    for spec in specs:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        get_component("smart_diagram").render(
            slide, spec, BBox(inch(0.5), inch(1), inch(9), inch(5)), ctx)
    out = tmp_path / "sa.pptx"
    prs.save(str(out))
    return out, ctx


def test_clone_writes_parts_no_drawing(tmp_path):
    out, ctx = _render_deck(tmp_path, [_spec("cycle")])
    with zipfile.ZipFile(out) as zf:
        names = zf.namelist()
        for part in ("data", "layout", "quickStyle", "colors"):
            assert f"ppt/diagrams/{part}1.xml" in names, part
        assert not any("drawing" in n and "diagrams" in n for n in names)
        slide_xml = zf.read("ppt/slides/slide1.xml").decode()
        assert "graphicFrame" in slide_xml and "relIds" in slide_xml
        ct = zf.read("[Content_Types].xml").decode()
        assert "diagramData" in ct
    # the preview caveat reached the report
    assert any("PowerPoint" in w for w in ctx.report.warnings)
    Presentation(str(out))  # reopens


def test_two_diagrams_get_distinct_part_indexes(tmp_path):
    out, _ = _render_deck(tmp_path, [_spec("cycle"), _spec("radial")])
    with zipfile.ZipFile(out) as zf:
        names = zf.namelist()
        assert "ppt/diagrams/data1.xml" in names
        assert "ppt/diagrams/data2.xml" in names


def test_hierarchy_nests_under_single_root(tmp_path):
    nodes = [SmartNodeSpec(label="Root"),
             SmartNodeSpec(label="A", children=["A1", "A2"]),
             SmartNodeSpec(label="B")]
    out, _ = _render_deck(tmp_path, [_spec("org_chart", nodes=nodes)])
    with zipfile.ZipFile(out) as zf:
        data = etree.fromstring(zf.read("ppt/diagrams/data1.xml"))
    pts = {p.get("modelId"): p for p in data.iter(f"{{{_DGM}}}pt")}
    texts = {}
    for pid, p in pts.items():
        t = "".join(x.text or "" for x in p.iter(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}t"))
        if t:
            texts[pid] = t
    cxns = [(c.get("srcId"), c.get("destId"))
            for c in data.iter(f"{{{_DGM}}}cxn")]
    doc_id = next(p.get("modelId") for p in pts.values()
                  if p.get("type") == "doc")
    # exactly ONE node hangs off the doc (the root)
    top = [d for s, d in cxns if s == doc_id]
    assert len(top) == 1 and texts[top[0]] == "Root"
    root_children = [texts[d] for s, d in cxns if s == top[0]]
    assert root_children == ["A", "B"]
    a_id = next(d for s, d in cxns if s == top[0] and texts[d] == "A")
    assert [texts[d] for s, d in cxns if s == a_id] == ["A1", "A2"]


def test_colors_recolored_to_theme_primary(tmp_path):
    out, ctx = _render_deck(tmp_path, [_spec("cycle")])
    with zipfile.ZipFile(out) as zf:
        colors = zf.read("ppt/diagrams/colors1.xml").decode()
    assert ctx.theme.color("primary") in colors
    assert 'schemeClr val="accent1"' not in colors


def test_cycle_flat_warns_on_children(tmp_path):
    nodes = [SmartNodeSpec(label="A", children=["x"]),
             SmartNodeSpec(label="B")]
    _, ctx = _render_deck(tmp_path, [_spec("cycle", nodes=nodes)])
    assert any("flat" in w for w in ctx.report.warnings)


def test_clone_failure_falls_back_to_drawn(tmp_path, monkeypatch):
    import deckengine.render.smartart_clone as sc

    def boom(*a, **k):
        raise RuntimeError("template missing")
    monkeypatch.setattr(sc, "add_smart_diagram", boom)
    ctx = make_ctx()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    nodes = [SmartNodeSpec(label="Root"),
             SmartNodeSpec(label="A"), SmartNodeSpec(label="B")]
    get_component("smart_diagram").render(
        slide, _spec("issue_tree", nodes=nodes),
        BBox(0, 0, inch(9), inch(5)), ctx)
    assert any("clone failed" in w for w in ctx.report.warnings)
    texts = " ".join(r.text for s in slide.shapes if s.has_text_frame
                     for p in s.text_frame.paragraphs for r in p.runs)
    assert "Root" in texts  # the drawn tree rendered instead


def test_smartart_toggle_threads_directive(monkeypatch):
    from deckengine.api import main as api_main
    from deckengine.llm.variants import SMARTART_DIRECTIVE
    captured = {}

    def fake_gds(prompt, **kw):
        captured.update(kw)
        from deckengine.schema.slide_types import DeckSpec
        return DeckSpec.model_validate({
            "schema_version": 1, "theme": "consulting_navy",
            "meta": {"title": "t", "date": "x", "footer_org": "d"},
            "slides": [{"slide_type": "bullet_content", "title": "T",
                        "bullets": [{"text": "a"}]}]})
    import deckengine.llm.spec_generator as sg
    monkeypatch.setattr(sg, "generate_deck_spec", fake_gds)
    monkeypatch.setattr(api_main, "_run_render", lambda job_id, spec: None)
    api_main._JOBS["tjob"] = {"status": "running", "owner": "t"}
    monkeypatch.setattr(api_main, "_persist", lambda j: None)
    req = api_main.GenerateFromPrompt(prompt="p", smartart=True,
                                      auto_approve=True)
    api_main._run_generate("tjob", req, outline=None)
    assert captured.get("design_directive") == SMARTART_DIRECTIVE
