"""S4 credibility wiring: speaker notes, per-chart source line, methodology
appendix, legend, and CSV marker injection."""
from __future__ import annotations

from pptx import Presentation

from deckengine.components.base import BuildReport, RenderContext, \
    get_component
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.llm.facts import FactTable
from deckengine.llm.provenance import (LEGEND, append_legend_once,
                                       inject_fact_markers,
                                       methodology_appendix)
from deckengine.render.deck_builder import build_deck
from deckengine.schema.components import ChartSeriesSpec, NativeChartSpec
from deckengine.schema.slide_types import (BulletContentSpec, ChartSlideSpec,
                                           DeckSpec, SectionDividerSpec)


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


# -- speaker notes -----------------------------------------------------------

def test_notes_written_and_readable(tmp_path):
    spec = DeckSpec.model_validate({
        "schema_version": 1, "theme": "consulting_navy",
        "meta": {"title": "n", "date": "13 Aug 2026", "footer_org": "DE"},
        "slides": [
            {"slide_type": "bullet_content", "title": "One",
             "bullets": [{"text": "a"}],
             "notes": "Open with the water number. Point at the right bar."},
            {"slide_type": "section_divider", "title": "Part two",
             "notes": "Pause here; invite questions on part one."},
            {"slide_type": "bullet_content", "title": "Three",
             "bullets": [{"text": "b"}]},
        ]})
    out = tmp_path / "notes.pptx"
    build_deck(spec, out)
    prs = Presentation(str(out))
    assert prs.slides[0].notes_slide.notes_text_frame.text.startswith(
        "Open with the water number")
    assert "Pause here" in prs.slides[1].notes_slide.notes_text_frame.text
    assert not prs.slides[2].has_notes_slide  # no empty notes parts


def test_overlong_notes_trim_instead_of_failing():
    long = ("Point at the water chart and hold there. " * 12)  # ~500 chars
    s = BulletContentSpec(title="T", bullets=[{"text": "a"}], notes=long)
    assert s.notes is not None and len(s.notes) <= 350
    assert s.notes.endswith(".")  # trimmed at a sentence boundary


# -- per-chart source line ---------------------------------------------------

def _chart(**kw):
    kw.setdefault("chart_type", "bar")
    kw.setdefault("categories", ["A", "B"])
    kw.setdefault("series", [ChartSeriesSpec(name="s", values=[3.0, 1.0])])
    kw.setdefault("sort", "desc")
    kw.setdefault("highlight", None)
    kw.setdefault("annotation", None)
    return NativeChartSpec(**kw)


def test_chart_source_measure_render_parity():
    ctx = make_ctx()
    comp = get_component("native_chart")
    data = _chart(annotation="A leads by 3x.",
                  source="Source: FAOSTAT 2025 [[src:official]]")
    width = inch(8)
    measured = comp.measure(data, width, ctx)
    plain_h = comp.measure(_chart(annotation="A leads by 3x."), width, ctx)
    assert measured > plain_h  # source line adds real height
    from pptx import Presentation as P
    prs = P()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    consumed = comp.render(slide, data,
                           BBox(inch(0.5), inch(1), width, inch(6)), ctx)
    assert abs(measured - consumed) <= pt(1)
    runs = [r for s in slide.shapes if s.has_text_frame
            for p in s.text_frame.paragraphs for r in p.runs]
    assert any("FAOSTAT" in r.text for r in runs)
    assert any(r.text == "●" for r in runs)  # the marker glyph rendered


# -- methodology appendix ----------------------------------------------------

def _marked_slide():
    return BulletContentSpec(
        title="Coverage reached 34% [[src:recon]] of irrigated area",
        bullets=[{"text": "Yields rose 18% [[src:est]] on treated plots."}])


def test_methodology_appendix_lists_non_official():
    appendix = methodology_appendix([_marked_slide()])
    assert appendix is not None
    labels = [g.label for g in appendix.table.groups]
    assert "Reconstructed figures" in labels and "Estimates" in labels
    all_rows = [r for g in appendix.table.groups for r in g.rows]
    assert any(r[0] == "34" for r in all_rows)
    assert any(r[0] == "18" for r in all_rows)


def test_methodology_appendix_none_when_all_official():
    s = BulletContentSpec(
        title="Outlay hit $310M [[src:official]] in FY25",
        bullets=[{"text": "Published in the state budget."}])
    assert methodology_appendix([s]) is None


# -- legend -------------------------------------------------------------------

def test_legend_appended_once_and_idempotent():
    slides = [SectionDividerSpec(title="Part"), _marked_slide(),
              _marked_slide()]
    append_legend_once(slides)
    append_legend_once(slides)  # idempotent
    foots = [getattr(s, "footnote", None) for s in slides]
    assert sum(1 for f in foots if f and LEGEND in f) == 1
    assert foots[1] and LEGEND in foots[1]  # divider has no footnote field


# -- CSV marker injection ------------------------------------------------------

def _facts():
    t = FactTable()
    t.add("income", "$1,216", "avg income", "survey.csv")
    t.add("delta", "+2.7%", "income delta", "survey.csv")
    t.add("twelve", "12", "small count", "survey.csv")
    return t


def test_inject_fact_markers_appends_official():
    s = BulletContentSpec(
        title="Income reached $1,216 against a +2.7% delta",
        bullets=[{"text": "Baseline holds at $1,216 statewide."}])
    out = inject_fact_markers(s, _facts())
    assert out.title.count("[[src:official]]") == 2
    assert "$1,216 [[src:official]]" in out.bullets[0].text
    again = inject_fact_markers(out, _facts())  # idempotent
    assert again.title.count("[[src:official]]") == 2


def test_inject_skips_substring_matches():
    t = FactTable()
    t.add("n", "12", "count", "x.csv")
    s = BulletContentSpec(title="Fleet of 312M units moved",
                          bullets=[{"text": "312M is the fleet size."}])
    out = inject_fact_markers(s, t)
    assert "[[src:official]]" not in out.title


def test_chart_slide_source_field_accepts_marker():
    s = ChartSlideSpec(title="T", chart=_chart(
        source="Source: IEA 2025 [[src:official]]"))
    assert s.chart.source.endswith("[[src:official]]")
