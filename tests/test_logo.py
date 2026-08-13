"""Brand logo: rendered top-right on every slide + the upload endpoint."""
import io

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from deckengine.core.units import SLIDE_W_16_9, inch
from deckengine.render.deck_builder import build_deck
from deckengine.schema.slide_types import (BulletContentSpec, DeckMeta,
                                           DeckSpec, TitleSlideSpec)


def _png(tmp_path, name="logo.png"):
    p = tmp_path / name
    Image.new("RGBA", (120, 40), (200, 30, 30, 255)).save(p, "PNG")
    return p


def _spec(logo):
    return DeckSpec(
        theme="consulting_navy",
        meta=DeckMeta(title="Deck", logo=logo, date="Jan 2026"),
        slides=[TitleSlideSpec(title="Cover", date="Jan 2026"),
                BulletContentSpec(title="A point that carries a verb",
                                  bullets=[{"text": "evidence one"}])])


def _pics(slide):
    return [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]


def test_logo_rendered_top_right_on_every_slide(tmp_path):
    logo = _png(tmp_path)
    out = tmp_path / "deck.pptx"
    report = build_deck(_spec(str(logo)), out)
    assert not [w for w in report.warnings if "logo" in w.lower()]
    prs = Presentation(str(out))
    assert len(prs.slides._sldIdLst) == 2
    for slide in prs.slides:
        pics = _pics(slide)
        assert len(pics) == 1, "each slide carries exactly the logo picture"
        pic = pics[0]
        assert pic.left > SLIDE_W_16_9 // 2   # right half
        assert pic.top < inch(1.0)            # top band


def test_missing_logo_warns_and_skips(tmp_path):
    out = tmp_path / "deck.pptx"
    report = build_deck(_spec("nope_does_not_exist.png"), out)
    assert any("logo" in w.lower() for w in report.warnings)
    prs = Presentation(str(out))
    for slide in prs.slides:
        assert _pics(slide) == []  # nothing drawn, render still succeeded


def test_no_logo_no_picture(tmp_path):
    out = tmp_path / "deck.pptx"
    report = build_deck(_spec(None), out)
    assert not [w for w in report.warnings if "logo" in w.lower()]
    prs = Presentation(str(out))
    for slide in prs.slides:
        assert _pics(slide) == []


def test_upload_asset_endpoint(tmp_path, monkeypatch):
    from fastapi.testclient import TestClient

    from deckengine.core import assets
    from deckengine.api.main import app
    monkeypatch.setenv("DECKENGINE_AUTH", "0")  # dev bypass for this test
    monkeypatch.setattr(assets, "ASSETS_DIR", tmp_path)

    buf = io.BytesIO()
    Image.new("RGBA", (60, 20), (0, 90, 160, 255)).save(buf, "PNG")
    client = TestClient(app)
    r = client.post("/assets",
                    files={"file": ("My Logo!.png", buf.getvalue(),
                                    "image/png")})
    assert r.status_code == 200
    asset = r.json()["asset"]
    assert asset == "uploads/My_Logo_.png"        # sanitized
    assert (tmp_path / "uploads" / "My_Logo_.png").is_file()
    assert assets.resolve_asset(asset) is not None  # render can find it
