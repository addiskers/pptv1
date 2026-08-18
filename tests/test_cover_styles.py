"""Cover compositions: four hand-built title-slide styles so a variant
batch never opens five decks on the same-looking slide."""
from __future__ import annotations

from pptx import Presentation

from deckengine.layout.zones import SLIDE
from deckengine.render.deck_builder import build_deck
from deckengine.schema.slide_types import DeckSpec, TitleSlideSpec

_FULL = {
    "slide_type": "title",
    "title": "Indonesia is the strongest near-term premium entry bet",
    "subtitle": "A distributor-led, digital-first launch beats the "
                "alternatives on speed, capital, and reversibility.",
    "date": "18 Aug 2026", "org": "SkyQuest Insights", "wordmark": "SkyQuest",
}

STYLES = ("dark_hero", "split_panel", "light_minimal", "band_statement")


def _build(style, tmp_path):
    spec = DeckSpec.model_validate({
        "schema_version": 1, "theme": "consulting_navy",
        "meta": {"title": "c", "date": "18 Aug 2026", "footer_org": "DE"},
        "slides": [{**_FULL, "style": style}]})
    out = tmp_path / f"{style}.pptx"
    report = build_deck(spec, out)
    return Presentation(str(out)).slides[0], report


def test_default_style_is_dark_hero():
    assert TitleSlideSpec.model_validate(
        {"slide_type": "title", "title": "T"}).style == "dark_hero"


def test_all_styles_render_clean(tmp_path):
    for style in STYLES:
        slide, report = _build(style, tmp_path)
        assert report.warnings == [], style
        texts = " ".join(s.text_frame.text for s in slide.shapes
                         if s.has_text_frame)
        assert "Indonesia" in texts, style
        assert "SkyQuest" in texts, style


def _autoshape_geoms(slide):
    return sorted((int(s.width), int(s.height))
                  for s in slide.shapes
                  if str(s.shape_type) == "AUTO_SHAPE (1)")


def test_styles_are_visually_distinct(tmp_path):
    """Each style paints a DIFFERENT panel geometry — the whole point."""
    geoms = {s: _autoshape_geoms(_build(s, tmp_path)[0]) for s in STYLES}
    full_w, full_h = SLIDE.w, SLIDE.h
    # dark_hero: one full-bleed panel
    assert (full_w, full_h) in geoms["dark_hero"]
    # split_panel: a full-height panel narrower than half the slide + a bit
    assert any(h == full_h and w < full_w * 0.55
               for w, h in geoms["split_panel"])
    # band_statement: a full-width band far shorter than the slide
    assert any(w == full_w and h < full_h * 0.55
               for w, h in geoms["band_statement"])
    # light_minimal: NO full-bleed dark panel at all
    assert not any(h >= full_h * 0.9 for _, h in geoms["light_minimal"])


def test_minimal_spec_renders_every_style(tmp_path):
    """Covers must never break: title-only specs render all four styles."""
    for style in STYLES:
        spec = DeckSpec.model_validate({
            "schema_version": 1, "theme": "consulting_navy",
            "meta": {"title": "c", "date": "18 Aug 2026",
                     "footer_org": "DE"},
            "slides": [{"slide_type": "title", "title": "Just a headline",
                        "style": style}]})
        report = build_deck(spec, tmp_path / f"min_{style}.pptx")
        assert report.warnings == [], style
