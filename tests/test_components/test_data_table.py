from pptx import Presentation

from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.data_table import DataTable
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import (DataColumnSpec, DataGroupSpec,
                                          DataTableSpec)


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def simple_table() -> DataTableSpec:
    """3 columns / 2 groups / 5 total rows — text + number cells only."""
    return DataTableSpec(
        columns=[
            DataColumnSpec(label="Category", frac=0.30),
            DataColumnSpec(label="Product", frac=0.45),
            DataColumnSpec(label="Units sold", frac=0.25, cell_kind="number"),
        ],
        groups=[
            DataGroupSpec(label="Animal Breeds", rows=[
                ["Improved dairy cow", "1,200"],
                ["Dual-purpose chicken", "3,400"],
            ]),
            DataGroupSpec(label="Crops", rows=[
                ["Drought-tolerant maize", "890"],
                ["Iron-rich beans", "455"],
                ["Vitamin A cassava", "210"],
            ]),
        ],
    )


def wide_table() -> DataTableSpec:
    """6 columns / 3 groups with badge, heatmap and dot cells; zebra on."""
    return DataTableSpec(
        columns=[
            DataColumnSpec(label="Value chain", frac=0.18),
            DataColumnSpec(label="Product", frac=0.28),
            DataColumnSpec(label="Type", frac=0.10, cell_kind="badge"),
            DataColumnSpec(label="Readiness", frac=0.16, cell_kind="heatmap"),
            DataColumnSpec(label="Adoption", frac=0.16, cell_kind="heatmap"),
            DataColumnSpec(label="Priority", frac=0.12, cell_kind="dot"),
        ],
        heatmap_lo=0.0,
        heatmap_hi=100.0,
        zebra=True,
        groups=[
            DataGroupSpec(label="Livestock", rows=[
                ["Improved dairy cow", "LIV", "90%", "45%", "x"],
                ["Dual-purpose chicken", "LIV", "0%", "N/A", ""],
            ]),
            DataGroupSpec(label="Crops", rows=[
                ["Drought-tolerant maize", "CRO", "75%", "60%", "x"],
                ["Iron-rich beans", "CRO", "40%", "20%", ""],
                ["Vitamin A cassava"],  # short row: missing values -> ""
            ]),
            DataGroupSpec(label="Digital", rows=[
                ["Advisory app", "ADS", "15%", "5%", ""],
            ]),
        ],
    )


def test_measure_equals_render_simple():
    ctx = make_ctx()
    comp = DataTable()
    data = simple_table()
    width = inch(6)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(6)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_measure_equals_render_wide():
    ctx = make_ctx()
    comp = DataTable()
    data = wide_table()
    width = inch(9)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(inch(0.5), inch(0.5), width, inch(6)), ctx)
    assert abs(measured - consumed) <= pt(1)


def _shape_with_text(slide, text: str):
    for s in slide.shapes:
        if s.has_text_frame and s.text_frame.text == text:
            return s
    raise AssertionError(f"no shape with text {text!r}")


def test_heatmap_fill_differs_between_low_and_high():
    ctx = make_ctx()
    _, slide = blank_slide()
    DataTable().render(slide, wide_table(), BBox(0, 0, inch(9), inch(6)), ctx)
    low = _shape_with_text(slide, "0%")
    high = _shape_with_text(slide, "90%")
    assert str(low.fill.fore_color.rgb) != str(high.fill.fore_color.rgb)
    assert str(high.fill.fore_color.rgb) == \
        ctx.theme.heatmap_color(90.0, 0.0, 100.0)
    # high-value cell flips to inverse ink; low-value keeps ink
    high_run = high.text_frame.paragraphs[0].runs[0]
    low_run = low.text_frame.paragraphs[0].runs[0]
    assert str(high_run.font.color.rgb) == ctx.theme.color("inverse_ink")
    assert str(low_run.font.color.rgb) == ctx.theme.color("ink")


def test_merged_group_label_fewer_rects_than_rows():
    ctx = make_ctx()
    _, slide = blank_slide()
    data = simple_table()
    DataTable().render(slide, data, BBox(0, 0, inch(6), inch(6)), ctx)
    col0_rects = [s for s in slide.shapes if s.left == 0]
    total_rows = sum(len(g.rows) for g in data.groups)
    # one header cell + ONE merged rect per group, not one per data row
    assert len(col0_rects) == 1 + len(data.groups)
    assert len(col0_rects) < 1 + total_rows
    labels = {s.text_frame.text for s in col0_rects if s.has_text_frame}
    assert {"Animal Breeds", "Crops"} <= labels


def test_row_offsets_arithmetic():
    ctx = make_ctx()
    comp = DataTable()
    data = wide_table()
    width = inch(9)
    offsets = comp.row_offsets(data, width, ctx)
    assert len(offsets) == len(data.groups) + 1
    assert offsets == sorted(offsets)
    assert offsets[-1] == comp.measure(data, width, ctx)
    total_rows = sum(len(g.rows) for g in data.groups)
    row_h = (offsets[-1] - offsets[0]) // total_rows
    for i, group in enumerate(data.groups):
        assert offsets[i + 1] - offsets[i] == len(group.rows) * row_h


def test_badge_pill_uses_palette_fill():
    ctx = make_ctx()
    _, slide = blank_slide()
    DataTable().render(slide, wide_table(), BBox(0, 0, inch(9), inch(6)), ctx)
    pills = [s for s in slide.shapes
             if s.has_text_frame and s.text_frame.text == "LIV"]
    assert len(pills) == 2  # one per LIV row
    for p in pills:
        assert str(p.fill.fore_color.rgb) == ctx.theme.badge_palette["LIV"]
        run = p.text_frame.paragraphs[0].runs[0]
        assert run.font.bold
        assert str(run.font.color.rgb) == ctx.theme.color("inverse_ink")


def test_heatmap_autoscales_to_clustered_data():
    """Default (0,100) bounds + values clustered 1-28% must auto-scale so the
    shading has contrast (the review's 'nearly white shares' defect)."""
    from deckengine.components.data_table import DataTable
    from deckengine.schema.components import (DataColumnSpec, DataGroupSpec,
                                              DataTableSpec)
    comp = DataTable()
    spec = DataTableSpec(
        columns=[
            DataColumnSpec(label="Group", frac=0.3),
            DataColumnSpec(label="Player", frac=0.4),
            DataColumnSpec(label="Share", frac=0.3, cell_kind="heatmap"),
        ],
        groups=[DataGroupSpec(label="A", rows=[["x", "1%"], ["y", "18%"],
                                               ["z", "28%"]])],
    )
    assert comp._heat_bounds(spec) == (1.0, 28.0)
    # explicit bounds are respected untouched
    spec2 = spec.model_copy(update={"heatmap_lo": 10.0, "heatmap_hi": 50.0})
    assert comp._heat_bounds(spec2) == (10.0, 50.0)
