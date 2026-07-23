from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.harvey_balls import HarveyBalls
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import HarveyBallsSpec, HarveyItemSpec


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _autoshapes(slide, mso_shape):
    return [s for s in slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and s.auto_shape_type == mso_shape]


def _fill_hex(shape) -> str | None:
    try:
        return str(shape.fill.fore_color.rgb)
    except (TypeError, AttributeError):
        return None  # background (no-fill) shapes have no fore_color


def three_items() -> HarveyBallsSpec:
    return HarveyBallsSpec(items=[
        HarveyItemSpec(label="Cost to serve", score=1),
        HarveyItemSpec(label="Speed to market", score=3),
        HarveyItemSpec(label="Regulatory fit", score=4),
    ])


def six_items() -> HarveyBallsSpec:
    return HarveyBallsSpec(items=[
        HarveyItemSpec(label="Cost", score=0),
        HarveyItemSpec(label="Speed", score=1),
        HarveyItemSpec(label="Delivery risk", score=2),
        HarveyItemSpec(label="Scalability of the operating model", score=3),
        HarveyItemSpec(label="Strategic fit", score=4),
        HarveyItemSpec(label="Talent depth", score=2),
    ])


def test_measure_equals_render_three_items():
    ctx = make_ctx()
    comp = HarveyBalls()
    data = three_items()
    width = inch(9)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(2)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_measure_equals_render_six_items():
    ctx = make_ctx()
    comp = HarveyBalls()
    data = six_items()
    width = inch(12)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(2)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_score_zero_draws_no_pie_and_no_primary_fill():
    ctx = make_ctx()
    _, slide = blank_slide()
    HarveyBalls().render(slide, HarveyBallsSpec(items=[
        HarveyItemSpec(label="Cost", score=0),
        HarveyItemSpec(label="Risk", score=0),
    ]), BBox(0, 0, inch(6), inch(1.5)), ctx)
    assert not _autoshapes(slide, MSO_SHAPE.PIE)
    ovals = _autoshapes(slide, MSO_SHAPE.OVAL)
    assert len(ovals) == 2  # one unfilled outline per ball, nothing under
    primary = ctx.theme.color("primary")
    assert all(_fill_hex(o) != primary for o in ovals)


def test_score_four_solid_primary_oval_under_outline():
    ctx = make_ctx()
    _, slide = blank_slide()
    HarveyBalls().render(slide, HarveyBallsSpec(items=[
        HarveyItemSpec(label="Strategic fit", score=4),
        HarveyItemSpec(label="Cost", score=0),
    ]), BBox(0, 0, inch(6), inch(1.5)), ctx)
    assert not _autoshapes(slide, MSO_SHAPE.PIE)
    ovals = _autoshapes(slide, MSO_SHAPE.OVAL)
    fills = [_fill_hex(o) for o in ovals]
    assert fills.count(ctx.theme.color("primary")) == 1
    # fill drawn UNDER the rim: primary oval precedes its unfilled outline
    filled = fills.index(ctx.theme.color("primary"))
    outline = ovals[filled + 1]
    assert _fill_hex(outline) is None
    assert (outline.left, outline.top) == (ovals[filled].left,
                                           ovals[filled].top)


def test_score_two_exactly_one_pie():
    ctx = make_ctx()
    _, slide = blank_slide()
    HarveyBalls().render(slide, HarveyBallsSpec(items=[
        HarveyItemSpec(label="Delivery risk", score=2),
        HarveyItemSpec(label="Cost", score=0),
    ]), BBox(0, 0, inch(6), inch(1.5)), ctx)
    pies = _autoshapes(slide, MSO_SHAPE.PIE)
    assert len(pies) == 1
    assert _fill_hex(pies[0]) == ctx.theme.color("primary")


def test_pie_adjustment_xml_quarters():
    """Wedge = 12 o'clock start, one clockwise quarter-turn per score point.

    Raw OOXML pie angles are 60000ths of a degree AND the preset geometry
    pins adjustments to [0, 21599999] — negative angles are clamped to 0
    (the full-circle bug the Q3 probe caught). Verify scores 1-3 write the
    NORMALIZED angles: start = 270 deg, end = (270 + score * 90) mod 360.
    """
    ctx = make_ctx()
    _, slide = blank_slide()
    HarveyBalls().render(slide, HarveyBallsSpec(items=[
        HarveyItemSpec(label="Cost", score=1),
        HarveyItemSpec(label="Risk", score=2),
        HarveyItemSpec(label="Fit", score=3),
    ]), BBox(0, 0, inch(9), inch(2)), ctx)
    pies = _autoshapes(slide, MSO_SHAPE.PIE)
    assert len(pies) == 3
    for score, pie in zip((1, 2, 3), pies):
        av_lst = pie._element.spPr.find(qn("a:prstGeom")).find(qn("a:avLst"))
        raw = {gd.get("name"): int(gd.get("fmla").split()[1])
               for gd in av_lst.findall(qn("a:gd"))}
        assert raw["adj1"] == round(270 * 60000)
        assert raw["adj2"] == round(((270 + score * 90) % 360) * 60000)
        assert 0 <= raw["adj1"] <= 21599999
        assert 0 <= raw["adj2"] <= 21599999


def test_labels_present_and_all_outline_ovals_same_size():
    ctx = make_ctx()
    _, slide = blank_slide()
    HarveyBalls().render(slide, three_items(),
                         BBox(0, 0, inch(9), inch(2)), ctx)

    text = " ".join("".join(r.text for r in p.runs)
                    for s in slide.shapes if s.has_text_frame
                    for p in s.text_frame.paragraphs)
    assert "Cost to serve" in text
    assert "Speed to market" in text
    assert "Regulatory fit" in text

    ovals = _autoshapes(slide, MSO_SHAPE.OVAL)
    outlines = [o for o in ovals if _fill_hex(o) is None]
    assert len(outlines) == 3  # one crisp rim per ball
    sizes = {(int(o.width), int(o.height)) for o in outlines}
    assert sizes == {(inch(0.32), inch(0.32))}


# --- wedge angle regression (PIE adjustments pinned to [0, 360)) --------------

def test_score_wedges_use_normalized_angles():
    from deckengine.schema.components import HarveyBallsSpec, HarveyItemSpec
    ctx = make_ctx()
    comp = HarveyBalls()
    # min_length=2 guard: pair the scored item with a zero-score sibling
    data = HarveyBallsSpec(items=[HarveyItemSpec(label="two", score=2),
                                  HarveyItemSpec(label="zero", score=0)])
    _, slide = blank_slide()
    comp.render(slide, data, BBox(0, 0, inch(6), inch(1)), ctx)
    from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
    pies = [s for s in slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and s.auto_shape_type == MSO_SHAPE.PIE]
    assert len(pies) == 1
    start, end = pies[0].adjustments[0], pies[0].adjustments[1]
    assert abs(start - 162.0) < 0.01   # 12 o'clock => 270 deg normalized
    assert abs(end - 54.0) < 0.01      # half sweep ends at 6 o'clock (90 deg)
