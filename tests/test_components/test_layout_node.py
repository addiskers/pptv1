from pptx import Presentation

from deckengine.components.base import (BuildReport, RenderContext,
                                        get_component)
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import (StatRowSpec, StatSpec,
                                          TextBlockSpec)
from deckengine.schema.layout_tree import ColsNode, PanelNode, RowsNode


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def text(t="Some content line that wraps once or twice at narrow widths"):
    return TextBlockSpec(text=t)


def stats():
    return StatRowSpec(stats=[StatSpec(label="Direct", value="3M"),
                              StatSpec(label="Digital", value="66M")])


# --- parity (fill_hint off => natural height everywhere) ----------------------

def test_rows_measure_equals_render():
    ctx = make_ctx()
    comp = get_component("rows")
    data = RowsNode(children=[text(), stats(), text("second block")])
    width = inch(8)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(6)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_cols_measure_equals_render_and_track_is_max():
    ctx = make_ctx()
    comp = get_component("cols")
    long = text("A much longer text that will wrap over several lines when "
                "constrained to half of the available width in its column")
    data = ColsNode(children=[text("short"), long])
    width = inch(8)
    measured = comp.measure(data, width, ctx)
    h_short = get_component("text_block").measure(
        text("short"), (width - ctx.theme.spacing(0.4)) // 2, ctx)
    assert measured > h_short  # track height set by the tall child
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(4)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_panel_measure_equals_render_and_draws_border():
    ctx = make_ctx()
    comp = get_component("panel")
    data = PanelNode(child=text(), fill_role="surface",
                     border_role="positive", inset=0.8)
    width = inch(6)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(3)), ctx)
    assert abs(measured - consumed) <= pt(1)
    inner = get_component("text_block").measure(
        text(), width - 2 * ctx.theme.spacing(0.8), ctx)
    assert measured == inner + 2 * ctx.theme.spacing(0.8)
    rects = [s for s in slide.shapes if not s.has_text_frame or s.width == width]
    assert any(s.width == width for s in slide.shapes)


def test_nested_tree_parity():
    ctx = make_ctx()
    comp = get_component("rows")
    data = RowsNode(children=[
        ColsNode(children=[text("left cell"), text("right cell")]),
        PanelNode(fill_role="surface",
                  child=ColsNode(children=[stats(), text("note")])),
    ])
    width = inch(9)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(6)), ctx)
    assert abs(measured - consumed) <= pt(1)


# --- fill mode -----------------------------------------------------------------

def test_rows_fracs_fill_assigned_zone():
    ctx = make_ctx()
    comp = get_component("rows")
    data = RowsNode(fracs=[0.7, 0.3], children=[text(), stats()])
    bbox = BBox(0, 0, inch(9), inch(5))
    _, slide = blank_slide()
    ctx.fill_hint = True
    try:
        consumed = comp.render(slide, data, bbox, ctx)
    finally:
        ctx.fill_hint = False
    assert consumed == bbox.h  # bespoke mode fills the whole zone


def test_rows_fracs_natural_without_fill_hint():
    ctx = make_ctx()
    comp = get_component("rows")
    data = RowsNode(fracs=[0.7, 0.3], children=[text(), stats()])
    measured = comp.measure(data, inch(9), ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, inch(9), inch(5)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_rows_pin_last_anchors_bottom():
    ctx = make_ctx()
    comp = get_component("rows")
    data = RowsNode(pin_last=True, children=[text("body content"), stats()])
    bbox = BBox(0, inch(1), inch(8), inch(4))
    nat_last = get_component("stat_row").measure(stats(), bbox.w, ctx)
    _, slide = blank_slide()
    ctx.fill_hint = True
    try:
        consumed = comp.render(slide, data, bbox, ctx)
    finally:
        ctx.fill_hint = False
    assert consumed == bbox.h
    # some shape from the pinned child starts at the anchored offset
    pin_y = bbox.bottom - nat_last
    assert any(abs(s.top - pin_y) <= inch(0.06) for s in slide.shapes)


def test_cols_frac_widths():
    ctx = make_ctx()
    comp = get_component("cols")
    data = ColsNode(fracs=[0.6, 0.4], children=[text("a"), text("b")])
    bbox = BBox(0, 0, inch(10), inch(2))
    boxes = comp._boxes(data, bbox, ctx)
    gap = ctx.theme.spacing(0.4)
    avail = bbox.w - gap
    assert boxes[0].w == round(avail * 0.6)
    assert boxes[1].w == round(avail * 0.4)
    assert boxes[1].x == boxes[0].w + gap


def test_cols_middle_align_centers_children():
    ctx = make_ctx()
    comp = get_component("cols")
    long = text("A much longer text that wraps over many lines when squeezed "
                "into a narrow column so the track grows tall around it and "
                "the short sibling needs vertical centering to look right")
    data = ColsNode(align="middle", children=[text("short"), long])
    bbox = BBox(0, inch(1), inch(8), inch(4))
    _, slide = blank_slide()
    comp.render(slide, data, bbox, ctx)
    # the short child's textbox must start below the track top
    tops = sorted(s.top for s in slide.shapes if s.has_text_frame)
    assert tops[-1] > bbox.y  # someone got offset down
