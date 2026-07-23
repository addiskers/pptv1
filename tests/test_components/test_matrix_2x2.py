"""matrix_2x2 tests — shared track height, gutter labels, dashed midlines."""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.matrix_2x2 import Matrix2x2
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import Matrix2x2Spec, QuadrantSpec

Y_GUTTER = inch(0.32)
X_GUTTER = inch(0.30)


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def shape_text(shape) -> str:
    """Reassemble visible text: wrap tokens carry their own spaces."""
    return "".join(r.text for p in shape.text_frame.paragraphs
                   for r in p.runs)


def slide_text(slide) -> str:
    return " ".join(shape_text(s) for s in slide.shapes if s.has_text_frame)


def short_spec(highlight: int | None = None) -> Matrix2x2Spec:
    return Matrix2x2Spec(
        x_label="Effort to implement →",
        y_label="Impact on yield →",
        quadrants=[
            QuadrantSpec(title="Quick wins",
                         items=["Seed treatment drives", "SMS advisories"]),
            QuadrantSpec(title="Big bets",
                         items=["Irrigation retrofits", "FPO aggregation"]),
            QuadrantSpec(title="Fill-ins", items=["District field days"]),
            QuadrantSpec(title="Deprioritise", items=["Blanket subsidies"]),
        ],
        highlight=highlight)


def long_spec() -> Matrix2x2Spec:
    long_item = ("a deliberately long line of quadrant evidence that must "
                 "wrap across the panel width more than once")
    return Matrix2x2Spec(
        x_label="Ease of delivery →",
        y_label="Contribution to income →",
        quadrants=[
            QuadrantSpec(title="**Scale now** — proven economics",
                         items=[long_item, long_item, long_item, long_item]),
            QuadrantSpec(title="Pilot with co-investment",
                         items=[long_item, long_item]),
            QuadrantSpec(title="Bundle with ongoing programs",
                         items=[long_item]),
            QuadrantSpec(title="Park until evidence improves",
                         items=[long_item, long_item, long_item]),
        ],
        highlight=0)


# --- measure == render ---------------------------------------------------------

def test_measure_equals_render_short():
    ctx = make_ctx()
    comp = Matrix2x2()
    data = short_spec()
    width = inch(8)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data,
                           BBox(inch(0.5), inch(1.2), width, inch(5.5)), ctx)
    assert measured == consumed  # same _plan arithmetic — exact
    assert abs(measured - consumed) <= pt(1)


def test_measure_equals_render_long_wrapping_content():
    ctx = make_ctx()
    comp = Matrix2x2()
    data = long_spec()
    width = inch(6)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data,
                           BBox(inch(0.4), inch(1), width, inch(6)), ctx)
    assert measured == consumed
    assert abs(measured - consumed) <= pt(1)


# --- grid geometry ---------------------------------------------------------------

def test_four_surface_panels_share_one_track_height():
    ctx = make_ctx()
    comp = Matrix2x2()
    bbox = BBox(inch(0.5), inch(1.2), inch(8), inch(5.5))
    _, slide = blank_slide()
    consumed = comp.render(slide, short_spec(), bbox, ctx)

    panels = [s for s in slide.shapes
              if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert len(panels) == 4
    heights = {int(s.height) for s in panels}
    assert len(heights) == 1  # one shared track height
    track_h = heights.pop()
    gap = ctx.theme.spacing(0.35)
    assert 2 * track_h + gap + gap + X_GUTTER == consumed

    surface = ctx.theme.color("surface")
    assert all(str(s.fill.fore_color.rgb) == surface for s in panels)
    # grid sits right of the y-label gutter, two rows x two columns
    assert min(s.left for s in panels) == bbox.x + Y_GUTTER
    assert len({int(s.top) for s in panels}) == 2
    assert len({int(s.left) for s in panels}) == 2


def test_highlight_accents_exactly_one_panel():
    ctx = make_ctx()
    comp = Matrix2x2()
    bbox = BBox(inch(0.5), inch(1.2), inch(8), inch(5.5))
    _, slide = blank_slide()
    comp.render(slide, short_spec(highlight=1), bbox, ctx)

    panels = [s for s in slide.shapes
              if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    accent = ctx.theme.color("accent")
    accented = []
    for s in panels:
        try:
            if str(s.line.color.rgb) == accent:
                accented.append(s)
        except (TypeError, AttributeError):
            pass
    assert len(accented) == 1
    # index 1 = top-right: top row, right column
    hit = accented[0]
    assert hit.top == min(p.top for p in panels)
    assert hit.left == max(p.left for p in panels)


def test_fill_hint_stretches_tracks_within_bbox():
    ctx = make_ctx()
    comp = Matrix2x2()
    data = short_spec()
    width = inch(8)
    natural = comp.measure(data, width, ctx)
    ctx.fill_hint = True
    bbox = BBox(inch(0.5), inch(0.8), width, inch(6.2))
    _, slide = blank_slide()
    consumed = comp.render(slide, data, bbox, ctx)
    assert natural < consumed <= bbox.h
    panels = [s for s in slide.shapes
              if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert len({int(s.height) for s in panels}) == 1  # tracks stay equal


# --- axis labels ------------------------------------------------------------------

def test_rotated_y_label_and_x_label_present():
    ctx = make_ctx()
    comp = Matrix2x2()
    _, slide = blank_slide()
    comp.render(slide, short_spec(),
                BBox(inch(0.5), inch(1.2), inch(8), inch(5.5)), ctx)

    rotated = [s for s in slide.shapes
               if s.has_text_frame and s.rotation == 270]
    assert len(rotated) == 1
    assert "Impact on yield" in shape_text(rotated[0])

    flat_text = " ".join(shape_text(s) for s in slide.shapes
                         if s.has_text_frame and s.rotation == 0)
    assert "Effort to implement" in flat_text


# --- dashed midlines ---------------------------------------------------------------

def test_dashed_midlines_span_grid_only():
    ctx = make_ctx()
    comp = Matrix2x2()
    bbox = BBox(inch(0.5), inch(1.2), inch(8), inch(5.5))
    _, slide = blank_slide()
    consumed = comp.render(slide, short_spec(), bbox, ctx)

    lines = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.LINE]
    assert len(lines) == 2
    for s in lines:
        dash = s.line._get_or_add_ln().find(qn("a:prstDash"))
        assert dash is not None and dash.get("val") == "dash"

    gap = ctx.theme.spacing(0.35)
    grid_h = consumed - gap - X_GUTTER
    hline = next(s for s in lines if s.height == 0)
    vline = next(s for s in lines if s.width == 0)
    assert hline.left == bbox.x + Y_GUTTER
    assert hline.width == bbox.w - Y_GUTTER
    assert vline.top == bbox.y
    assert vline.height == grid_h
    assert bbox.x + Y_GUTTER < vline.left < bbox.right


# --- content -------------------------------------------------------------------------

def test_all_quadrant_titles_and_bullets_in_runs():
    ctx = make_ctx()
    comp = Matrix2x2()
    data = short_spec()
    _, slide = blank_slide()
    comp.render(slide, data,
                BBox(inch(0.5), inch(1.2), inch(8), inch(5.5)), ctx)

    text = slide_text(slide)
    for title in ("Quick wins", "Big bets", "Fill-ins", "Deprioritise"):
        assert title in text
    assert "Seed treatment drives" in text
    assert "•" in text
