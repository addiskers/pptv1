from pptx import Presentation

from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.text_block import TextBlock
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import TextBlockSpec


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def test_measure_equals_render():
    ctx = make_ctx()
    comp = TextBlock()
    data = TextBlockSpec(text="A paragraph with **bold numbers 42%** that wraps "
                              "across a couple of lines at body size.")
    width = inch(3)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(3)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_render_writes_runs_with_bold():
    ctx = make_ctx()
    _, slide = blank_slide()
    TextBlock().render(slide, TextBlockSpec(text="norm **bold** tail"),
                       BBox(0, 0, inch(4), inch(1)), ctx)
    runs = [r for shape in slide.shapes for p in shape.text_frame.paragraphs
            for r in p.runs]
    assert any(r.font.bold for r in runs)
    assert "".join(r.text for r in runs).startswith("norm")


def test_truncation_reported():
    ctx = make_ctx()
    _, slide = blank_slide()
    TextBlock().render(slide, TextBlockSpec(text="word " * 150),
                       BBox(0, 0, inch(1.5), inch(0.4)), ctx)
    assert ctx.report.truncations
