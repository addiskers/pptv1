"""comparison_columns tests — row-track grid alignment is the whole point."""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE

# Child components resolve through the registry; importing the modules here
# (never inside the component module) triggers registration.
import deckengine.components.text_block  # noqa: F401

try:
    import deckengine.components.stat_row  # noqa: F401
except ImportError:  # real component not built yet — stub registered below
    pass
try:
    import deckengine.components.mini_table  # noqa: F401
except ImportError:
    pass

from deckengine.components.base import (BuildReport, Component, RenderContext,
                                        get_component, register)
from deckengine.components.comparison_columns import ComparisonColumns
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import Span, TextMeasurer
from deckengine.core.pptx_shapes import add_shape
from deckengine.core.pptx_text import (add_text_box, make_text_frame,
                                       write_spans_paragraph)
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import (ComparisonColumnSpec,
                                          ComparisonColumnsSpec, MiniTableSpec,
                                          StatRowSpec, StatSpec, TextBlockSpec)


# --- registry stand-ins (used only until the real components land) ----------

class _StubStatRow(Component):
    spec_model = StatRowSpec

    def _h(self, ctx: RenderContext) -> int:
        value_h = ctx.measurer.line_height_emu(
            [Span("0", bold=True)], ctx.font("body"), ctx.size("stat"))
        label_h = ctx.measurer.line_height_emu(
            [Span("x")], ctx.font("body"), ctx.size("micro"))
        return value_h + label_h

    def measure(self, data, width, ctx):
        return self._h(ctx)

    def render(self, slide, data, bbox, ctx):
        h = self._h(ctx)
        cells = bbox.with_height(h).cols(len(data.stats),
                                         gap=ctx.theme.spacing(0.3))
        for cell, stat in zip(cells, data.stats):
            box = add_text_box(slide, cell, align="center")
            tf = box.text_frame
            write_spans_paragraph(tf, [Span(stat.value, bold=True)],
                                  ctx.size("stat"), ctx.theme,
                                  family=ctx.font("body"), align="center")
            p = tf.add_paragraph()
            write_spans_paragraph(tf, [Span(stat.label)], ctx.size("micro"),
                                  ctx.theme, family=ctx.font("body"),
                                  align="center", paragraph=p)
        return h


class _StubMiniTable(Component):
    spec_model = MiniTableSpec

    def _row_h(self, ctx: RenderContext) -> int:
        return (ctx.measurer.line_height_emu([Span("Hg")], ctx.font("body"),
                                             ctx.size("micro"))
                + ctx.theme.spacing(0.3))

    def measure(self, data, width, ctx):
        return self._row_h(ctx) * (1 + len(data.rows))

    def render(self, slide, data, bbox, ctx):
        row_h = self._row_h(ctx)
        y = bbox.y
        ncols = len(data.headers)
        for r, texts in enumerate([data.headers] + data.rows):
            row_box = BBox(bbox.x, y, bbox.w, row_h)
            for cell_box, text in zip(row_box.cols(ncols), texts):
                s = add_shape(slide, cell_box, ctx.theme, shape="rect",
                              fill_role="primary" if r == 0 else "surface")
                tf = make_text_frame(s, align="center", anchor="middle")
                write_spans_paragraph(
                    tf,
                    [Span(text, color_role="inverse_ink" if r == 0 else "ink")],
                    ctx.size("micro"), ctx.theme, family=ctx.font("body"),
                    align="center")
            y += row_h
        return row_h * (1 + len(data.rows))


def _ensure_registered(kind: str, cls: type[Component]) -> None:
    try:
        get_component(kind)
    except KeyError:
        register(kind)(cls)


_ensure_registered("stat_row", _StubStatRow)
_ensure_registered("mini_table", _StubMiniTable)


# --- helpers -----------------------------------------------------------------

def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _state_column(header: str, area: str, context: str,
                  constraint: str) -> ComparisonColumnSpec:
    return ComparisonColumnSpec(header=header, cells=[
        StatRowSpec(stats=[StatSpec(label="Area (m ha)", value=area),
                           StatSpec(label="Yield (t/ha)", value="2.4")]),
        TextBlockSpec(text=context, size_role="small"),
        TextBlockSpec(text=constraint, size_role="small"),
        MiniTableSpec(headers=["Crop", "Share"],
                      rows=[["Rice", "48%"], ["Wheat", "31%"]]),
    ])


def reference_spec() -> ComparisonColumnsSpec:
    """Mirrors the UP / Bihar / Odisha reference slide body."""
    return ComparisonColumnsSpec(
        columns=[
            _state_column("UTTAR PRADESH", "5.9",
                          "Largest rice-wheat belt; **irrigated** systems "
                          "dominate the western districts.",
                          "Groundwater stress and late sowing compress the "
                          "rabi window."),
            _state_column("BIHAR", "3.2",
                          "Smallholder-dense flood plains; kharif rice "
                          "under **rainfed** risk.",
                          "Fragmented plots limit mechanization uptake."),
            _state_column("ODISHA", "4.1",
                          "Coastal rice systems exposed to cyclonic "
                          "losses most seasons.",
                          "Upland tracts remain low-input with weak market "
                          "links."),
        ],
        row_labels=["Scale", "Context", "Constraints", "Crop mix"],
        header_fill_role="primary",
        dashed_separators=True,
    )


# --- measure == render (2 fixtures) ------------------------------------------

def test_measure_equals_render_reference_three_columns():
    ctx = make_ctx()
    comp = ComparisonColumns()
    spec = reference_spec()
    width = inch(9)
    measured = comp.measure(spec, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, spec, BBox(inch(0.5), inch(1.2), width, inch(5.8)),
                           ctx)
    assert measured == consumed  # same track arithmetic — exact
    assert abs(measured - consumed) <= pt(1)


def test_measure_equals_render_two_columns_uneven_cells():
    ctx = make_ctx()
    comp = ComparisonColumns()
    spec = ComparisonColumnsSpec(
        columns=[
            ComparisonColumnSpec(header="WITH", cells=[
                TextBlockSpec(text="Program districts saw yields climb."),
                MiniTableSpec(headers=["Metric", "Value"],
                              rows=[["Adoption", "62%"]]),
                TextBlockSpec(text="Extension access improved.",
                              size_role="small"),
            ]),
            ComparisonColumnSpec(header="WITHOUT", cells=[
                TextBlockSpec(text="Control districts held flat."),
                MiniTableSpec(headers=["Metric", "Value"],
                              rows=[["Adoption", "18%"]]),
            ]),  # one fewer cell: last track stays empty in this column
        ],
        row_labels=None,
        dashed_separators=True,
    )
    width = inch(6)
    measured = comp.measure(spec, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, spec, BBox(inch(0.5), inch(1), width, inch(5)),
                           ctx)
    assert measured == consumed
    assert abs(measured - consumed) <= pt(1)


# --- CRITICAL: row tracks never desynchronize --------------------------------

def test_mini_tables_land_on_same_y_despite_uneven_text():
    """One column's text_block wraps far longer than the other's; the shared
    row track must still put both columns' mini_table cells at the same y."""
    ctx = make_ctx()
    comp = ComparisonColumns()
    long_text = ("Bihar's smallholder base is deliberately described at much "
                 "greater length here so that this cell wraps across many "
                 "more lines than its neighbour, forcing the shared row "
                 "track to grow to the tallest column's measured height.")
    spec = ComparisonColumnsSpec(
        columns=[
            ComparisonColumnSpec(header="UP", cells=[
                TextBlockSpec(text="Short note."),
                MiniTableSpec(headers=["A", "B"], rows=[["1", "2"]]),
            ]),
            ComparisonColumnSpec(header="BIHAR", cells=[
                TextBlockSpec(text=long_text),
                MiniTableSpec(headers=["A", "B"], rows=[["1", "2"]]),
            ]),
        ],
        row_labels=None,
        dashed_separators=False,
    )
    bbox = BBox(inch(0.5), inch(1), inch(6), inch(5.5))
    _, slide = blank_slide()
    comp.render(slide, spec, bbox, ctx)

    # mini_table cells are the only autoshape rects below the header bands
    rects = [s for s in slide.shapes
             if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and s.top > bbox.y]
    assert rects, "mini_table rendered no cell rects"
    mid = bbox.x + bbox.w // 2
    col0_min_y = min(s.top for s in rects if s.left < mid)
    col1_min_y = min(s.top for s in rects if s.left >= mid)
    assert col0_min_y == col1_min_y, (
        f"mini_table rows desynchronized: {col0_min_y} != {col1_min_y}")

    # sanity: the long column really did wrap longer than the short one, so
    # the shared track y sits well below a single short line
    short_h = get_component("text_block").measure(spec.columns[0].cells[0],
                                                  inch(2.9), ctx)
    long_h = get_component("text_block").measure(spec.columns[1].cells[0],
                                                 inch(2.9), ctx)
    assert long_h > short_h * 2
    assert col0_min_y >= bbox.y + long_h  # track grew to the tallest cell


# --- behavior: rail pentagons + dashed separators -----------------------------

def test_rail_pentagons_and_dashed_separators():
    ctx = make_ctx()
    comp = ComparisonColumns()
    spec = reference_spec()
    bbox = BBox(inch(0.4), inch(1.2), inch(9), inch(5.8))
    _, slide = blank_slide()
    consumed = comp.render(slide, spec, bbox, ctx)

    pentagons = [s for s in slide.shapes
                 if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                 and s.auto_shape_type == MSO_SHAPE.PENTAGON]
    assert len(pentagons) == 4  # one per row track
    assert all(p.height <= inch(0.62) for p in pentagons)
    rail_right = bbox.x + min(round(bbox.w * 0.14), inch(1.6))
    assert all(p.left == bbox.x and p.left + p.width <= rail_right
               for p in pentagons)

    # vertical dashed separators: full height (header top -> last track
    # bottom), centered in the two gaps between the three columns
    vlines = [s for s in slide.shapes
              if s.shape_type == MSO_SHAPE_TYPE.LINE
              and s.width == 0 and s.height == consumed]
    assert len(vlines) == 2
    assert all(s.top == bbox.y for s in vlines)
    assert all(rail_right < s.left < bbox.x + bbox.w for s in vlines)
