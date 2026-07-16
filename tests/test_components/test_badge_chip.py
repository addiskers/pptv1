from pptx import Presentation

from deckengine.components.badge_chip import BadgeChip
from deckengine.components.base import BuildReport, RenderContext
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import BadgeChipSpec


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def test_measure_equals_render_known_code():
    ctx = make_ctx()
    comp = BadgeChip()
    data = BadgeChipSpec(code="LIV")
    width = inch(2)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(1)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_measure_equals_render_unknown_code():
    ctx = make_ctx()
    comp = BadgeChip()
    data = BadgeChipSpec(code="ZZZ")
    width = inch(1.5)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(1)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_palette_fill_and_inverse_bold_text():
    ctx = make_ctx()
    _, slide = blank_slide()
    BadgeChip().render(slide, BadgeChipSpec(code="CRO"),
                       BBox(0, 0, inch(2), inch(1)), ctx)
    chip = slide.shapes[0]
    assert str(chip.fill.fore_color.rgb) == ctx.theme.badge_palette["CRO"]
    runs = [r for p in chip.text_frame.paragraphs for r in p.runs]
    assert [r.text for r in runs] == ["CRO"]
    assert runs[0].font.bold
    assert str(runs[0].font.color.rgb) == ctx.theme.color("inverse_ink")
    assert not ctx.report.warnings
    # pill hugs its text: narrower than the generous bbox
    assert chip.width < inch(2)


def test_unknown_code_warns_and_falls_back_to_muted():
    ctx = make_ctx()
    _, slide = blank_slide()
    BadgeChip().render(slide, BadgeChipSpec(code="ZZZ"),
                       BBox(0, 0, inch(2), inch(1)), ctx)
    assert ctx.report.warnings
    chip = slide.shapes[0]
    assert str(chip.fill.fore_color.rgb) == ctx.theme.color("ink_muted")


def test_width_capped_at_bbox_width():
    ctx = make_ctx()
    _, slide = blank_slide()
    narrow = pt(10)
    BadgeChip().render(slide, BadgeChipSpec(code="DIR"),
                       BBox(0, 0, narrow, inch(1)), ctx)
    assert slide.shapes[0].width == narrow
