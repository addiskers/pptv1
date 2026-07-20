from pptx import Presentation
from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.icon_tile_row import IconTileRow
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import IconTileRowSpec, IconTileSpec


def make_ctx():
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def spec(n=4):
    return IconTileRowSpec(tiles=[
        IconTileSpec(icon="people", stat=f"**{i*10}%**",
                     text=f"metric number {i} description")
        for i in range(1, n + 1)])


def test_measure_equals_render():
    ctx = make_ctx()
    comp = IconTileRow()
    data = spec(4)
    m = comp.measure(data, inch(6), ctx)
    slide = blank_slide()
    c = comp.render(slide, data, BBox(0, 0, inch(6), inch(5)), ctx)
    assert abs(m - c) <= pt(1)


def test_measure_equals_render_three():
    ctx = make_ctx()
    comp = IconTileRow()
    data = spec(3)
    m = comp.measure(data, inch(5), ctx)
    slide = blank_slide()
    c = comp.render(slide, data, BBox(0, 0, inch(5), inch(5)), ctx)
    assert abs(m - c) <= pt(1)


def test_renders_tiles_and_bands():
    ctx = make_ctx()
    slide = blank_slide()
    IconTileRow().render(slide, spec(4), BBox(0, 0, inch(6), inch(5)), ctx)
    # 4 bands + 4 dark tiles = 8 autoshapes minimum, plus icon pictures
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert len(shapes) >= 8
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 4  # one icon per tile
