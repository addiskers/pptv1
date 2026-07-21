from pptx import Presentation

from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.kpi_card_strip import KpiCardStrip
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import (DonutStatSpec, KpiCardSpec,
                                          KpiCardStripSpec, ProgressPillSpec)


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def three_cards() -> KpiCardStripSpec:
    return KpiCardStripSpec(cards=[
        KpiCardSpec(title="**$2.2B**", body="leveraged in public and private "
                                            "co-investment since 2008"),
        KpiCardSpec(title="**60+**", body="innovations advanced to scale"),
        KpiCardSpec(title="**150+**", body="companies engaged across the value "
                                           "chain in focus geographies"),
    ])


def five_cards() -> KpiCardStripSpec:
    return KpiCardStripSpec(cards=[
        KpiCardSpec(title="**$2.2B**", body="capital leveraged"),
        KpiCardSpec(title="**60+**", body="innovations scaled"),
        KpiCardSpec(title="**150+**", body="companies engaged"),
        KpiCardSpec(title="**35**", body="countries with active programs and "
                                         "delivery partners on the ground"),
        KpiCardSpec(title="**9x**", body="return on catalytic funding"),
    ])


def test_measure_equals_render_three_cards():
    ctx = make_ctx()
    comp = KpiCardStrip()
    data = three_cards()
    width = inch(9)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(3)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_measure_equals_render_five_cards():
    ctx = make_ctx()
    comp = KpiCardStrip()
    data = five_cards()
    width = inch(11)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(3)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_all_cards_same_height_and_text_present():
    ctx = make_ctx()
    comp = KpiCardStrip()
    data = three_cards()
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, inch(9), inch(3)), ctx)

    rects = list(slide.shapes)
    assert len(rects) == 3
    heights = {int(s.height) for s in rects}
    assert len(heights) == 1
    assert heights.pop() == consumed
    assert consumed >= inch(0.9)

    text = " ".join(r.text for s in rects
                    for p in s.text_frame.paragraphs for r in p.runs)
    assert "$2.2B" in text
    assert "innovations" in text
    assert "150+" in text


# --- micro-viz combo tiles (reference FMD tile) -------------------------------

def viz_cards() -> KpiCardStripSpec:
    return KpiCardStripSpec(cards=[
        KpiCardSpec(title="**$2.2B**", body="capital leveraged"),
        KpiCardSpec(title="FMD doses administered",
                    viz=ProgressPillSpec(label="2025 → 2030", value_pct=52,
                                         display="312 M",
                                         target_display="600 M")),
        KpiCardSpec(title="Immunisation coverage",
                    viz=DonutStatSpec(value_pct=50, center_text="50%",
                                      label="of livestock covered")),
    ])


def test_measure_equals_render_with_viz_cards():
    ctx = make_ctx()
    comp = KpiCardStrip()
    data = viz_cards()
    width = inch(10)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(3)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_viz_card_nests_component_with_inverse_colors():
    ctx = make_ctx()
    comp = KpiCardStrip()
    data = viz_cards()
    _, slide = blank_slide()
    comp.render(slide, data, BBox(0, 0, inch(10), inch(3)), ctx)

    # the donut hole must match the card fill, not the slide bg — with cards
    # and hole both primary_dark, at least 4 shapes carry that fill
    # (3 cards + 1 hole)
    card_fill = ctx.theme.color("primary_dark")
    fills = []
    for s in slide.shapes:
        try:
            fills.append(str(s.fill.fore_color.rgb))
        except (TypeError, AttributeError):
            pass
    assert fills.count(card_fill) >= 4

    text = " ".join(r.text for s in slide.shapes if s.has_text_frame
                    for p in s.text_frame.paragraphs for r in p.runs)
    assert "FMD" in text
    assert "50%" in text
    assert "312 M" in text


def test_body_optional_card():
    ctx = make_ctx()
    comp = KpiCardStrip()
    data = KpiCardStripSpec(cards=[
        KpiCardSpec(title="**9x**"),
        KpiCardSpec(title="**$1B**", body="deployed"),
    ])
    width = inch(6)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(2)), ctx)
    assert abs(measured - consumed) <= pt(1)
