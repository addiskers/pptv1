"""native_chart variant renderers — reopen the built pptx and assert the
chart XML is valid and carries the requested variant markup."""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.native_chart import NativeChart
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import (ChartSeriesSpec, ChartStyleSpec,
                                          NativeChartSpec)


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_paper"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _chart(**kw) -> NativeChartSpec:
    kw.setdefault("sort", "none")
    kw.setdefault("highlight", None)
    kw.setdefault("annotation", None)
    return NativeChartSpec(**kw)


def _graphic_frames(slide):
    return [s for s in slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.CHART]


def _render(data, ctx=None, h=inch(3)):
    ctx = ctx or make_ctx()
    _, slide = blank_slide()
    consumed = NativeChart().render(slide, data, BBox(0, 0, inch(9), h), ctx)
    return slide, consumed, ctx


# --- parity holds for every variant -------------------------------------------

def test_measure_render_parity_variants():
    ctx = make_ctx()
    specs = [
        _chart(chart_type="bar", categories=["A", "B", "C"],
               series=[ChartSeriesSpec(name="s", values=[3, 2, 1])],
               sort="desc", value_suffix="%"),
        _chart(chart_type="line", categories=["19", "20", "21", "22"],
               series=[ChartSeriesSpec(name="s", values=[1, 2, 3, 4])],
               style=ChartStyleSpec(endpoint_labels=True, cagr_chip=True)),
        _chart(chart_type="waterfall",
               categories=["Open", "Up", "Down", "Close"],
               series=[ChartSeriesSpec(name="s", values=[100, 20, -20, 100])]),
    ]
    for data in specs:
        measured = NativeChart().measure(data, inch(9), ctx)
        _, consumed, _ = _render(data)
        assert abs(measured - consumed) <= pt(1), data.chart_type


# --- value labels: auto vs explicit -------------------------------------------

def _has_data_labels(frame) -> bool:
    xml = frame.chart._chartSpace.xml
    return "dLbls" in xml


def test_value_labels_auto_single_series_bar():
    slide, _, _ = _render(_chart(
        chart_type="bar", categories=["A", "B"],
        series=[ChartSeriesSpec(name="s", values=[2, 1])], sort="desc"))
    assert _has_data_labels(_graphic_frames(slide)[0])  # historic default


def test_value_labels_off_when_false():
    slide, _, _ = _render(_chart(
        chart_type="bar", categories=["A", "B"],
        series=[ChartSeriesSpec(name="s", values=[2, 1])], sort="desc",
        style=ChartStyleSpec(value_labels=False)))
    assert not _has_data_labels(_graphic_frames(slide)[0])


def test_value_labels_on_line_when_true():
    slide, _, _ = _render(_chart(
        chart_type="line", categories=["A", "B", "C"],
        series=[ChartSeriesSpec(name="s", values=[1, 2, 3])],
        style=ChartStyleSpec(value_labels=True)))
    assert _has_data_labels(_graphic_frames(slide)[0])


# --- endpoint labels: exactly first + last idx --------------------------------

def test_endpoint_labels_only_first_and_last():
    slide, _, _ = _render(_chart(
        chart_type="line", categories=["19", "20", "21", "22"],
        series=[ChartSeriesSpec(name="s", values=[1, 2, 3, 4])],
        style=ChartStyleSpec(endpoint_labels=True)))
    xml = _graphic_frames(slide)[0].chart._chartSpace
    idxs = sorted(int(el.get("val"))
                  for dl in xml.iter(qn("c:dLbl"))
                  for el in dl.iter(qn("c:idx")))
    assert idxs == [0, 3]  # first + last of a 4-point series


# --- highlight_series mutes the rest ------------------------------------------

def test_highlight_series_mutes_others():
    ctx = make_ctx()
    slide, _, _ = _render(_chart(
        chart_type="line", categories=["A", "B", "C"],
        series=[ChartSeriesSpec(name="Acme", values=[1, 2, 3]),
                ChartSeriesSpec(name="Rival", values=[1, 1, 1])],
        style=ChartStyleSpec(highlight_series="Acme")), ctx)
    chart = _graphic_frames(slide)[0].chart
    accent = ctx.theme.color("accent")
    grid = ctx.theme.color("grid")
    colors = []
    for ser in chart.series:
        ln = ser._element.find(qn("c:spPr"))
        srgb = ln.iter(qn("a:srgbClr"))
        colors.append(next(srgb).get("val"))
    assert accent in colors and grid in colors


def test_unknown_highlight_series_warns():
    ctx = make_ctx()
    _render(_chart(
        chart_type="line", categories=["A", "B"],
        series=[ChartSeriesSpec(name="Acme", values=[1, 2]),
                ChartSeriesSpec(name="Rival", values=[1, 1])],
        style=ChartStyleSpec(highlight_series="Nope")), ctx)
    assert any("highlight_series" in w for w in ctx.report.warnings)


# --- forecast split -----------------------------------------------------------

def test_forecast_splits_into_two_series_with_dash():
    ctx = make_ctx()
    slide, _, _ = _render(_chart(
        chart_type="line", categories=["19", "20", "21", "22"],
        series=[ChartSeriesSpec(name="Rev", values=[1, 2, 3, 4])],
        style=ChartStyleSpec(forecast_from="21")), ctx)
    chart = _graphic_frames(slide)[0].chart
    assert len(chart.series) == 2  # solid + dashed halves
    xml = chart._chartSpace.xml
    assert "prstDash" in xml
    assert not chart.has_legend  # split would read as two series


def test_forecast_unknown_category_warns():
    ctx = make_ctx()
    _render(_chart(
        chart_type="line", categories=["19", "20"],
        series=[ChartSeriesSpec(name="Rev", values=[1, 2])],
        style=ChartStyleSpec(forecast_from="99")), ctx)
    assert any("forecast_from" in w for w in ctx.report.warnings)


# --- benchmark pins the axis + draws a dashed overlay -------------------------

def test_benchmark_pins_axis_and_draws_line():
    slide, _, _ = _render(_chart(
        chart_type="bar", categories=["A", "B", "C"],
        series=[ChartSeriesSpec(name="s", values=[90, 70, 50])], sort="desc",
        value_suffix="%",
        style=ChartStyleSpec(benchmark={"value": 74, "label": "avg"})))
    chart = _graphic_frames(slide)[0].chart
    va = chart.value_axis
    assert va.minimum_scale == 0
    assert va.maximum_scale is not None and va.maximum_scale >= 90
    # dashed overlay connector present on the slide
    dashed = [s for s in slide.shapes
              if s.shape_type == MSO_SHAPE_TYPE.LINE]
    assert len(dashed) == 1


# --- cagr chip: engine-computed value -----------------------------------------

def test_cagr_chip_text_is_engine_computed():
    slide, _, _ = _render(_chart(
        chart_type="line", categories=["19", "20", "21", "22", "23"],
        series=[ChartSeriesSpec(name="s", values=[10, 12, 14, 17, 21])],
        style=ChartStyleSpec(cagr_chip=True)))
    # 10 -> 21 over 4 periods = 20.4%
    text = " ".join(r.text for s in slide.shapes if s.has_text_frame
                    for p in s.text_frame.paragraphs for r in p.runs)
    assert "CAGR +20.4%" in text


def test_cagr_chip_skips_on_bad_data():
    ctx = make_ctx()
    _render(_chart(
        chart_type="line", categories=["A", "B"],
        series=[ChartSeriesSpec(name="s", values=[0, 5])],
        style=ChartStyleSpec(cagr_chip=True)), ctx)
    assert any("cagr" in w.lower() for w in ctx.report.warnings)


# --- waterfall: bases float, deltas colored by sign, arithmetic checked -------

def test_waterfall_two_series_and_balances():
    ctx = make_ctx()
    slide, _, _ = _render(_chart(
        chart_type="waterfall",
        categories=["Open", "Up", "Down", "Close"],
        series=[ChartSeriesSpec(name="s", values=[100, 30, -30, 100])]), ctx)
    chart = _graphic_frames(slide)[0].chart
    assert len(chart.series) == 2  # invisible base + visible
    assert not any("balance" in w for w in ctx.report.warnings)


def test_waterfall_imbalance_warns():
    ctx = make_ctx()
    _render(_chart(
        chart_type="waterfall",
        categories=["Open", "Up", "Close"],
        series=[ChartSeriesSpec(name="s", values=[100, 30, 200])]), ctx)
    assert any("balance" in w for w in ctx.report.warnings)


# --- diverging: single-series bar with negatives colors by sign ---------------

def test_diverging_colors_by_sign():
    ctx = make_ctx()
    slide, _, _ = _render(_chart(
        chart_type="bar", categories=["A", "B", "C"],
        series=[ChartSeriesSpec(name="s", values=[5, -3, 2])], sort="none"),
        ctx)
    chart = _graphic_frames(slide)[0].chart
    pos, neg = ctx.theme.color("positive"), ctx.theme.color("negative")
    pt_colors = [next(p.format._element.iter(qn("a:srgbClr"))).get("val")
                 for p in chart.series[0].points]
    assert pos in pt_colors and neg in pt_colors


# --- horizontal reverses category order + maps enum ---------------------------

def test_horizontal_reverses_categories():
    slide, _, _ = _render(_chart(
        chart_type="bar", categories=["First", "Second", "Third"],
        series=[ChartSeriesSpec(name="s", values=[3, 2, 1])], sort="none",
        style=ChartStyleSpec(direction="horizontal")))
    chart = _graphic_frames(slide)[0].chart
    cats = list(chart.plots[0].categories)
    assert cats == ["Third", "Second", "First"]  # reversed for top-first read


# --- compact drops the legend -------------------------------------------------

def test_compact_no_legend():
    slide, _, _ = _render(_chart(
        chart_type="line", categories=["A", "B", "C"],
        series=[ChartSeriesSpec(name="x", values=[1, 2, 3]),
                ChartSeriesSpec(name="y", values=[2, 2, 2])],
        style=ChartStyleSpec(compact=True)))
    assert not _graphic_frames(slide)[0].chart.has_legend


# --- percent_100 maps to the 100% stacked enum --------------------------------

def test_percent_100_enum():
    from pptx.enum.chart import XL_CHART_TYPE
    slide, _, _ = _render(_chart(
        chart_type="stacked_bar", categories=["A", "B"],
        series=[ChartSeriesSpec(name="x", values=[60, 55]),
                ChartSeriesSpec(name="y", values=[40, 45])],
        style=ChartStyleSpec(percent_100=True)))
    assert (_graphic_frames(slide)[0].chart.chart_type
            == XL_CHART_TYPE.COLUMN_STACKED_100)
