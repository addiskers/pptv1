from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.image_block import ImageBlock
from deckengine.core.assets import ASSETS_DIR
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import ImageBlockSpec

SAMPLE = "hero_field.png"  # 1200x800 under assets/samples/


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def pictures(slide):
    return [s for s in slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.PICTURE]


def all_text(slide) -> str:
    return " ".join(r.text for s in slide.shapes if s.has_text_frame
                    for p in s.text_frame.paragraphs for r in p.runs)


def test_measure_equals_render_without_caption():
    ctx = make_ctx()
    comp = ImageBlock()
    data = ImageBlockSpec(src=SAMPLE)
    width = inch(6)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(4)), ctx)
    assert abs(measured - consumed) <= pt(1)
    assert measured == inch(2.2)  # frame only, never the image file


def test_measure_equals_render_with_caption():
    ctx = make_ctx()
    comp = ImageBlock()
    data = ImageBlockSpec(src=SAMPLE, fit="cover", height_in=1.8,
                          caption="Field trial sites, *2025 season*")
    width = inch(8)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(inch(1), inch(1), width,
                                             inch(4)), ctx)
    assert abs(measured - consumed) <= pt(1)
    assert measured > inch(1.8)  # caption + gap added below the frame
    text = all_text(slide)
    assert "sites," in text
    assert "season" in text


def test_measure_independent_of_asset_presence():
    ctx = make_ctx()
    comp = ImageBlock()
    width = inch(7)
    present = comp.measure(ImageBlockSpec(src=SAMPLE), width, ctx)
    missing = comp.measure(ImageBlockSpec(src="no_such_file.png"),
                           width, ctx)
    assert present == missing


def test_missing_asset_renders_placeholder():
    ctx = make_ctx()
    comp = ImageBlock()
    data = ImageBlockSpec(src="definitely_missing.png",
                          caption="Where the photo would go")
    width = inch(6)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(4)), ctx)

    assert abs(measured - consumed) <= pt(1)
    assert pictures(slide) == []
    rects = [s for s in slide.shapes
             if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert len(rects) == 1
    assert rects[0].width == inch(6)
    assert rects[0].height == inch(2.2)
    assert "definitely_missing.png" in all_text(slide)
    assert any("image_block" in w and "not found" in w
               for w in ctx.report.warnings)


def test_placeholder_reports_truncated_long_src():
    ctx = make_ctx()
    comp = ImageBlock()
    data = ImageBlockSpec(src="missing_" + "x" * 200 + ".png")
    _, slide = blank_slide()
    comp.render(slide, data, BBox(0, 0, inch(3), inch(4)), ctx)
    assert any("image_block placeholder" in t for t in ctx.report.truncations)


def test_contain_places_picture_inside_frame():
    ctx = make_ctx()
    comp = ImageBlock()
    data = ImageBlockSpec(src=SAMPLE, fit="contain")
    bbox = BBox(inch(1), inch(1.5), inch(6), inch(4))
    _, slide = blank_slide()
    comp.render(slide, data, bbox, ctx)

    pics = pictures(slide)
    assert len(pics) == 1
    pic = pics[0]
    frame_bottom = bbox.y + inch(2.2)
    assert pic.left >= bbox.x
    assert pic.top >= bbox.y
    assert pic.left + pic.width <= bbox.right + 1
    assert pic.top + pic.height <= frame_bottom + 1
    # 1200x800 source: aspect preserved, height-limited in a 6x2.2in frame
    assert abs(pic.width / pic.height - 1.5) < 0.01
    assert abs(pic.height - inch(2.2)) <= 2
    # centered horizontally
    assert abs((pic.left - bbox.x) - (bbox.right - pic.left - pic.width)) <= 2


def test_cover_fills_frame_exactly_and_caches_crop():
    ctx = make_ctx()
    comp = ImageBlock()
    data = ImageBlockSpec(src=SAMPLE, fit="cover")
    bbox = BBox(inch(0.5), inch(2), inch(6), inch(4))
    _, slide = blank_slide()
    comp.render(slide, data, bbox, ctx)

    pics = pictures(slide)
    assert len(pics) == 1
    pic = pics[0]
    assert pic.left == bbox.x
    assert pic.top == bbox.y
    assert pic.width == inch(6)
    assert pic.height == inch(2.2)
    crops = list((ASSETS_DIR / "cache_crop").glob("hero_field_*.png"))
    assert crops, "cover crop not cached under assets/cache_crop"
