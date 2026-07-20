from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.two_tone_header import TwoToneHeader
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import Span, TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import TwoToneHeaderSpec


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def rects(slide):
    return [s for s in slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]


LONG_RIGHT = ("Adoption of drought-tolerant varieties across all focus "
              "districts, with input credit bundled at the point of sale")


def test_measure_equals_render_single_line():
    ctx = make_ctx()
    comp = TwoToneHeader()
    data = TwoToneHeaderSpec(left="Objective", right="Raise smallholder yields")
    width = inch(6)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data,
                           BBox(inch(1), inch(1), width, inch(1)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_measure_equals_render_wrapping_right():
    ctx = make_ctx()
    comp = TwoToneHeader()
    data = TwoToneHeaderSpec(left="Target", right=LONG_RIGHT, left_frac=0.3)
    width = inch(4)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data,
                           BBox(inch(0.5), inch(0.5), width, inch(1.5)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_band_geometry_and_fills():
    ctx = make_ctx()
    _, slide = blank_slide()
    width = inch(6)
    data = TwoToneHeaderSpec(left="Phase 1", right="Foundation", left_frac=0.25)
    origin = BBox(inch(1), inch(1), width, inch(1))
    consumed = TwoToneHeader().render(slide, data, origin, ctx)
    shapes = rects(slide)
    assert len(shapes) == 2
    left = min(shapes, key=lambda s: int(s.left))
    right = max(shapes, key=lambda s: int(s.left))
    left_w = round(width * 0.25)
    assert int(left.width) == left_w
    assert int(left.left) == origin.x
    assert int(right.left) == origin.x + left_w
    assert int(right.width) == width - left_w
    assert abs(int(left.height) - consumed) <= pt(1)
    assert abs(int(right.height) - consumed) <= pt(1)
    assert left.fill.fore_color.rgb == RGBColor.from_string(
        ctx.theme.color("primary"))
    assert right.fill.fore_color.rgb == RGBColor.from_string(
        ctx.theme.color("surface_alt"))
    # single-line band: one small line plus spacing(0.5)
    one_line = max(
        ctx.measurer.line_height_emu([Span("Ag", bold=True)],
                                     ctx.font("body"), ctx.size("small")),
        ctx.measurer.line_height_emu([Span("Ag")],
                                     ctx.font("body"), ctx.size("small")))
    assert abs(consumed - (one_line + ctx.theme.spacing(0.5))) <= pt(1)


def test_text_styling_and_right_inset():
    ctx = make_ctx()
    _, slide = blank_slide()
    width = inch(6)
    data = TwoToneHeaderSpec(left="Workstream", right="Market linkages",
                             left_frac=0.35)
    origin = BBox(0, 0, width, inch(1))
    TwoToneHeader().render(slide, data, origin, ctx)
    left = min(rects(slide), key=lambda s: int(s.left))
    left_runs = [r for p in left.text_frame.paragraphs for r in p.runs]
    assert "".join(r.text for r in left_runs) == "Workstream"
    assert all(r.font.bold for r in left_runs)
    assert all(str(r.font.color.rgb) == ctx.theme.color("inverse_ink")
               for r in left_runs)
    boxes = [s for s in slide.shapes
             if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX]
    assert len(boxes) == 1
    left_w = round(width * 0.35)
    assert int(boxes[0].left) == origin.x + left_w + ctx.theme.spacing(0.5)
    right_runs = [r for p in boxes[0].text_frame.paragraphs for r in p.runs]
    assert "".join(r.text for r in right_runs) == "Market linkages"
    assert all(str(r.font.color.rgb) == ctx.theme.color("ink")
               for r in right_runs)


def test_wrapping_side_grows_band_to_two_lines():
    ctx = make_ctx()
    comp = TwoToneHeader()
    width = inch(4)
    short = comp.measure(
        TwoToneHeaderSpec(left="Target", right="Yields", left_frac=0.3),
        width, ctx)
    tall = comp.measure(
        TwoToneHeaderSpec(left="Target", right=LONG_RIGHT, left_frac=0.3),
        width, ctx)
    assert tall > short
    # never more than two lines: band capped at two small lines + padding
    two_lines = 2 * ctx.measurer.line_height_emu(
        [Span("Ag", bold=True)], ctx.font("body"), ctx.size("small"))
    assert tall <= two_lines + ctx.theme.spacing(0.5) + pt(1)
