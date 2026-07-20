from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.progress_pill import ProgressPill
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import Span, TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import ProgressPillSpec


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _rounded(slide):
    return [s for s in slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and s.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE]


def _runs(slide):
    return [r for s in slide.shapes if s.has_text_frame
            for p in s.text_frame.paragraphs for r in p.runs]


def test_measure_equals_render_with_target():
    ctx = make_ctx()
    comp = ProgressPill()
    data = ProgressPillSpec(label="**FMD vaccination** coverage",
                            value_pct=62, display="312M",
                            target_display="500M")
    width = inch(6)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(1)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_measure_equals_render_no_target_narrow():
    ctx = make_ctx()
    comp = ProgressPill()
    data = ProgressPillSpec(label="Cold-chain capacity upgraded in districts",
                            value_pct=3, display="9M")
    width = inch(3)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(1.5)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_track_and_segment_geometry():
    ctx = make_ctx()
    _, slide = blank_slide()
    width = inch(6)
    ProgressPill().render(slide, ProgressPillSpec(label="Coverage",
                                                  value_pct=62,
                                                  display="312M"),
                          BBox(0, 0, width, inch(1)), ctx)
    track, seg = _rounded(slide)
    assert str(track.fill.fore_color.rgb) == ctx.theme.color("surface_alt")
    assert str(seg.fill.fore_color.rgb) == ctx.theme.color("primary")
    # track spans the full width; expected track height = micro line + s(0.3)
    line_h = ctx.measurer.line_height_emu([Span("Ag", bold=True)],
                                          ctx.font("body"), ctx.size("micro"))
    assert track.width == width
    assert track.height == seg.height == line_h + ctx.theme.spacing(0.3)
    # segment = value_pct% of the track, aligned to the same band
    assert seg.width == round(width * 0.62)
    assert seg.top == track.top and seg.left == track.left


def test_display_inside_segment_bold_inverse():
    ctx = make_ctx()
    _, slide = blank_slide()
    ProgressPill().render(slide, ProgressPillSpec(label="Coverage",
                                                  value_pct=62,
                                                  display="312M"),
                          BBox(0, 0, inch(6), inch(1)), ctx)
    disp = [r for r in _runs(slide) if r.text == "312M"]
    assert len(disp) == 1
    assert disp[0].font.bold
    assert str(disp[0].font.color.rgb) == ctx.theme.color("inverse_ink")


def test_tiny_value_min_segment_and_display_outside_in_ink():
    ctx = make_ctx()
    _, slide = blank_slide()
    ProgressPill().render(slide, ProgressPillSpec(label="Coverage",
                                                  value_pct=1,
                                                  display="9M"),
                          BBox(0, 0, inch(3), inch(1)), ctx)
    _, seg = _rounded(slide)
    # 1% of 3in is thinner than the pill height -> clamped to a circle-ish pill
    assert seg.width == seg.height
    disp = [r for r in _runs(slide) if r.text == "9M"]
    assert len(disp) == 1
    assert str(disp[0].font.color.rgb) == ctx.theme.color("ink")


def test_target_display_right_aligned_muted():
    ctx = make_ctx()
    _, slide = blank_slide()
    ProgressPill().render(slide, ProgressPillSpec(label="Coverage",
                                                  value_pct=62,
                                                  display="312M",
                                                  target_display="500M"),
                          BBox(0, 0, inch(6), inch(1)), ctx)
    tgt_runs = [(s, p, r) for s in slide.shapes if s.has_text_frame
                for p in s.text_frame.paragraphs for r in p.runs
                if r.text == "500M"]
    assert len(tgt_runs) == 1
    shape, para, run = tgt_runs[0]
    assert str(run.font.color.rgb) == ctx.theme.color("ink_muted")
    assert para.alignment == PP_ALIGN.RIGHT
    # the target box ends at the track end (full component width)
    assert shape.left + shape.width == inch(6)
