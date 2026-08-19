"""Wave-3 diagram forms: cycle, tree (issue/driver/org), onion."""
from __future__ import annotations

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from deckengine.components.base import BuildReport, RenderContext, \
    get_component
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.rings import get_cycle_ring
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.llm.canon import CANON
from deckengine.llm.format_rules import first_rule, signals_from
from deckengine.schema.components import (CycleSpec, OnionSpec, TreeNodeSpec,
                                          TreeSpec)


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _cycle():
    return CycleSpec(stages=["Sellers", "Selection", "Buyers", "Prices"],
                     hub="Flywheel", highlight_index=2)


def _tree(variant="issue", **kw):
    return TreeSpec(root="Why is the gap widening?", variant=variant,
                    children=[
                        TreeNodeSpec(label="Inputs", children=["Feedstock"]),
                        TreeNodeSpec(label="Utilization",
                                     children=["Line 2", "Changeovers"]),
                        TreeNodeSpec(label="Logistics")], **kw)


def _onion():
    return OnionSpec(layers=["Core", "Members", "Dealers", "Retail"],
                     highlight_index=1)


def test_measure_render_parity():
    ctx = make_ctx()
    for kind, spec in (("cycle", _cycle()), ("tree", _tree()),
                       ("tree", _tree("driver")), ("tree", _tree("org")),
                       ("onion", _onion())):
        comp = get_component(kind)
        _, slide = blank_slide()
        m = comp.measure(spec, inch(9), ctx)
        c = comp.render(slide, spec, BBox(0, 0, inch(9), inch(6)), ctx)
        assert abs(m - c) <= pt(1), f"{kind}/{getattr(spec, 'variant', '')}"


def test_cycle_ring_cached_and_labels_present():
    ctx = make_ctx()
    _, slide = blank_slide()
    get_component("cycle").render(slide, _cycle(),
                                  BBox(0, 0, inch(9), inch(5)), ctx)
    pics = [s for s in slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1
    texts = " ".join(r.text for s in slide.shapes if s.has_text_frame
                     for p in s.text_frame.paragraphs for r in p.runs)
    for stage in ("Sellers", "Selection", "Buyers", "Prices", "Flywheel"):
        assert stage in texts
    p = get_cycle_ring(4, "2C3A47", "C05621", 2)
    assert p.is_file()
    assert get_cycle_ring(4, "2C3A47", "C05621", 2) == p  # cache hit


def test_driver_tree_has_operator_chips_issue_does_not():
    ctx = make_ctx()
    _, s_issue = blank_slide()
    get_component("tree").render(s_issue, _tree("issue"),
                                 BBox(0, 0, inch(9), inch(4)), ctx)
    _, s_driver = blank_slide()
    get_component("tree").render(s_driver, _tree("driver", operator="x"),
                                 BBox(0, 0, inch(9), inch(4)), ctx)

    def chip_texts(slide):
        return [r.text for s in slide.shapes if s.has_text_frame
                for p in s.text_frame.paragraphs for r in p.runs
                if r.text in ("+", "×")]
    assert chip_texts(s_issue) == []
    assert chip_texts(s_driver) == ["×"] * 3


def test_tree_highlight_and_range_guard():
    ctx = make_ctx()
    _, slide = blank_slide()
    get_component("tree").render(slide, _tree(highlight_index=1),
                                 BBox(0, 0, inch(9), inch(4)), ctx)
    fills = [str(s.fill.fore_color.rgb) for s in slide.shapes
             if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             and s.fill.type is not None]
    assert ctx.theme.color("accent") in fills
    bad = _tree(highlight_index=3)
    _, slide2 = blank_slide()
    get_component("tree").render(slide2, bad,
                                 BBox(0, 0, inch(9), inch(4)), ctx)
    assert any("out of range" in w for w in ctx.report.warnings)


def test_onion_core_inside_rail_outside():
    ctx = make_ctx()
    _, slide = blank_slide()
    get_component("onion").render(slide, _onion(),
                                  BBox(0, 0, inch(9), inch(4)), ctx)
    ovals = [s for s in slide.shapes
             if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert len(ovals) == 4
    # dashed leader lines: one per non-core layer
    lines = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.LINE]
    assert len(lines) == 3
    texts = " ".join(r.text for s in slide.shapes if s.has_text_frame
                     for p in s.text_frame.paragraphs for r in p.runs)
    for label in ("Core", "Members", "Dealers", "Retail"):
        assert label in texts


def _temple():
    from deckengine.schema.components import TempleSpec
    return TempleSpec(goal="Category leadership",
                      pillars=["Agronomy", "Linkages", "Financing"],
                      foundation="Execution discipline", highlight_index=1)


def _iceberg():
    from deckengine.schema.components import IcebergSpec
    return IcebergSpec(visible=["Falling prices"],
                       hidden=["Fragmented aggregation", "No grading",
                               "Distress sales"])


def test_temple_iceberg_parity_and_content():
    ctx = make_ctx()
    for kind, spec in (("temple", _temple()), ("iceberg", _iceberg())):
        comp = get_component(kind)
        _, slide = blank_slide()
        m = comp.measure(spec, inch(9), ctx)
        c = comp.render(slide, spec, BBox(0, 0, inch(9), inch(6)), ctx)
        assert abs(m - c) <= pt(1), kind
        import re
        texts = re.sub(r"\s+", " ", " ".join(
            r.text for s in slide.shapes if s.has_text_frame
            for p in s.text_frame.paragraphs for r in p.runs))
        if kind == "temple":
            for t in ("Category leadership", "Agronomy", "Linkages",
                      "Financing", "Execution discipline"):
                assert t in texts
        else:
            for t in ("Falling prices", "Fragmented aggregation",
                      "Distress sales"):
                assert t in texts


def test_temple_highlighted_pillar_accented():
    ctx = make_ctx()
    _, slide = blank_slide()
    get_component("temple").render(slide, _temple(),
                                   BBox(0, 0, inch(9), inch(4)), ctx)
    fills = [str(s.fill.fore_color.rgb) for s in slide.shapes
             if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             and s.fill.type is not None]
    assert ctx.theme.color("accent") in fills


def test_iceberg_mass_rotated_and_waterline_dashed():
    ctx = make_ctx()
    _, slide = blank_slide()
    get_component("iceberg").render(slide, _iceberg(),
                                    BBox(0, 0, inch(9), inch(4)), ctx)
    shapes = [s for s in slide.shapes
              if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert any(abs(s.rotation - 180) < 0.1 for s in shapes)  # sunken mass
    lines = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.LINE]
    assert len(lines) == 1                                   # waterline


# -- routing + canon ---------------------------------------------------------

def test_wave3_routing():
    assert first_rule(signals_from(
        "The flywheel compounds: more sellers bring more buyers")).id \
        == "reinforcing_cycle"
    assert first_rule(signals_from(
        "The cost gap breaks down into three MECE drivers")).id \
        == "decomposition_tree"
    assert first_rule(signals_from(
        "Adoption spreads in concentric layers from the core outward")).id \
        == "layers_onion"
    # 'value chain' must NOT hijack into tree/cycle
    r = first_rule(signals_from("Our value chain spans four activities"))
    assert r is None or r.id not in ("reinforcing_cycle",
                                     "decomposition_tree")


def test_wave3b_routing():
    assert first_rule(signals_from(
        "The strategy rests on three capability pillars")).id \
        == "pillars_temple"
    assert first_rule(signals_from(
        "Falling prices are the symptom; the drivers sit beneath the "
        "surface")).id == "hidden_iceberg"


def test_wave3_canon_flips():
    for form, engine in (("flywheel", "cycle"),
                         ("issue_tree", "tree"),
                         ("driver_tree", "tree"),
                         ("onion", "onion"),
                         ("temple", "temple"),
                         ("iceberg", "iceberg")):
        assert CANON[form].status == "primitive", form
        assert engine in (CANON[form].engine or ""), form
