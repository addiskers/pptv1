from pptx import Presentation
from pptx.dml.color import RGBColor

from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.bullet_list import BulletList
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import BulletItemSpec, BulletListSpec


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def shape_text(shape) -> str:
    return "".join(r.text for p in shape.text_frame.paragraphs for r in p.runs)


def test_measure_equals_render_flat_body():
    ctx = make_ctx()
    comp = BulletList()
    data = BulletListSpec(items=[
        BulletItemSpec(text="**Yields doubled** across 14 program countries"),
        BulletItemSpec(text="Smallholder adoption reached 62% by 2024, driven "
                            "by extension networks and input subsidy reform"),
        BulletItemSpec(text="Unit delivery costs fell 18% year over year"),
    ])
    width = inch(4)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(inch(1), inch(1), width, inch(4)),
                           ctx)
    assert abs(measured - consumed) <= pt(1)


def test_measure_equals_render_nested_small():
    ctx = make_ctx()
    comp = BulletList()
    data = BulletListSpec(size_role="small", items=[
        BulletItemSpec(text="**Portfolio:** three crop platforms", level=0),
        BulletItemSpec(text="Maize breeding pipeline spans 9 national "
                            "programs with shared trial protocols", level=1),
        BulletItemSpec(text="Stage-gate reviews twice a year", level=2),
    ])
    width = inch(3)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data,
                           BBox(inch(0.5), inch(0.5), width, inch(4)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_glyph_column_indent_and_color():
    ctx = make_ctx()
    _, slide = blank_slide()
    data = BulletListSpec(items=[
        BulletItemSpec(text="Top level point", level=0),
        BulletItemSpec(text="Nested supporting detail", level=1),
    ])
    BulletList().render(slide, data, BBox(0, 0, inch(4), inch(2)), ctx)
    shapes = list(slide.shapes)
    assert len(shapes) == 4  # one glyph box + one text box per item
    glyphs = [s for s in shapes if shape_text(s) == "•"]
    assert len(glyphs) == 2
    # level-1 glyph sits exactly one indent step (spacing(2)) further right
    lefts = sorted(int(s.left) for s in glyphs)
    assert lefts[1] - lefts[0] == ctx.theme.spacing(2)
    # bullet dot colored with the primary role
    run = glyphs[0].text_frame.paragraphs[0].runs[0]
    assert run.font.color.rgb == RGBColor.from_string(
        ctx.theme.color("primary"))


def test_bold_lead_in_preserved():
    ctx = make_ctx()
    _, slide = blank_slide()
    data = BulletListSpec(items=[
        BulletItemSpec(text="**Lead-in:** rest of the point in regular weight"),
    ])
    BulletList().render(slide, data, BBox(0, 0, inch(4), inch(1)), ctx)
    runs = [r for s in slide.shapes for p in s.text_frame.paragraphs
            for r in p.runs]
    bold_runs = [r for r in runs if r.font.bold and r.text != "•"]
    assert bold_runs and bold_runs[0].text.startswith("Lead-in:")
