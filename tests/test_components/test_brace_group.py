from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE

from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.brace_group import BraceGroup
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import BraceGroupSpec, MiniTableSpec


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def reference_spec() -> BraceGroupSpec:
    return BraceGroupSpec(
        content=[MiniTableSpec(
            headers=["State", "Direct", "Digital"],
            rows=[["Bihar", "7.0 Mn", "18.5 Mn"],
                  ["Odisha", "2.0 Mn", "7.6 Mn"],
                  ["Total", "10.6 Mn", "37.7 Mn"]])],
        takeaway="**10.6 Mn** SSPs benefitted directly")


def test_measure_equals_render():
    ctx = make_ctx()
    comp = BraceGroup()
    data = reference_spec()
    width = inch(9)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(4)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_brace_is_line_only_and_spans_content():
    ctx = make_ctx()
    comp = BraceGroup()
    data = reference_spec()
    bbox = BBox(inch(0.5), inch(1.5), inch(9), inch(4))
    _, slide = blank_slide()
    comp.render(slide, data, bbox, ctx)

    braces = [s for s in slide.shapes
              if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
              and s.auto_shape_type == MSO_SHAPE.RIGHT_BRACE]
    assert len(braces) == 1
    brace = braces[0]
    assert brace.width == inch(0.22)
    assert brace.fill.type == MSO_FILL.BACKGROUND  # line-only, no fill
    # brace sits between content (left) and takeaway (right)
    takeaway_w = round(bbox.w * data.takeaway_frac)
    assert brace.left + brace.width < bbox.x + bbox.w - takeaway_w

    # runs concatenate within a paragraph (wrap tokens carry their spaces)
    text = " ".join("".join(r.text for r in p.runs)
                    for s in slide.shapes if s.has_text_frame
                    for p in s.text_frame.paragraphs)
    assert "10.6 Mn" in text
    assert "Bihar" in text


def test_takeaway_height_can_dominate():
    ctx = make_ctx()
    comp = BraceGroup()
    from deckengine.schema.components import TextBlockSpec
    data = BraceGroupSpec(
        content=[TextBlockSpec(text="short line")],
        takeaway="A very long takeaway sentence that wraps across many lines "
                 "in its narrow column and therefore sets the group height "
                 "instead of the single short content line on the left side",
        takeaway_frac=0.2)
    width = inch(8)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(4)), ctx)
    assert abs(measured - consumed) <= pt(1)
