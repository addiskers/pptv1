from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.donut_stat import DonutStat
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import DonutStatSpec


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _autoshapes(slide, mso_shape):
    return [s for s in slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and s.auto_shape_type == mso_shape]


def test_measure_equals_render_wide():
    ctx = make_ctx()
    comp = DonutStat()
    data = DonutStatSpec(value_pct=50, center_text="50%",
                         label="Households with **improved** water access")
    width = inch(3)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(2)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_measure_equals_render_narrow_label_below():
    ctx = make_ctx()
    comp = DonutStat()
    data = DonutStatSpec(value_pct=72, center_text="72%",
                         label="Coverage across all program districts")
    width = inch(1)  # too narrow for a side label -> label goes below
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(2.5)), ctx)
    assert abs(measured - consumed) <= pt(1)
    # ring (0.4in) plus a stacked label is taller than the ring alone
    assert measured > inch(0.4)


def test_ring_geometry_and_stacked_fills():
    ctx = make_ctx()
    _, slide = blank_slide()
    DonutStat().render(slide, DonutStatSpec(value_pct=40, center_text="40%",
                                            label="On-time delivery"),
                       BBox(0, 0, inch(3), inch(2)), ctx)
    ovals = _autoshapes(slide, MSO_SHAPE.OVAL)
    pies = _autoshapes(slide, MSO_SHAPE.PIE)
    assert len(ovals) == 2 and len(pies) == 1
    track, hole = ovals
    # outer diameter = min(bbox.w * 0.4, 1.1in) -> capped at 1.1in here
    assert track.width == track.height == inch(1.1)
    assert hole.width == hole.height == round(inch(1.1) * 0.55)
    assert str(track.fill.fore_color.rgb) == ctx.theme.color("surface_alt")
    assert str(pies[0].fill.fore_color.rgb) == ctx.theme.color("primary")
    assert str(hole.fill.fore_color.rgb) == ctx.theme.color("bg")
    # hole is concentric with the track
    assert hole.left - track.left == (track.width - hole.width) // 2
    assert hole.top - track.top == (track.height - hole.height) // 2


def test_pie_adjustment_xml_angles():
    """Empirical contract: OOXML pie angles are 60000ths of a degree AND the
    preset pins adjustments to [0, 21599999] — negative values are clamped
    to 0 by PowerPoint (full-circle bug). The component must write degrees
    normalized into [0, 360): start = 270, end = (270 + pct * 3.6) mod 360.
    """
    ctx = make_ctx()
    _, slide = blank_slide()
    DonutStat().render(slide, DonutStatSpec(value_pct=40, center_text="40%",
                                            label="On-time delivery"),
                       BBox(0, 0, inch(3), inch(2)), ctx)
    pie = _autoshapes(slide, MSO_SHAPE.PIE)[0]
    av_lst = pie._element.spPr.find(qn("a:prstGeom")).find(qn("a:avLst"))
    raw = {gd.get("name"): int(gd.get("fmla").split()[1])
           for gd in av_lst.findall(qn("a:gd"))}
    assert raw["adj1"] == round(270 * 60000)                      # 16200000
    assert raw["adj2"] == round(((270 + 40 * 3.6) % 360) * 60000)  # 3240000
    assert 0 <= raw["adj1"] <= 21599999
    assert 0 <= raw["adj2"] <= 21599999


def test_full_value_renders_solid_primary_oval_not_pie():
    ctx = make_ctx()
    _, slide = blank_slide()
    DonutStat().render(slide, DonutStatSpec(value_pct=100, center_text="100%",
                                            label="Complete"),
                       BBox(0, 0, inch(3), inch(2)), ctx)
    ovals = _autoshapes(slide, MSO_SHAPE.OVAL)
    assert not _autoshapes(slide, MSO_SHAPE.PIE)
    assert len(ovals) == 3  # track + full primary + hole
    assert str(ovals[1].fill.fore_color.rgb) == ctx.theme.color("primary")


def test_zero_value_draws_no_wedge():
    ctx = make_ctx()
    _, slide = blank_slide()
    DonutStat().render(slide, DonutStatSpec(value_pct=0, center_text="0%",
                                            label="Not started"),
                       BBox(0, 0, inch(3), inch(2)), ctx)
    assert len(_autoshapes(slide, MSO_SHAPE.OVAL)) == 2  # track + hole only
    assert not _autoshapes(slide, MSO_SHAPE.PIE)


def test_center_text_bold_ink_and_label_present():
    ctx = make_ctx()
    _, slide = blank_slide()
    DonutStat().render(slide, DonutStatSpec(value_pct=50, center_text="50%",
                                            label="Vaccination coverage"),
                       BBox(0, 0, inch(3), inch(2)), ctx)
    runs = [r for s in slide.shapes if s.has_text_frame
            for p in s.text_frame.paragraphs for r in p.runs]
    center = [r for r in runs if r.text == "50%"]
    assert len(center) == 1
    assert center[0].font.bold
    assert str(center[0].font.color.rgb) == ctx.theme.color("ink")
    assert any("Vaccination" in r.text for r in runs)


# --- pie angle normalization (regression: PIE pins adjustments to [0, 360)) ---

def test_pie_angles_normalized_into_pin_range():
    """25% must be a top-right quarter: start 270deg, end 0deg — never a
    negative adjustment (the preset pins those to 0 => full-circle bug)."""
    ctx = make_ctx()
    comp = DonutStat()
    data = DonutStatSpec(value_pct=25, center_text="25%", label="quarter")
    _, slide = blank_slide()
    comp.render(slide, data, BBox(0, 0, inch(3), inch(2)), ctx)
    pies = _autoshapes(slide, MSO_SHAPE.PIE)
    assert len(pies) == 1
    start, end = pies[0].adjustments[0], pies[0].adjustments[1]
    assert start >= 0 and end >= 0
    assert abs(start - 270 * 0.6) < 0.01   # 12 o'clock, normalized
    assert abs(end - 0.0) < 0.01           # 3 o'clock


def test_pie_half_sweeps_right_half():
    ctx = make_ctx()
    comp = DonutStat()
    data = DonutStatSpec(value_pct=50, center_text="50%", label="half")
    _, slide = blank_slide()
    comp.render(slide, data, BBox(0, 0, inch(3), inch(2)), ctx)
    pies = _autoshapes(slide, MSO_SHAPE.PIE)
    start, end = pies[0].adjustments[0], pies[0].adjustments[1]
    assert abs(start - 162.0) < 0.01       # 270 deg
    assert abs(end - 54.0) < 0.01          # 90 deg (6 o'clock)
