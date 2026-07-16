from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.mini_table import MiniTable
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import MiniTableSpec


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def test_measure_equals_render_equal_columns():
    ctx = make_ctx()
    comp = MiniTable()
    data = MiniTableSpec(headers=["Crop", "Yield", "Area"],
                         rows=[["Maize", "2.1", "35%"],
                               ["Cassava", "8.6", "22%"]])
    width = inch(3)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(3)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_measure_equals_render_col_fracs():
    ctx = make_ctx()
    comp = MiniTable()
    data = MiniTableSpec(headers=["Country", "2008", "2015", "CAGR"],
                         rows=[["Ethiopia", "12", "19", "6.8%"],
                               ["Tanzania", "9", "14", "6.5%"],
                               ["Nigeria", "31", "40", "3.7%"],
                               ["Ghana", "5", "8", "6.9%"]],
                         col_fracs=[2.0, 1.0, 1.0, 1.0], align="left")
    width = inch(4.5)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(4)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_shape_grid_not_native_table():
    ctx = make_ctx()
    _, slide = blank_slide()
    data = MiniTableSpec(headers=["A", "B", "C"],
                         rows=[["1", "2", "3"], ["4", "5", "6"]])
    MiniTable().render(slide, data, BBox(0, 0, inch(3), inch(2)), ctx)
    # one rect per cell: (1 header + 2 rows) x 3 cols
    autoshapes = [s for s in slide.shapes
                  if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert len(autoshapes) == 9
    assert not any(s.has_table for s in slide.shapes)


def test_header_bold_body_regular_and_text_placed():
    ctx = make_ctx()
    _, slide = blank_slide()
    data = MiniTableSpec(headers=["Crop", "Yield"],
                         rows=[["Maize", "2.1"]])
    MiniTable().render(slide, data, BBox(0, 0, inch(2.5), inch(1)), ctx)
    shapes = [s for s in slide.shapes
              if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    texts = {}
    for s in shapes:
        runs = [r for p in s.text_frame.paragraphs for r in p.runs]
        assert len(runs) == 1
        texts[runs[0].text] = runs[0].font.bold
    assert texts == {"Crop": True, "Yield": True, "Maize": False, "2.1": False}


def test_rows_truncated_when_bbox_too_short():
    ctx = make_ctx()
    comp = MiniTable()
    data = MiniTableSpec(headers=["A", "B"],
                         rows=[[str(i), str(i)] for i in range(8)])
    _, slide = blank_slide()
    short = BBox(0, 0, inch(2), pt(30))  # far less than 9 rows need
    consumed = comp.render(slide, data, short, ctx)
    assert ctx.report.truncations
    assert consumed <= short.h + pt(1)
    assert consumed < comp.measure(data, short.w, ctx)
