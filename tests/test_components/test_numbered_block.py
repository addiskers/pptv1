from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.numbered_block import NumberedBlock
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import Span, TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import NumberedBlockSpec


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def autoshapes(slide):
    return [s for s in slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]


def text_boxes(slide):
    return [s for s in slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX]


def expected_diameter(ctx: RenderContext) -> int:
    two_lines = 2 * ctx.measurer.line_height_emu(
        [Span("Ag", bold=True)], ctx.font("body"), ctx.size("small"))
    return min(inch(0.42), two_lines)


def test_measure_equals_render_title_only():
    ctx = make_ctx()
    comp = NumberedBlock()
    data = NumberedBlockSpec(number="1", title="Strengthen delivery channels")
    width = inch(4)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data,
                           BBox(inch(1), inch(1), width, inch(2)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_measure_equals_render_title_and_body():
    ctx = make_ctx()
    comp = NumberedBlock()
    data = NumberedBlockSpec(
        number="2B",
        title="**Scale** smallholder access",
        body="Extend agro-dealer networks into underserved districts and "
             "bundle inputs with credit so adoption survives the first "
             "season of poor rainfall.")
    width = inch(3.5)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data,
                           BBox(inch(0.5), inch(0.5), width, inch(3)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_circle_geometry_fill_and_number_styling():
    ctx = make_ctx()
    _, slide = blank_slide()
    data = NumberedBlockSpec(number="1A", title="Pillar one")
    origin = BBox(inch(1), inch(1), inch(4), inch(2))
    NumberedBlock().render(slide, data, origin, ctx)
    shapes = autoshapes(slide)
    assert len(shapes) == 1
    circle = shapes[0]
    d = expected_diameter(ctx)
    assert int(circle.width) == d and int(circle.height) == d
    assert int(circle.left) == origin.x
    assert circle.fill.fore_color.rgb == RGBColor.from_string(
        ctx.theme.color("accent"))
    runs = [r for p in circle.text_frame.paragraphs for r in p.runs]
    assert [r.text for r in runs] == ["1A"]
    assert runs[0].font.bold
    assert str(runs[0].font.color.rgb) == ctx.theme.color("inverse_ink")


def test_text_column_position_colors_and_centering():
    ctx = make_ctx()
    _, slide = blank_slide()
    data = NumberedBlockSpec(
        number="3",
        title="Unlock financing",
        body="Blended capital reduces first-loss exposure for local banks "
             "and unlocks working capital for input suppliers across the "
             "value chain in every focus geography.")
    origin = BBox(inch(1), inch(1), inch(3), inch(3))
    consumed = NumberedBlock().render(slide, data, origin, ctx)
    d = expected_diameter(ctx)
    gap = ctx.theme.spacing(0.5)

    boxes = text_boxes(slide)
    assert len(boxes) == 2  # title + body
    assert all(int(b.left) == origin.x + d + gap for b in boxes)
    title_box = min(boxes, key=lambda b: int(b.top))
    body_box = max(boxes, key=lambda b: int(b.top))
    title_runs = [r for p in title_box.text_frame.paragraphs for r in p.runs]
    assert str(title_runs[0].font.color.rgb) == ctx.theme.color("primary")
    body_runs = [r for p in body_box.text_frame.paragraphs for r in p.runs]
    assert str(body_runs[0].font.color.rgb) == ctx.theme.color("ink")

    # circle vertically centered against the (taller) text stack
    circle = autoshapes(slide)[0]
    assert abs((int(circle.top) - origin.y) - (consumed - d) // 2) <= pt(1)
    # multi-line body makes the stack taller than the circle
    assert consumed > d


def test_short_stack_height_floors_at_circle_diameter():
    ctx = make_ctx()
    comp = NumberedBlock()
    data = NumberedBlockSpec(number="4", title="Go")
    measured = comp.measure(data, inch(4), ctx)
    assert measured >= expected_diameter(ctx)
