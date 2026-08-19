"""Census-form wave: hub_spoke, staircase, venn, xy_chart, combo chart.

Build order came from corpus/work/form_census.json — these are the forms
big-firm decks actually use that the engine lacked.
"""
from __future__ import annotations

from pptx import Presentation
from pptx.oxml.ns import qn

from deckengine.components.base import BuildReport, RenderContext, \
    get_component
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.llm.canon import CANON
from deckengine.llm.format_rules import first_rule, signals_from
from deckengine.schema.components import (ChartSeriesSpec, ChartStyleSpec,
                                          HubSpokeSpec, NativeChartSpec,
                                          SpokeSpec, StaircaseSpec,
                                          StairStepSpec, VennSpec,
                                          XYChartSpec, XYPointSpec)


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _hub():
    return HubSpokeSpec(hub="Core", spokes=[
        SpokeSpec(label="Payments", sub="12M users"),
        SpokeSpec(label="Logistics"), SpokeSpec(label="Credit"),
        SpokeSpec(label="Ads")], highlight_index=2)


def _stairs():
    return StaircaseSpec(steps=[
        StairStepSpec(label="Stabilize", value="Y1"),
        StairStepSpec(label="Scale", value="Y2"),
        StairStepSpec(label="Lead", value="Y3")])


def _venn():
    return VennSpec(circles=["Trust", "Reach", "Economics"],
                    intersection="The winning zone")


def _xy(**kw):
    kw.setdefault("points", [
        XYPointSpec(label="A", x=1.0, y=10.0, size=3.0),
        XYPointSpec(label="B", x=4.0, y=22.0, size=6.0),
        XYPointSpec(label="C", x=6.0, y=8.0, size=2.0)])
    kw.setdefault("x_label", "Growth")
    kw.setdefault("y_label", "Size")
    return XYChartSpec(**kw)


def test_measure_render_parity_all_forms():
    ctx = make_ctx()
    for kind, spec in (("hub_spoke", _hub()), ("staircase", _stairs()),
                       ("venn", _venn()), ("xy_chart", _xy())):
        comp = get_component(kind)
        _, slide = blank_slide()
        m = comp.measure(spec, inch(9), ctx)
        c = comp.render(slide, spec, BBox(0, 0, inch(9), inch(6)), ctx)
        assert abs(m - c) <= pt(1), f"{kind}: measured {m} consumed {c}"


def test_hub_spoke_highlight_and_range_guard():
    ctx = make_ctx()
    comp = get_component("hub_spoke")
    _, slide = blank_slide()
    comp.render(slide, _hub(), BBox(0, 0, inch(9), inch(5)), ctx)
    fills = [str(s.fill.fore_color.rgb) for s in slide.shapes
             if s.shape_type is not None and s.has_text_frame
             and s.fill.type is not None]
    assert ctx.theme.color("accent") in fills  # highlighted chip
    bad = _hub().model_copy(update={"highlight_index": 7})
    _, slide2 = blank_slide()
    comp.render(slide2, bad, BBox(0, 0, inch(9), inch(5)), ctx)
    assert any("out of range" in w for w in ctx.report.warnings)


def test_staircase_steps_ascend_and_last_accented():
    ctx = make_ctx()
    comp = get_component("staircase")
    _, slide = blank_slide()
    comp.render(slide, _stairs(), BBox(0, 0, inch(9), inch(4)), ctx)
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    bars = sorted((s for s in slide.shapes
                   if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE),
                  key=lambda s: s.left)
    heights = [s.height for s in bars[:3]]
    assert heights == sorted(heights)          # ascending
    assert str(bars[2].fill.fore_color.rgb) == ctx.theme.color("accent")
    bottoms = {s.top + s.height for s in bars[:3]}
    assert len(bottoms) == 1                   # bottom-aligned


def test_venn_translucent_fills_written():
    ctx = make_ctx()
    comp = get_component("venn")
    _, slide = blank_slide()
    comp.render(slide, _venn(), BBox(0, 0, inch(9), inch(5)), ctx)
    ovals = [s for s in slide.shapes
             if s._element.find(qn("p:spPr")) is not None
             and b"alpha" in s._element.xml.encode()
             if isinstance(s._element.xml, str)]
    alphas = [s for s in slide.shapes if "alpha" in s._element.xml]
    assert len(alphas) >= 3                    # every circle is translucent


def test_xy_quadrants_draw_midlines_and_captions():
    ctx = make_ctx()
    comp = get_component("xy_chart")
    _, slide = blank_slide()
    comp.render(slide, _xy(quadrants=True,
                           quadrant_labels=["TL", "TR", "BL", "BR"],
                           highlight="B"),
                BBox(0, 0, inch(9), inch(5)), ctx)
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    lines = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.LINE]
    assert len(lines) == 2                     # the two dashed midlines
    texts = " ".join(r.text for s in slide.shapes if s.has_text_frame
                     for p in s.text_frame.paragraphs for r in p.runs)
    for cap in ("TL", "TR", "BL", "BR"):
        assert cap in texts


def test_xy_bad_highlight_warns():
    ctx = make_ctx()
    _, slide = blank_slide()
    get_component("xy_chart").render(slide, _xy(highlight="Nope"),
                                     BBox(0, 0, inch(9), inch(5)), ctx)
    assert any("matches no point" in w for w in ctx.report.warnings)


def test_combo_chart_splits_line_to_secondary_axis():
    ctx = make_ctx()
    spec = NativeChartSpec(
        chart_type="combo", categories=["FY24", "FY25", "FY26"],
        series=[ChartSeriesSpec(name="Revenue", values=[100.0, 130.0, 170.0]),
                ChartSeriesSpec(name="Margin %", values=[8.0, 10.0, 12.0])],
        sort="none", highlight=None, annotation=None,
        style=ChartStyleSpec(combo_line_series="Margin %"))
    _, slide = blank_slide()
    get_component("native_chart").render(slide, spec,
                                         BBox(0, 0, inch(9), inch(5)), ctx)
    chart = next(s.chart for s in slide.shapes if s.has_chart)
    xml = chart._chartSpace.xml
    assert "<c:lineChart>" in xml or "lineChart" in xml
    plotArea = chart._chartSpace.find(qn("c:chart")).find(qn("c:plotArea"))
    assert plotArea.find(qn("c:lineChart")) is not None
    assert len(plotArea.findall(qn("c:valAx"))) == 2   # secondary axis
    bar_sers = plotArea.find(qn("c:barChart")).findall(qn("c:ser"))
    line_sers = plotArea.find(qn("c:lineChart")).findall(qn("c:ser"))
    assert len(bar_sers) == 1 and len(line_sers) == 1


def test_combo_unknown_line_series_degrades():
    ctx = make_ctx()
    spec = NativeChartSpec(
        chart_type="combo", categories=["A", "B"],
        series=[ChartSeriesSpec(name="X", values=[1.0, 2.0]),
                ChartSeriesSpec(name="Y", values=[3.0, 4.0])],
        sort="none", highlight=None, annotation=None,
        style=ChartStyleSpec(combo_line_series="Ghost"))
    _, slide = blank_slide()
    get_component("native_chart").render(slide, spec,
                                         BBox(0, 0, inch(9), inch(4)), ctx)
    assert any("combo_line_series" in w for w in ctx.report.warnings)


# -- routing + canon ---------------------------------------------------------

def test_new_forms_route():
    assert first_rule(signals_from(
        "The sweet spot sits where trust overlaps distribution")) \
        .id == "overlap_venn"
    assert first_rule(signals_from(
        "The platform core anchors a partner ecosystem")).id \
        == "ecosystem_hub"
    assert first_rule(signals_from(
        "A three-level maturity ladder climbs from pilots to scale")).id \
        == "value_staircase"
    r = first_rule(signals_from("Adoption scales with income"))
    assert r.id == "correlation_scatter" and "xy_chart" in r.then


def test_combo_rule_needs_trend_and_periods():
    from deckengine.llm.facts import FactTable
    t = FactTable()
    t.add("y1", "2024", "period 2024")
    t.add("y2", "2026", "period 2026")
    s = signals_from("EBITDA margin expanded alongside revenue growth "
                     "since 2024", t)
    assert first_rule(s).id == "combo_two_units"
    # without periods it must NOT hijack the trend family
    s2 = signals_from("Margins support revenue quality")
    assert (first_rule(s2) or type("R", (), {"id": ""})).id \
        != "combo_two_units"


def test_canon_gaps_flipped():
    for form in ("scatter", "bubble", "quadrant_scatter", "venn",
                 "staircase", "hub_spoke", "combo_line_column"):
        assert CANON[form].status == "primitive", form
