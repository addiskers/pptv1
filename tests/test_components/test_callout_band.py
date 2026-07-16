from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.shapes.connector import Connector

from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.callout_band import CalloutBand
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import CalloutBandSpec


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def test_measure_equals_render_full_band():
    ctx = make_ctx()
    comp = CalloutBand()
    data = CalloutBandSpec(label="TOTAL", icon="$",
                           segments=["**$1.2B** committed to date",
                                     "Across **14** countries",
                                     "Since program launch in 2008"])
    width = inch(11)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(2)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_measure_equals_render_minimal_band():
    ctx = make_ctx()
    comp = CalloutBand()
    data = CalloutBandSpec(segments=[
        "A single closing statement that summarizes the total impact of the "
        "portfolio across every region and program cycle to date."])
    width = inch(5)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(3)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_label_adds_exact_tab_overhang_and_min_height():
    ctx = make_ctx()
    comp = CalloutBand()
    plain = CalloutBandSpec(segments=["Alpha", "Beta"])
    labeled = CalloutBandSpec(label="TOTAL", segments=["Alpha", "Beta"])
    m_plain = comp.measure(plain, inch(9), ctx)
    m_labeled = comp.measure(labeled, inch(9), ctx)
    assert m_labeled - m_plain == ctx.theme.spacing(0.6)
    assert m_plain >= inch(0.9)


def test_vlines_between_segments():
    ctx = make_ctx()
    _, slide = blank_slide()
    CalloutBand().render(slide, CalloutBandSpec(segments=["One", "Two", "Three"]),
                         BBox(0, 0, inch(10), inch(1.5)), ctx)
    vlines = [s for s in slide.shapes if isinstance(s, Connector)]
    assert len(vlines) == 2  # n segments -> n-1 separators


def test_tab_and_icon_rendered():
    ctx = make_ctx()
    _, slide = blank_slide()
    data = CalloutBandSpec(label="TOTAL", icon="$", segments=["Only segment"])
    CalloutBand().render(slide, data, BBox(0, inch(1), inch(8), inch(2)), ctx)
    autoshapes = [s for s in slide.shapes
                  if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    kinds = {s.auto_shape_type for s in autoshapes}
    assert MSO_SHAPE.RECTANGLE in kinds          # the band panel
    assert MSO_SHAPE.OVAL in kinds               # the icon circle
    assert MSO_SHAPE.ROUNDED_RECTANGLE in kinds  # the label tab
    tab = next(s for s in autoshapes
               if s.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE)
    band = next(s for s in autoshapes
                if s.auto_shape_type == MSO_SHAPE.RECTANGLE)
    # tab's vertical center sits on the band's top edge, tab top at bbox top
    assert tab.top + tab.height // 2 == band.top
    assert tab.top == inch(1)
    runs = [r for s in autoshapes for p in s.text_frame.paragraphs
            for r in p.runs]
    label_runs = [r for r in runs if r.text == "TOTAL"]
    assert label_runs and label_runs[0].font.bold
    assert any(r.text == "$" for r in runs)
