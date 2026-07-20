from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.chevron_pathway import ChevronPathway
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import ChevronPathwaySpec


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def four_steps() -> ChevronPathwaySpec:
    return ChevronPathwaySpec(
        steps=["Discover", "Validate", "Scale", "Sustain"],
        highlight_index=2)


def two_steps() -> ChevronPathwaySpec:
    return ChevronPathwaySpec(
        steps=["Country-led diagnosis and priority setting",
               "Coordinated donor investment"])


def test_measure_equals_render_four_steps():
    ctx = make_ctx()
    comp = ChevronPathway()
    data = four_steps()
    width = inch(9)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(1.5)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_measure_equals_render_two_long_steps():
    ctx = make_ctx()
    comp = ChevronPathway()
    data = two_steps()
    width = inch(7)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(1.5)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_highlight_chevron_fill_differs_from_others():
    ctx = make_ctx()
    _, slide = blank_slide()
    ChevronPathway().render(slide, four_steps(),
                            BBox(0, 0, inch(10), inch(1.5)), ctx)
    chevrons = [s for s in slide.shapes
                if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert len(chevrons) == 4
    fills = [str(c.fill.fore_color.rgb) for c in chevrons]
    assert fills[2] == ctx.theme.color("accent")
    for i in (0, 1, 3):
        assert fills[i] == ctx.theme.color("primary_dark")
        assert fills[i] != fills[2]


def test_no_highlight_means_uniform_primary_dark():
    ctx = make_ctx()
    _, slide = blank_slide()
    ChevronPathway().render(slide, two_steps(),
                            BBox(0, 0, inch(8), inch(1.5)), ctx)
    chevrons = [s for s in slide.shapes
                if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert len(chevrons) == 2
    assert {str(c.fill.fore_color.rgb) for c in chevrons} == \
        {ctx.theme.color("primary_dark")}


def test_step_text_bold_inverse_and_within_two_lines():
    ctx = make_ctx()
    _, slide = blank_slide()
    consumed = ChevronPathway().render(slide, two_steps(),
                                       BBox(0, 0, inch(8), inch(1.5)), ctx)
    chevrons = [s for s in slide.shapes
                if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    for chev in chevrons:
        assert chev.height == consumed
        paras = [p for p in chev.text_frame.paragraphs if p.runs]
        assert 1 <= len(paras) <= 2
        for p in paras:
            for r in p.runs:
                assert r.font.bold
                assert (str(r.font.color.rgb)
                        == ctx.theme.color("inverse_ink"))
    text = " ".join(r.text for c in chevrons
                    for p in c.text_frame.paragraphs for r in p.runs)
    assert "Country-led" in text
    assert "donor" in text


def test_out_of_range_highlight_warns_and_falls_back():
    ctx = make_ctx()
    _, slide = blank_slide()
    data = ChevronPathwaySpec(steps=["One", "Two"], highlight_index=5)
    ChevronPathway().render(slide, data, BBox(0, 0, inch(6), inch(1.5)), ctx)
    assert any("highlight_index" in w for w in ctx.report.warnings)
    chevrons = [s for s in slide.shapes
                if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert {str(c.fill.fore_color.rgb) for c in chevrons} == \
        {ctx.theme.color("primary_dark")}
