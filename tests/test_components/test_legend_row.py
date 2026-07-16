from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.legend_row import LegendRow
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import LegendItemSpec, LegendRowSpec


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def test_measure_equals_render_two_items():
    ctx = make_ctx()
    comp = LegendRow()
    data = LegendRowSpec(items=[LegendItemSpec(code="LIV", label="Livestock"),
                                LegendItemSpec(code="CRO", label="Crops")])
    width = inch(5)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(1)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_measure_equals_render_many_items():
    ctx = make_ctx()
    comp = LegendRow()
    data = LegendRowSpec(items=[
        LegendItemSpec(code="LIV", label="Livestock"),
        LegendItemSpec(code="CRO", label="Crops"),
        LegendItemSpec(code="ADS", label="Access & delivery systems"),
        LegendItemSpec(code="PAC", label="Policy & advocacy"),
        LegendItemSpec(code="DIR", label="Direct"),
    ])
    width = inch(11)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(0.5)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_pill_uses_badge_color_with_white_bold_code():
    ctx = make_ctx()
    _, slide = blank_slide()
    data = LegendRowSpec(items=[LegendItemSpec(code="CRO", label="Crops")])
    LegendRow().render(slide, data, BBox(0, 0, inch(4), inch(0.5)), ctx)
    pills = [s for s in slide.shapes
             if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert len(pills) == 1
    pill = pills[0]
    assert str(pill.fill.fore_color.rgb) == ctx.theme.badge_palette["CRO"]
    run = pill.text_frame.paragraphs[0].runs[0]
    assert run.text == "CRO"
    assert run.font.bold
    assert str(run.font.color.rgb) == ctx.theme.color("inverse_ink")
    # label follows in its own text box
    boxes = [s for s in slide.shapes
             if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX]
    assert len(boxes) == 1
    label_run = boxes[0].text_frame.paragraphs[0].runs[0]
    assert label_run.text == "Crops"
    assert str(label_run.font.color.rgb) == ctx.theme.color("ink")
    assert boxes[0].left >= pill.left + pill.width  # label sits after the pill


def test_overflow_warns_but_still_renders():
    ctx = make_ctx()
    _, slide = blank_slide()
    data = LegendRowSpec(items=[
        LegendItemSpec(code="LIV", label="A very long legend label indeed"),
        LegendItemSpec(code="CRO", label="Another very long legend label"),
    ])
    LegendRow().render(slide, data, BBox(0, 0, inch(1), inch(0.5)), ctx)
    assert any("legend_row" in w for w in ctx.report.warnings)
    assert len([s for s in slide.shapes
                if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]) == 2


def test_unknown_codes_get_distinct_cycle_colors():
    """Unknown codes take DISTINCT colors from a deterministic cycle — an
    all-navy legend defeats the point of a legend."""
    ctx = make_ctx()
    _, slide = blank_slide()
    data = LegendRowSpec(items=[LegendItemSpec(code="ZZZ", label="Mystery"),
                                LegendItemSpec(code="YYY", label="Enigma")])
    LegendRow().render(slide, data, BBox(0, 0, inch(6), inch(0.5)), ctx)
    pills = [s for s in slide.shapes
             if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert len(pills) == 2
    fills = [str(p.fill.fore_color.rgb) for p in pills]
    assert fills[0] == ctx.theme.color("primary")
    assert fills[0] != fills[1]
