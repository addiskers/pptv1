from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.highlight_box import HighlightBox
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import HighlightBoxSpec


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def rects(slide):
    return [s for s in slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]


def test_measure_equals_render_title_and_body():
    ctx = make_ctx()
    comp = HighlightBox()
    data = HighlightBoxSpec(
        title="**Key takeaway:** scale-up is on track",
        body="Adoption reached 62% of target districts by 2024, and unit "
             "costs fell 18% as delivery moved to partner networks.")
    width = inch(4)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data,
                           BBox(inch(1), inch(1), width, inch(3)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_measure_equals_render_title_only_no_bar():
    ctx = make_ctx()
    comp = HighlightBox()
    data = HighlightBoxSpec(title="Program economics improve with scale",
                            accent_bar=False)
    width = inch(3)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data,
                           BBox(inch(0.5), inch(0.5), width, inch(2)), ctx)
    assert abs(measured - consumed) <= pt(1)
    assert len(rects(slide)) == 1  # panel only, no accent bar


def test_accent_bar_geometry_and_fills():
    ctx = make_ctx()
    _, slide = blank_slide()
    data = HighlightBoxSpec(title="Takeaway", body="Supporting sentence.")
    origin = BBox(inch(1), inch(1), inch(4), inch(3))
    consumed = HighlightBox().render(slide, data, origin, ctx)
    shapes = rects(slide)
    assert len(shapes) == 2
    panel = max(shapes, key=lambda s: int(s.width))
    bar = min(shapes, key=lambda s: int(s.width))
    # bar: spacing(0.35) wide, flush left, full panel height, accent fill
    assert int(bar.width) == ctx.theme.spacing(0.35)
    assert int(bar.left) == origin.x
    assert abs(int(bar.height) - consumed) <= pt(1)
    assert bar.fill.fore_color.rgb == RGBColor.from_string(
        ctx.theme.color("accent"))
    # panel: spec fill role, height equals consumed height
    assert panel.fill.fore_color.rgb == RGBColor.from_string(
        ctx.theme.color("surface_alt"))
    assert abs(int(panel.height) - consumed) <= pt(1)


def test_title_and_body_text_written_inside_padding():
    ctx = make_ctx()
    _, slide = blank_slide()
    data = HighlightBoxSpec(title="So what", body="Detail line.")
    origin = BBox(0, 0, inch(4), inch(2))
    HighlightBox().render(slide, data, origin, ctx)
    boxes = [s for s in slide.shapes
             if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX]
    assert len(boxes) == 2
    texts = ["".join(r.text for p in b.text_frame.paragraphs for r in p.runs)
             for b in boxes]
    assert "So what" in texts and "Detail line." in texts
    pad = ctx.theme.spacing(0.8)
    assert all(int(b.left) == origin.x + pad for b in boxes)
