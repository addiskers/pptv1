from pptx import Presentation

from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.footnote_strip import FootnoteStrip
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import FootnoteStripSpec


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _text_runs(slide):
    return [r for sh in slide.shapes if sh.has_text_frame
            for p in sh.text_frame.paragraphs for r in p.runs]


def test_measure_equals_render_single_note():
    ctx = make_ctx()
    comp = FootnoteStrip()
    data = FootnoteStripSpec(notes=["Source: World Bank commodity data, 2025"])
    width = inch(6)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(1)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_measure_equals_render_wrapping_notes():
    ctx = make_ctx()
    comp = FootnoteStrip()
    data = FootnoteStripSpec(notes=[
        "Source: FAO agricultural census 2024",
        "Excludes smallholdings under 0.5 ha",
        "Currency figures in constant 2020 USD",
    ])
    width = inch(2.5)  # narrow: forces multi-line wrap
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data,
                           BBox(inch(0.3), inch(6), width, inch(1.2)), ctx)
    assert abs(measured - consumed) <= pt(1)
    assert measured > comp.measure(
        FootnoteStripSpec(notes=["short"]), width, ctx)


def test_notes_joined_muted_and_micro():
    ctx = make_ctx()
    _, slide = blank_slide()
    FootnoteStrip().render(
        slide,
        FootnoteStripSpec(notes=["Note one", "Note two", "Note three"]),
        BBox(0, 0, inch(8), inch(0.5)), ctx)
    runs = _text_runs(slide)
    text = "".join(r.text for r in runs)
    assert "Note one | Note two | Note three" == text
    muted = ctx.theme.color("ink_muted").upper()
    assert runs and all(str(r.font.color.rgb).upper() == muted for r in runs)
    assert all(r.font.size.pt <= ctx.size("micro") for r in runs)
