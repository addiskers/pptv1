from pptx import Presentation

from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.stat_row import StatRow
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import StatRowSpec, StatSpec


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


THREE_STATS = StatRowSpec(stats=[
    StatSpec(label="Portfolio value", value="$1.2B"),
    StatSpec(label="Active grants", value="312"),
    StatSpec(label="Countries", value="14"),
])

ONE_STAT_NO_RULE = StatRowSpec(stats=[
    StatSpec(label="Yield gain", value="+38%", underline=False),
])


def test_measure_equals_render_three_stats():
    ctx = make_ctx()
    comp = StatRow()
    width = inch(9)
    measured = comp.measure(THREE_STATS, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, THREE_STATS, BBox(0, 0, width, inch(2)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_measure_equals_render_no_underline():
    ctx = make_ctx()
    comp = StatRow()
    width = inch(3)
    measured = comp.measure(ONE_STAT_NO_RULE, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, ONE_STAT_NO_RULE,
                           BBox(0, 0, width, inch(2)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_renders_labels_values_and_rules():
    ctx = make_ctx()
    _, slide = blank_slide()
    StatRow().render(slide, THREE_STATS, BBox(0, 0, inch(9), inch(2)), ctx)
    text_shapes = [s for s in slide.shapes if s.has_text_frame]
    rules = [s for s in slide.shapes if not s.has_text_frame]
    assert len(text_shapes) == 6   # 3 labels + 3 values
    assert len(rules) == 3         # one underline per stat
    bold_runs = [r for s in text_shapes for p in s.text_frame.paragraphs
                 for r in p.runs if r.font.bold]
    assert {r.text for r in bold_runs} == {"$1.2B", "312", "14"}


def test_no_underline_means_no_rules_and_shorter_stack():
    ctx = make_ctx()
    with_rule = StatRowSpec(stats=[StatSpec(label="Yield gain", value="+38%")])
    h_no = StatRow().measure(ONE_STAT_NO_RULE, inch(3), ctx)
    h_yes = StatRow().measure(with_rule, inch(3), ctx)
    assert h_no < h_yes
    _, slide = blank_slide()
    StatRow().render(slide, ONE_STAT_NO_RULE, BBox(0, 0, inch(3), inch(2)), ctx)
    assert all(s.has_text_frame for s in slide.shapes)  # no connectors drawn
