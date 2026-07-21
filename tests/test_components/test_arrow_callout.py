from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE

from deckengine.components.arrow_callout import ArrowCallout
from deckengine.components.base import BuildReport, RenderContext
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import ArrowCalloutSpec, ArrowStatSpec


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def reference_spec() -> ArrowCalloutSpec:
    return ArrowCalloutSpec(
        title="AgDev has deployed **$336M** in India",
        sub="'Paid' + 'unpaid' investments marked as India, 2015 to current",
        stats=[ArrowStatSpec(value="185", label="AgDev funded investments"),
               ArrowStatSpec(value="123", label="Managed by Seattle POs"),
               ArrowStatSpec(value="62", label="Managed by India POs")])


def test_measure_equals_render_with_stats():
    ctx = make_ctx()
    comp = ArrowCallout()
    data = reference_spec()
    width = inch(11)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(3)), ctx)
    assert abs(measured - consumed) <= pt(1)
    assert measured >= inch(1.0)


def test_measure_equals_render_without_stats():
    ctx = make_ctx()
    comp = ArrowCallout()
    data = ArrowCalloutSpec(title="One clear statement, no trailing stats")
    width = inch(9)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(2)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_arrow_head_and_dashed_separators():
    ctx = make_ctx()
    comp = ArrowCallout()
    data = reference_spec()
    bbox = BBox(inch(0.5), inch(5.0), inch(11), inch(2))
    _, slide = blank_slide()
    consumed = comp.render(slide, data, bbox, ctx)

    pentagons = [s for s in slide.shapes
                 if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                 and s.auto_shape_type == MSO_SHAPE.PENTAGON]
    assert len(pentagons) == 1
    arrow = pentagons[0]
    box_right = bbox.x + round(bbox.w * 0.42)
    assert arrow.left == box_right          # protrudes from the box edge
    assert arrow.width == inch(0.5)
    assert arrow.height <= inch(0.55)

    # dashed vlines between the 3 stat columns
    vlines = [s for s in slide.shapes
              if s.shape_type == MSO_SHAPE_TYPE.LINE and s.width == 0]
    assert len(vlines) == 2
    assert all(s.left > arrow.left + arrow.width for s in vlines)
    assert all(s.height < consumed for s in vlines)  # vertically inset

    text = " ".join(r.text for s in slide.shapes if s.has_text_frame
                    for p in s.text_frame.paragraphs for r in p.runs)
    assert "$336M" in text
    assert "185" in text
    assert "Seattle" in text


def test_no_stats_box_spans_width_minus_arrow():
    ctx = make_ctx()
    comp = ArrowCallout()
    data = ArrowCalloutSpec(title="Statement only")
    bbox = BBox(0, 0, inch(9), inch(2))
    _, slide = blank_slide()
    comp.render(slide, data, bbox, ctx)
    rects = [s for s in slide.shapes
             if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             and s.auto_shape_type == MSO_SHAPE.RECTANGLE]
    assert len(rects) == 1
    assert rects[0].width == inch(9) - inch(0.5)
