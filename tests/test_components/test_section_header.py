from pptx import Presentation
from pptx.shapes.connector import Connector

from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.section_header import SectionHeader
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import SectionHeaderSpec


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _text_runs(slide):
    return [r for sh in slide.shapes if sh.has_text_frame
            for p in sh.text_frame.paragraphs for r in p.runs]


def test_measure_equals_render_title_only():
    ctx = make_ctx()
    comp = SectionHeader()
    data = SectionHeaderSpec(
        title="Smallholder productivity gains concentrated in 6 geographies")
    width = inch(9)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(2)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_measure_equals_render_with_subtitle():
    ctx = make_ctx()
    comp = SectionHeader()
    data = SectionHeaderSpec(
        title="Crop yields **doubled** across focus geographies since 2010",
        subtitle="Preliminary read: **142** country programs, FAO census basis")
    width = inch(6)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data,
                           BBox(inch(0.5), inch(0.4), width, inch(3)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_rule_adds_height_and_draws_connector():
    ctx = make_ctx()
    comp = SectionHeader()
    width = inch(8)
    with_rule = SectionHeaderSpec(title="Portfolio overview", rule=True)
    without = SectionHeaderSpec(title="Portfolio overview", rule=False)
    assert comp.measure(with_rule, width, ctx) > comp.measure(without, width, ctx)

    _, slide = blank_slide()
    comp.render(slide, with_rule, BBox(0, 0, width, inch(2)), ctx)
    assert any(isinstance(sh, Connector) for sh in slide.shapes)

    _, bare = blank_slide()
    comp.render(bare, without, BBox(0, 0, width, inch(2)), ctx)
    assert not any(isinstance(sh, Connector) for sh in bare.shapes)


def test_subtitle_color_follows_accent_flag():
    ctx = make_ctx()
    comp = SectionHeader()
    accent = ctx.theme.color("accent").upper()
    muted = ctx.theme.color("ink_muted").upper()

    _, slide = blank_slide()
    comp.render(slide,
                SectionHeaderSpec(title="T", subtitle="grant pipeline view",
                                  rule=False, accent_subtitle=True),
                BBox(0, 0, inch(8), inch(2)), ctx)
    colors = {str(r.font.color.rgb).upper() for r in _text_runs(slide)}
    assert accent in colors

    _, slide2 = blank_slide()
    comp.render(slide2,
                SectionHeaderSpec(title="T", subtitle="grant pipeline view",
                                  rule=False, accent_subtitle=False),
                BBox(0, 0, inch(8), inch(2)), ctx)
    colors2 = {str(r.font.color.rgb).upper() for r in _text_runs(slide2)}
    assert muted in colors2
    assert accent not in colors2
