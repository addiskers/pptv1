"""The fit contract's last stragglers: components that used to paint past a
too-short canvas cell now CLAMP into it (shrink -> ellipsize/drop, always
reported) — 'never overlaps' holds for every component."""
from __future__ import annotations

from pptx import Presentation

from deckengine.components.base import (BuildReport, RenderContext,
                                        get_component)
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch
from deckengine.schema.components import (BraceGroupSpec, CalloutBandSpec,
                                          IconStatRowSpec, NumberedBlockSpec)


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


_LONG = ("A rather long segment of evidence text that would naturally wrap "
         "onto several lines and demand far more height than the tiny box")


def _starved_render(kind, spec, h):
    ctx = make_ctx()
    slide = blank_slide()
    box = BBox(0, 0, inch(8), h)
    consumed = get_component(kind).render(slide, spec, box, ctx)
    return consumed, ctx.report, box


def test_callout_band_clamps():
    spec = CalloutBandSpec(label="SO WHAT", segments=[_LONG, _LONG, _LONG])
    natural = get_component("callout_band").measure(spec, inch(8), make_ctx())
    assert natural > inch(0.6)              # genuinely starved below
    consumed, report, box = _starved_render("callout_band", spec, inch(0.6))
    assert consumed <= box.h + inch(0.02)   # never paints past the cell
    assert not any("taller than provided bbox" in w for w in report.warnings)
    assert report.truncations                # the cut is REPORTED


def test_icon_stat_row_clamps():
    spec = IconStatRowSpec(icon="growth", stat="8.8%", text=_LONG)
    consumed, report, box = _starved_render("icon_stat_row", spec, inch(0.5))
    assert consumed <= box.h + inch(0.02)
    assert not any("taller than provided bbox" in w for w in report.warnings)


def test_numbered_block_clamps():
    spec = NumberedBlockSpec(number="1", title="A titled step", body=_LONG)
    consumed, report, box = _starved_render("numbered_block", spec, inch(0.6))
    assert consumed <= box.h + inch(0.02)
    assert not any("taller than provided bbox" in w for w in report.warnings)


def test_brace_group_drops_trailing_children():
    spec = BraceGroupSpec(
        content=[{"kind": "text_block", "text": _LONG},
                 {"kind": "text_block", "text": _LONG},
                 {"kind": "text_block", "text": _LONG}],
        takeaway="The takeaway that matters most")
    consumed, report, box = _starved_render("brace_group", spec, inch(1.0))
    assert consumed <= box.h + inch(0.05)
    assert any("dropped" in t for t in report.truncations)


def test_natural_boxes_unchanged():
    """Stacker paths hand components their measured height — the clamp must
    be a strict no-op there (measure==render parity holds)."""
    for kind, spec in [
        ("callout_band", CalloutBandSpec(label="SO WHAT",
                                         segments=["a point", "b point"])),
        ("icon_stat_row", IconStatRowSpec(icon="growth", stat="8.8%",
                                          text="short text")),
        ("numbered_block", NumberedBlockSpec(number="1", title="Step",
                                             body="short body")),
    ]:
        ctx = make_ctx()
        natural = get_component(kind).measure(spec, inch(8), ctx)
        slide = blank_slide()
        consumed = get_component(kind).render(
            slide, spec, BBox(0, 0, inch(8), natural), ctx)
        assert abs(consumed - natural) <= inch(0.01), kind
        assert ctx.report.warnings == [], kind
