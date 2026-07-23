from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.funnel import Funnel
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import FunnelSpec, FunnelStageSpec


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def four_stages() -> FunnelSpec:
    return FunnelSpec(stages=[
        FunnelStageSpec(label="Total addressable market", value="$4.8B"),
        FunnelStageSpec(label="Serviceable market", value="$2.1B"),
        FunnelStageSpec(label="Qualified pipeline", value="$640M"),
        FunnelStageSpec(label="Closed-won revenue", value="$95M"),
    ])


def three_stages_fracs() -> FunnelSpec:
    # 1.4 and 0.1 are out of range: clamp to 1.0 and 0.25
    return FunnelSpec(stages=[
        FunnelStageSpec(label="Leads generated", value="12,400"),
        FunnelStageSpec(label="Opportunities", value="3,100"),
        FunnelStageSpec(label="Customers", value="410"),
    ], fracs=[1.4, 0.6, 0.1])


def bands_by_top(slide):
    return sorted((s for s in slide.shapes
                   if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE),
                  key=lambda s: s.top)


def all_run_text(slide) -> str:
    # runs concatenate within a paragraph (wrap tokens carry their spaces)
    return " ".join("".join(r.text for r in p.runs)
                    for s in slide.shapes if s.has_text_frame
                    for p in s.text_frame.paragraphs)


def test_measure_equals_render_default_taper():
    ctx = make_ctx()
    comp = Funnel()
    data = four_stages()
    width = inch(9)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data,
                           BBox(inch(0.5), inch(1), width, inch(4)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_measure_equals_render_explicit_fracs():
    ctx = make_ctx()
    comp = Funnel()
    data = three_stages_fracs()
    width = inch(8)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data,
                           BBox(inch(0.4), inch(1.2), width, inch(3)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_bands_narrow_centered_with_connectors():
    ctx = make_ctx()
    comp = Funnel()
    data = four_stages()
    bbox = BBox(inch(0.5), inch(1), inch(9), inch(4))
    _, slide = blank_slide()
    consumed = comp.render(slide, data, bbox, ctx)

    bands = bands_by_top(slide)
    assert len(bands) == 4
    widths = [s.width for s in bands]
    assert all(a > b for a, b in zip(widths, widths[1:]))
    lefts = [s.left for s in bands]
    assert all(a < b for a, b in zip(lefts, lefts[1:]))
    mid = bbox.x + bbox.w // 2
    for s in bands:
        assert abs((s.left + s.width // 2) - mid) <= 2
    assert widths[0] == bbox.w                    # taper starts at 1.0
    assert widths[-1] == round(bbox.w * 0.35)     # and ends at 0.35
    assert all(s.height >= inch(0.45) for s in bands)
    assert bands[-1].top + bands[-1].height == bbox.y + consumed

    # dashed connectors: one per gap, centered on the funnel axis
    vlines = [s for s in slide.shapes
              if s.shape_type == MSO_SHAPE_TYPE.LINE]
    assert len(vlines) == 3
    assert all(s.left == mid and s.width == 0 for s in vlines)


def test_explicit_fracs_clamped():
    ctx = make_ctx()
    comp = Funnel()
    data = three_stages_fracs()
    bbox = BBox(inch(0.4), inch(1.2), inch(8), inch(3))
    _, slide = blank_slide()
    comp.render(slide, data, bbox, ctx)

    bands = bands_by_top(slide)
    assert len(bands) == 3
    assert bands[0].width == bbox.w               # 1.4 clamped to 1.0
    assert bands[1].width == round(bbox.w * 0.6)
    assert bands[2].width == round(bbox.w * 0.25)  # 0.1 clamped to 0.25


def test_last_band_uses_accent_fill():
    ctx = make_ctx()
    comp = Funnel()
    _, slide = blank_slide()
    comp.render(slide, four_stages(), BBox(0, 0, inch(9), inch(4)), ctx)

    bands = bands_by_top(slide)
    assert str(bands[-1].fill.fore_color.rgb) == ctx.theme.color("accent")
    for s in bands[:-1]:
        assert str(s.fill.fore_color.rgb) == ctx.theme.color("primary")


def test_labels_and_values_present():
    ctx = make_ctx()
    comp = Funnel()
    data = four_stages()
    _, slide = blank_slide()
    comp.render(slide, data, BBox(inch(0.5), inch(1), inch(9), inch(4)), ctx)

    text = all_run_text(slide)
    for stage in data.stages:
        assert stage.label in text
        assert stage.value in text
    assert not ctx.report.truncations


def test_value_placement_inside_vs_outside():
    # top band spans the full width: its value must fall back inside the
    # band (inverse ink); narrower bands hang their value outside (ink)
    ctx = make_ctx()
    comp = Funnel()
    data = four_stages()
    bbox = BBox(inch(0.5), inch(1), inch(9), inch(4))
    _, slide = blank_slide()
    comp.render(slide, data, bbox, ctx)

    colors = {}
    for s in slide.shapes:
        if not s.has_text_frame:
            continue
        for p in s.text_frame.paragraphs:
            for r in p.runs:
                colors[r.text] = str(r.font.color.rgb)
    assert colors["$4.8B"] == ctx.theme.color("inverse_ink")
    assert colors["$95M"] == ctx.theme.color("ink")


def test_value_too_wide_for_band_warns():
    # bbox so narrow the value fits neither beside nor inside its band
    ctx = make_ctx()
    comp = Funnel()
    data = FunnelSpec(stages=[
        FunnelStageSpec(label="A", value="$1,234,567,890"),
        FunnelStageSpec(label="B", value="$1,234,567,890"),
        FunnelStageSpec(label="C", value="$1,234,567,890"),
    ])
    _, slide = blank_slide()
    comp.render(slide, data, BBox(inch(0.5), inch(1), inch(1), inch(4)), ctx)
    assert any("fits neither" in w for w in ctx.report.warnings)


def test_fill_hint_stretches_bands_capped():
    ctx = make_ctx()
    comp = Funnel()
    data = four_stages()
    width = inch(9)
    natural = comp.measure(data, width, ctx)
    ctx.fill_hint = True
    _, slide = blank_slide()
    consumed = comp.render(slide, data,
                           BBox(0, 0, width, natural * 3), ctx)
    assert consumed > natural
    assert consumed <= round(natural * 1.5) + pt(1)
