from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE

from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.icon_stat_row import IconStatRow
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import IconStatRowSpec


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def test_measure_equals_render_short():
    ctx = make_ctx()
    comp = IconStatRow()
    data = IconStatRowSpec(icon="+", stat="**2.1M**",
                           text="Farmers reached with improved seed varieties.")
    width = inch(6)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(2)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_measure_equals_render_wrapping_description():
    ctx = make_ctx()
    comp = IconStatRow()
    data = IconStatRowSpec(
        icon="W", stat="**48%**",
        text="Increase in smallholder productivity across program districts, "
             "measured against the 2019 baseline cohort of participants.")
    width = inch(4)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(3)), ctx)
    assert abs(measured - consumed) <= pt(1)
    # a wrapping description grows the row past the single-line minimum
    assert measured > inch(0.42)


def test_min_row_height():
    ctx = make_ctx()
    data = IconStatRowSpec(icon="+", stat="7", text="Countries.")
    assert IconStatRow().measure(data, inch(9), ctx) == inch(0.42)


def test_render_draws_circle_matching_row_height_and_content():
    ctx = make_ctx()
    data = IconStatRowSpec(icon="+", stat="**2.1M**",
                           text="Farmers reached with improved seed varieties.")
    _, slide = blank_slide()
    consumed = IconStatRow().render(slide, data, BBox(0, 0, inch(6), inch(2)),
                                    ctx)
    ovals = [s for s in slide.shapes
             if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             and s.auto_shape_type == MSO_SHAPE.OVAL]
    assert len(ovals) == 1
    # circle: diameter == row height, on both axes
    assert ovals[0].width == ovals[0].height == consumed
    assert str(ovals[0].fill.fore_color.rgb) == ctx.theme.color("primary_dark")
    text = " ".join(r.text for s in slide.shapes if s.has_text_frame
                    for p in s.text_frame.paragraphs for r in p.runs)
    assert "+" in text          # icon glyph
    assert "2.1M" in text       # stat
    assert "Farmers" in text    # description
