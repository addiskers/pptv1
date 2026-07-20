from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.connector import Connector
from pptx.util import Pt

from deckengine.components.base import BuildReport, RenderContext
from deckengine.components.timeline_row import TimelineRow
from deckengine.core.bbox import BBox
from deckengine.core.fit_text import TextMeasurer
from deckengine.core.theme import load_theme
from deckengine.core.units import inch, pt
from deckengine.schema.components import MilestoneSpec, TimelineRowSpec


def make_ctx() -> RenderContext:
    return RenderContext(theme=load_theme("consulting_navy"),
                         measurer=TextMeasurer(), report=BuildReport())


def blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def three_milestones() -> TimelineRowSpec:
    return TimelineRowSpec(milestones=[
        MilestoneSpec(label="Pilot launch", date="Q1 2026", done=True),
        MilestoneSpec(label="Scale to 3 regions", date="Q3 2026", done=False),
        MilestoneSpec(label="National rollout", date="Q2 2027", done=False),
    ])


def six_milestones() -> TimelineRowSpec:
    return TimelineRowSpec(milestones=[
        MilestoneSpec(label="Design", date="Jan", done=True),
        MilestoneSpec(label="Procurement", date="Mar", done=True),
        MilestoneSpec(label="Field trials", date="Jun", done=True),
        MilestoneSpec(label="Partner onboarding", date="Sep", done=False),
        MilestoneSpec(label="Regulatory approval", date="Nov", done=False),
        MilestoneSpec(label="Launch", date="Dec", done=False),
    ])


def test_measure_equals_render_three_milestones():
    ctx = make_ctx()
    comp = TimelineRow()
    data = three_milestones()
    width = inch(9)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(2)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_measure_equals_render_six_milestones():
    ctx = make_ctx()
    comp = TimelineRow()
    data = six_milestones()
    width = inch(11)
    measured = comp.measure(data, width, ctx)
    _, slide = blank_slide()
    consumed = comp.render(slide, data, BBox(0, 0, width, inch(2)), ctx)
    assert abs(measured - consumed) <= pt(1)


def test_done_and_pending_markers_get_different_fills():
    ctx = make_ctx()
    _, slide = blank_slide()
    TimelineRow().render(slide, three_milestones(),
                         BBox(0, 0, inch(9), inch(2)), ctx)
    markers = [s for s in slide.shapes
               if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert len(markers) == 3  # one oval per milestone, in order
    done, pending = markers[0], markers[1]
    assert str(done.fill.fore_color.rgb) == ctx.theme.color("primary")
    assert str(pending.fill.fore_color.rgb) == ctx.theme.color("bg")
    assert str(pending.line.color.rgb) == ctx.theme.color("primary")
    assert pending.line.width == Pt(1.25)
    assert str(done.fill.fore_color.rgb) != str(pending.fill.fore_color.rgb)
    # all markers are spacing(1.4) circles
    d = ctx.theme.spacing(1.4)
    assert all(m.width == d and m.height == d for m in markers)


def test_rule_spans_width_through_marker_band_middle():
    ctx = make_ctx()
    _, slide = blank_slide()
    x, y, w = inch(1), inch(1), inch(9)
    TimelineRow().render(slide, three_milestones(),
                         BBox(x, y, w, inch(2)), ctx)
    rules = [s for s in slide.shapes if isinstance(s, Connector)]
    assert len(rules) == 1
    rule = rules[0]
    assert rule.left == x
    assert rule.width == w
    assert rule.line.width == Pt(1.5)
    assert str(rule.line.color.rgb) == ctx.theme.color("grid")
    # rule passes through the vertical middle of every marker
    markers = [s for s in slide.shapes
               if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    for m in markers:
        mid = m.top + m.height // 2
        assert abs(rule.top - mid) <= pt(1)


def test_markers_evenly_spaced_between_6_and_94_percent():
    ctx = make_ctx()
    _, slide = blank_slide()
    x, w = inch(0.5), inch(10)
    TimelineRow().render(slide, three_milestones(),
                         BBox(x, 0, w, inch(2)), ctx)
    markers = [s for s in slide.shapes
               if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    centers = [m.left + m.width // 2 for m in markers]
    assert abs(centers[0] - (x + round(w * 0.06))) <= pt(1)
    assert abs(centers[-1] - (x + round(w * 0.94))) <= pt(1)
    assert abs((centers[1] - centers[0]) - (centers[2] - centers[1])) <= pt(1)


def test_dates_muted_above_and_labels_bold_below():
    ctx = make_ctx()
    _, slide = blank_slide()
    TimelineRow().render(slide, three_milestones(),
                         BBox(0, 0, inch(9), inch(2)), ctx)
    boxes = [s for s in slide.shapes
             if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX]
    assert len(boxes) == 6  # date + label per milestone
    runs = {r.text: r for b in boxes
            for p in b.text_frame.paragraphs for r in p.runs}
    date_run = runs["Q1 2026"]
    assert not date_run.font.bold
    assert str(date_run.font.color.rgb) == ctx.theme.color("ink_muted")
    label_run = next(r for t, r in runs.items() if "Pilot" in t)
    assert label_run.font.bold
    assert str(label_run.font.color.rgb) == ctx.theme.color("ink")
    # dates sit above the markers, labels below
    markers = [s for s in slide.shapes
               if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    marker_top = min(m.top for m in markers)
    marker_bottom = max(m.top + m.height for m in markers)
    date_boxes = [b for b in boxes
                  if any(r.text.startswith("Q") for p in b.text_frame.paragraphs
                         for r in p.runs)]
    label_boxes = [b for b in boxes if b not in date_boxes]
    assert all(b.top + b.height <= marker_top for b in date_boxes)
    assert all(b.top >= marker_bottom for b in label_boxes)
