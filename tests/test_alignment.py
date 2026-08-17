"""The alignment engine: rendered-output audit, grid snap, fill rail."""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches

from deckengine.render.audit import audit_slide
from deckengine.render.deck_builder import build_deck
from deckengine.llm.canvas_rules import check_canvas_slide
from deckengine.schema.canvas import CanvasSlideSpec
from deckengine.schema.slide_types import DeckSpec
from deckengine.slides.canvas_slide import _snap

REPO = Path(__file__).resolve().parents[1]
SW, SH = Inches(13.333), Inches(7.5)


def _slide():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH
    return prs.slides.add_slide(prs.slide_layouts[6])


def _box(slide, x, y, w, h, text, align=PP_ALIGN.LEFT):
    b = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    b.text_frame.text = text
    b.text_frame.paragraphs[0].alignment = align
    return b


# -- the audit ---------------------------------------------------------------

def test_audit_catches_text_clip():
    s = _slide()
    _box(s, 1, 1, 4, 1, "a long list row that gets covered")
    _box(s, 1, 1.4, 4, 1, "the band that covers it")
    res = audit_slide(s, int(SW), int(SH))
    assert res["overlaps"]


def test_audit_exempts_colocated_label_value():
    s = _slide()
    _box(s, 1, 1, 4, 0.5, "Innovations screened", PP_ALIGN.LEFT)
    _box(s, 1, 1, 4, 0.5, "115", PP_ALIGN.RIGHT)  # same box, right-aligned
    assert audit_slide(s, int(SW), int(SH))["overlaps"] == []


def test_audit_exempts_tagged_shapes():
    s = _slide()
    a = _box(s, 1, 1, 2, 0.5, "TOTAL")
    a.name = "de:tab"
    _box(s, 1, 1.2, 4, 1, "segment text under the tab")
    assert audit_slide(s, int(SW), int(SH))["overlaps"] == []


def test_edges_count_left_aligned_only():
    s = _slide()
    for i, x in enumerate((1.0, 1.06, 1.13, 2.5)):  # ragged cluster + one
        _box(s, x, 1 + i, 3, 0.5, f"row {i}")
    _box(s, 5.37, 1, 2, 0.5, "centered pill", PP_ALIGN.CENTER)  # ignored
    res = audit_slide(s, int(SW), int(SH))
    assert res["left_edges"] == 4  # 1.0/1.06/1.13 quantize to 3 + 2.5
    # (0.05in steps: 1.0 and 1.06 differ; the metric surfaces raggedness)


def test_deck_builder_reports_overlap_defect(tmp_path):
    spec = DeckSpec.model_validate({
        "schema_version": 1, "theme": "consulting_navy",
        "meta": {"title": "a", "date": "17 Aug 2026", "footer_org": "DE"},
        "slides": [{
            "slide_type": "canvas", "title": "Clip test",
            "elements": [
                {"x": 0.05, "y": 0.2, "w": 0.5, "h": 0.3, "z": 1,
                 "content": {"kind": "canvas_text",
                             "text": "first block of body text that is "
                                     "long enough to matter"}},
                {"x": 0.05, "y": 0.25, "w": 0.5, "h": 0.3, "z": 1,
                 "content": {"kind": "canvas_text",
                             "text": "second block painted over the first "
                                     "one"}},
            ]}]})
    report = build_deck(spec, tmp_path / "clip.pptx")
    assert any("text overlap" in w for w in report.warnings)


# -- grid snap ---------------------------------------------------------------

def _els(coords):
    return CanvasSlideSpec.model_validate({
        "slide_type": "canvas", "title": "T",
        "elements": [{"x": x, "y": y, "w": w, "h": h,
                      "content": {"kind": "canvas_text", "text": "t"}}
                     for x, y, w, h in coords]}).elements


def test_snap_quantizes_to_grid():
    geoms = _snap(_els([(0.26, 0.11, 0.3, 0.3), (0.51, 0.11, 0.3, 0.3)]))
    assert geoms[0][0] == 0.25          # 0.26 -> 6/24
    assert geoms[1][0] == 0.5           # 0.51 -> 12/24
    assert geoms[0][1] == geoms[1][1]   # same row line


def test_mutual_snap_before_grid_keeps_pairs_together():
    # 0.100 vs 0.108 straddle a 24-grid boundary; mutual snap first keeps
    # them on ONE shared edge
    geoms = _snap(_els([(0.100, 0.1, 0.3, 0.2), (0.108, 0.4, 0.3, 0.2)]))
    assert geoms[0][0] == geoms[1][0]


def test_thin_rules_keep_size():
    geoms = _snap(_els([(0.1, 0.1, 0.3, 0.2), (0.1, 0.5, 0.3, 0.02)]))
    assert abs(geoms[1][3] - 0.02) < 1e-9  # sub-cell height untouched


# -- fill rail ---------------------------------------------------------------

def test_dead_canvas_flagged():
    s = CanvasSlideSpec.model_validate({
        "slide_type": "canvas", "title": "Sparse",
        "elements": [
            {"x": 0.3, "y": 0.4, "w": 0.3, "h": 0.15,
             "content": {"kind": "canvas_text", "text": "lonely"}},
            {"x": 0.3, "y": 0.6, "w": 0.3, "h": 0.1,
             "content": {"kind": "canvas_text", "text": "also lonely"}},
        ]})
    assert any("covers only" in p for p in check_canvas_slide(s))


def test_golden_canvas_passes_fill_rail():
    golden = json.loads((REPO / "eval" / "goldens" / "canvas_hero.json")
                        .read_text(encoding="utf-8"))
    s = CanvasSlideSpec.model_validate(golden["slides"][0])
    assert not any("covers only" in p for p in check_canvas_slide(s))
