"""The Canvas: freeform slide spec, renderer, and design rails."""
from __future__ import annotations

import json
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Emu

from deckengine.llm.canvas_rules import check_canvas_slide
from deckengine.render.deck_builder import build_deck
from deckengine.schema.canvas import CanvasSlideSpec
from deckengine.schema.slide_types import DeckSpec

REPO = Path(__file__).resolve().parents[1]


def _canvas(elements, **kw):
    kw.setdefault("slide_type", "canvas")
    kw.setdefault("title", "The claim this slide argues")
    return CanvasSlideSpec.model_validate({**kw, "elements": elements})


def _text(x, y, w, h, text="hello", z=0, **c):
    c.setdefault("kind", "canvas_text")
    c["text"] = text
    return {"x": x, "y": y, "w": w, "h": h, "z": z, "content": c}


def _shape(x, y, w, h, z=0, **c):
    c.setdefault("kind", "canvas_shape")
    return {"x": x, "y": y, "w": w, "h": h, "z": z, "content": c}


# -- schema rails ------------------------------------------------------------

def test_offslide_element_rejected():
    with pytest.raises(Exception, match="off the slide"):
        _canvas([_text(0.8, 0.1, 0.5, 0.2), _text(0, 0, 0.2, 0.2)])


def test_full_slide_mode_requires_headline():
    with pytest.raises(Exception, match="headline"):
        _canvas([_text(0, 0, 0.5, 0.1, size_role="micro"),
                 _text(0, 0.5, 0.5, 0.1, size_role="micro")],
                render_title=False)
    # with a big text element it validates
    _canvas([_text(0, 0.1, 0.9, 0.2, size_pt=28),
             _text(0, 0.5, 0.5, 0.1)], render_title=False)


# -- design rails ------------------------------------------------------------

def test_text_collision_flagged_and_panel_exempt():
    s = _canvas([_text(0.1, 0.1, 0.4, 0.2, "alpha"),
                 _text(0.2, 0.15, 0.4, 0.2, "beta")])
    assert any("overlaps" in p for p in check_canvas_slide(s))
    # same geometry, but the lower element is a label-less PANEL: legal
    # (the fill rail may still note the sparse fixture — filter to overlaps)
    s2 = _canvas([_shape(0.05, 0.05, 0.6, 0.4, fill_role="surface"),
                  _text(0.1, 0.1, 0.4, 0.2, "alpha", z=1)])
    assert not [p for p in check_canvas_slide(s2) if "overlaps" in p]


def test_contrast_rail():
    s = _canvas([_shape(0.0, 0.0, 0.5, 1.0, fill_role="primary_dark"),
                 _text(0.05, 0.2, 0.35, 0.2, "dark on dark", z=1,
                       color_role="ink")])
    probs = check_canvas_slide(s)
    assert any("inverse_ink" in p for p in probs)
    s2 = _canvas([_shape(0.0, 0.0, 0.5, 1.0, fill_role="primary_dark"),
                  _text(0.05, 0.2, 0.35, 0.2, "readable", z=1,
                        color_role="inverse_ink")])
    assert check_canvas_slide(s2) == []


def test_dark_bg_contrast_rail():
    s = _canvas(
        [_text(0.1, 0.2, 0.5, 0.2, "lost text"),
         _text(0.1, 0.5, 0.5, 0.2, "more", color_role="inverse_ink")],
        bg_fill_role="primary_dark")
    assert any("bg_fill_role" in p for p in check_canvas_slide(s))


def test_sliver_text_flagged():
    s = _canvas([_text(0.1, 0.1, 0.04, 0.5, "a rather long sentence here"),
                 _text(0.5, 0.1, 0.3, 0.2, "fine")])
    assert any("widen" in p for p in check_canvas_slide(s))


def test_dark_shape_label_contrast():
    s = _canvas([_shape(0.1, 0.1, 0.3, 0.2, fill_role="primary",
                        label="unreadable", label_color_role="ink"),
                 _text(0.6, 0.6, 0.3, 0.2, "x")])
    assert any("label_color_role" in p for p in check_canvas_slide(s))


def test_light_label_on_light_fill_flagged():
    s = _canvas([_shape(0.1, 0.1, 0.3, 0.2, fill_role="surface",
                        label="ghost text", label_color_role="inverse_ink"),
                 _text(0.6, 0.6, 0.3, 0.2, "x")])
    assert any("label_color_role='ink'" in p for p in check_canvas_slide(s))


def test_light_text_on_light_bg_flagged():
    s = _canvas([_text(0.1, 0.1, 0.5, 0.2, "ghost", color_role="inverse_ink"),
                 _text(0.1, 0.5, 0.3, 0.2, "fine")])
    assert any("light background" in p for p in check_canvas_slide(s))
    # the same text contained in a dark panel is readable: no problem
    s2 = _canvas([_shape(0.05, 0.05, 0.6, 0.4, fill_role="primary"),
                  _text(0.1, 0.1, 0.4, 0.2, "fine", z=1,
                        color_role="inverse_ink")])
    assert not [p for p in check_canvas_slide(s2) if "light background" in p]


def test_component_inside_dark_panel_flagged():
    comp = {"x": 0.1, "y": 0.15, "w": 0.35, "h": 0.3, "z": 1,
            "content": {"kind": "bullet_list",
                        "items": [{"text": "ink text on navy panel"}]}}
    s = _canvas([_shape(0.05, 0.05, 0.5, 0.5, fill_role="primary_dark"),
                 comp, _text(0.6, 0.6, 0.3, 0.2, "x")])
    assert any("unreadable on dark fills" in p for p in check_canvas_slide(s))
    # the same component on a light panel is the classic legal pattern
    s2 = _canvas([_shape(0.05, 0.05, 0.5, 0.5, fill_role="surface"),
                  dict(comp), _text(0.6, 0.6, 0.3, 0.2, "x")])
    assert not [p for p in check_canvas_slide(s2)
                if "unreadable on dark fills" in p]


# -- renderer ----------------------------------------------------------------

def _build(spec_dict, tmp_path):
    deck = DeckSpec.model_validate({
        "schema_version": 1, "theme": "consulting_navy",
        "meta": {"title": "c", "date": "17 Aug 2026", "footer_org": "DE"},
        "slides": [spec_dict]})
    out = tmp_path / "canvas.pptx"
    report = build_deck(deck, out)
    return out, report


def test_golden_renders_clean(tmp_path):
    golden = json.loads((REPO / "eval" / "goldens" / "canvas_hero.json")
                        .read_text(encoding="utf-8"))
    out, report = _build(golden["slides"][0], tmp_path)
    assert report.warnings == [] and report.truncations == []
    slide = Presentation(str(out)).slides[0]
    texts = " ".join(
        " ".join(" ".join((r.text or "") for p in s.text_frame.paragraphs
                          for r in p.runs).split())
        for s in slide.shapes if s.has_text_frame)
    assert "3.4x" in texts and "less water" in texts
    # the golden passes its own rails
    assert check_canvas_slide(
        CanvasSlideSpec.model_validate(golden["slides"][0])) == []


def test_z_order_paints_low_first(tmp_path):
    spec = {"slide_type": "canvas", "title": "Z",
            "elements": [
                _text(0.1, 0.1, 0.3, 0.2, "on top", z=5,
                      color_role="inverse_ink"),
                _shape(0.05, 0.05, 0.5, 0.4, z=0, fill_role="primary_dark"),
            ]}
    out, _ = _build(spec, tmp_path)
    slide = Presentation(str(out)).slides[0]
    names = [(s.shape_type, s.has_text_frame and "".join(
        r.text for p in s.text_frame.paragraphs for r in p.runs))
        for s in slide.shapes]
    # the dark panel autoshape must be added BEFORE the text box
    panel_idx = next(i for i, (t, _) in enumerate(names)
                     if str(t) == "AUTO_SHAPE (1)")
    text_idx = next(i for i, (_, txt) in enumerate(names)
                    if txt == "on top")
    assert panel_idx < text_idx


def test_overflow_text_is_fitted_not_spilled(tmp_path):
    long = "This sentence is far too long for a tiny box " * 6
    spec = {"slide_type": "canvas", "title": "Fit",
            "elements": [_text(0.1, 0.1, 0.25, 0.08, long),
                         _text(0.1, 0.5, 0.3, 0.2, "second")]}
    out, report = _build(spec, tmp_path)
    assert report.truncations  # reported, not silently spilled
    slide = Presentation(str(out)).slides[0]
    box = next(s for s in slide.shapes if s.has_text_frame and "far too long"
               in s.text_frame.text)
    # every written line must fit the frame width (the safety-net promise)
    from deckengine.core.fit_text import TextMeasurer
    m = TextMeasurer()
    for para in box.text_frame.paragraphs:
        w = sum(m.width_pt(r.text, r.font.name or "Segoe UI",
                           r.font.size.pt, bold=bool(r.font.bold))
                for r in para.runs if r.text)
        assert w * 12700 <= int(box.width) + 12700  # ≤ frame + 1pt


def test_snap_aligns_near_edges(tmp_path):
    spec = {"slide_type": "canvas", "title": "Snap",
            "elements": [
                _shape(0.100, 0.1, 0.3, 0.2, fill_role="surface"),
                _shape(0.108, 0.4, 0.3, 0.2, fill_role="surface"),
            ]}
    out, _ = _build(spec, tmp_path)
    slide = Presentation(str(out)).slides[0]
    lefts = sorted(int(s.left) for s in slide.shapes
                   if str(s.shape_type) == "AUTO_SHAPE (1)")
    assert abs(lefts[0] - lefts[1]) <= Emu(2)  # snapped to a shared edge


def test_component_leaf_renders(tmp_path):
    spec = {"slide_type": "canvas", "title": "Chart in a box",
            "elements": [
                {"x": 0.1, "y": 0.05, "w": 0.8, "h": 0.6, "z": 0,
                 "content": {"kind": "native_chart", "chart_type": "bar",
                             "categories": ["A", "B"],
                             "series": [{"name": "s", "values": [3, 1]}],
                             "sort": "desc", "highlight": None,
                             "annotation": None}},
                _text(0.1, 0.7, 0.8, 0.1, "takeaway line"),
            ]}
    out, _ = _build(spec, tmp_path)
    slide = Presentation(str(out)).slides[0]
    assert any(s.has_chart for s in slide.shapes)
